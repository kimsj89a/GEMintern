import sys
import os

# Add parent directory to path to import core_ppt_updater
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from core_ppt_updater import InvestmentHistoryUpdater
from pptx import Presentation

def main():
    if len(sys.argv) < 3:
        print("Usage: python update_investment_history.py <input_pptx> <output_pptx>")
        # Create a dummy PPTX for testing if arguments not provided
        print("Creating dummy input.pptx for testing...")
        prs = Presentation()
        # Create 9 empty slides
        for _ in range(9):
            prs.slides.add_slide(prs.slide_layouts[6]) # Blind layout
        prs.save("input.pptx")
        input_path = "input.pptx"
        output_path = "output.pptx"
    else:
        input_path = sys.argv[1]
        output_path = sys.argv[2]
        
    print(f"Processing {input_path} -> {output_path}")
    
    updater = InvestmentHistoryUpdater(input_path)
    
    # 1. Extract (Simulation)
    data = updater.extract_data()
    print("Extracted Data (Mock/Template):")
    for row in data:
        print(row)
        
    # 2. Update Slide
    # We assume slide index 8 (9th slide) exists because we created 9 slides or expect the user provided one.
    try:
        updater.update_slide(8)
        print("Slide 9 updated.")
    except Exception as e:
        print(f"Error updating slide: {e}")
        
    # 3. Save
    updater.save(output_path)
    print(f"Saved to {output_path}")

if __name__ == "__main__":
    main()
