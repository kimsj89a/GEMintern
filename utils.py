import io
import re
import os
import tempfile
import pandas as pd
import fitz  # PyMuPDF
from docx import Document
from docx.shared import Inches, Pt
from docx.oxml.ns import qn
from pptx import Presentation
import ocr

# OpenDataLoader PDF 지원 확인
OPENDATALOADER_AVAILABLE = False
try:
    import opendataloader_pdf
    OPENDATALOADER_AVAILABLE = True
except ImportError:
    pass

# MarkItDown 지원 확인
MARKITDOWN_AVAILABLE = False
try:
    from markitdown import MarkItDown
    MARKITDOWN_AVAILABLE = True
except ImportError:
    pass

# Document AI OCR 지원 확인
DOCAI_AVAILABLE = False
try:
    import utils_docai
    DOCAI_AVAILABLE = True
except ImportError:
    pass

def get_ocr_status():
    """OCR 상태 확인 (UI에서 사용)"""
    return ocr.get_ocr_status()

def _docx_to_ppt_markdown(doc: Document, filename: str) -> str:
    """
    Convert a Word document into PPT-friendly Markdown.

    Mapping:
    - Heading 1 -> section cover (#)
    - Heading 2 -> slide title (##)
    - Heading 3+ -> emphasized line inside a slide (###)
    - Normal paragraphs -> bullets (-)
    """
    title = os.path.splitext(filename)[0]
    lines = [f"# {title}"]

    def has_slide_title() -> bool:
        return any(line.startswith("## ") for line in lines)

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue

        style_name = ""
        try:
            style_name = (para.style.name or "").lower()
        except Exception:
            style_name = ""

        if style_name.startswith("heading"):
            level = 2
            match = re.search(r"heading\s*(\d+)", style_name)
            if match:
                try:
                    level = int(match.group(1))
                except Exception:
                    level = 2

            if level <= 1:
                lines.append(f"# {text}")
            elif level == 2:
                lines.append(f"## {text}")
            else:
                lines.append(f"### {text}")
            continue

        if not has_slide_title():
            lines.append("## Overview")

        lines.append(f"- {text}")

    return "\n".join(lines).strip() + "\n\n"


MAX_FILE_SIZE_MB = 50       # 파일 크기 제한 (MB)
MAX_PDF_PAGES = 200         # PDF 최대 페이지 수
MAX_OCR_PAGES = 50          # Gemini Vision OCR 최대 페이지 수


def parse_uploaded_file(uploaded_file, api_key=None, docai_config=None, template_option=None):
    """파일 형태별 텍스트 추출 (전체 시트 지원 + OCR 지원)

    Args:
        uploaded_file: Streamlit 업로드 파일 객체
        api_key: Google API 키 (PDF OCR용, 선택사항)
        docai_config: Document AI 설정 dict (선택사항)
            - project_id: GCP 프로젝트 ID
            - location: 위치 (us/eu)
            - processor_id: 프로세서 ID
            - credentials_json: 서비스 계정 JSON 문자열
    """
    if uploaded_file is None:
        return ""

    # --- 파일 크기 제한 ---
    file_size_mb = uploaded_file.size / (1024 * 1024) if hasattr(uploaded_file, 'size') else 0
    if file_size_mb > MAX_FILE_SIZE_MB:
        return (
            f"### [파일명: {uploaded_file.name} - SKIPPED]\n"
            f"⚠️ 파일 크기 초과: {file_size_mb:.1f}MB (제한: {MAX_FILE_SIZE_MB}MB)\n"
            f"파일이 너무 커서 처리를 건너뛰었습니다.\n\n"
        )

    file_type = uploaded_file.name.split('.')[-1].lower()

    # [Document AI OCR] PDF/이미지 우선 처리
    if DOCAI_AVAILABLE and docai_config and file_type in utils_docai.get_supported_extensions():
        try:
            uploaded_file.seek(0)
            file_bytes = uploaded_file.read()
            mime_type = utils_docai.get_mime_type(uploaded_file.name)

            ocr_result = utils_docai.process_document(
                file_bytes=file_bytes,
                mime_type=mime_type,
                project_id=docai_config['project_id'],
                location=docai_config.get('location', 'us'),
                processor_id=docai_config['processor_id'],
                credentials_json=docai_config.get('credentials_json')
            )

            uploaded_file.seek(0)
            if ocr_result and ocr_result.get('text'):
                return f"### [파일명: {uploaded_file.name} (Document AI OCR)]\n{ocr_result['text']}\n\n"
        except Exception as e:
            uploaded_file.seek(0)
            # Document AI 실패 시 다음 방법으로 진행

    # [OpenDataLoader PDF] PDF 전용 고정밀 파서 (표 추출 0.93 정확도)
    if OPENDATALOADER_AVAILABLE and file_type == 'pdf':
        try:
            import logging as _logging
            _logging.getLogger(__name__).info(f"OpenDataLoader PDF: parsing {uploaded_file.name}")
            suffix = os.path.splitext(uploaded_file.name)[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                uploaded_file.seek(0)
                tmp.write(uploaded_file.read())
                tmp_path = tmp.name
            uploaded_file.seek(0)

            out_dir = tempfile.mkdtemp(prefix="odl_")
            try:
                import concurrent.futures
                def _run_odl(path, out):
                    opendataloader_pdf.convert(
                        input_path=[path],
                        output_dir=out,
                        format="markdown",
                    )
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(_run_odl, tmp_path, out_dir)
                    future.result(timeout=120)  # 2분 타임아웃

                # 결과 .md 파일 읽기
                md_files = [f for f in os.listdir(out_dir) if f.endswith('.md')]
                if md_files:
                    md_path = os.path.join(out_dir, md_files[0])
                    with open(md_path, 'r', encoding='utf-8') as mf:
                        md_content = mf.read()
                    if md_content.strip():
                        return f"### [파일명: {uploaded_file.name} (OpenDataLoader)]\n{md_content}\n\n"
            finally:
                # 임시 파일 정리
                import shutil
                if os.path.exists(tmp_path):
                    try: os.unlink(tmp_path)
                    except OSError: pass
                if os.path.exists(out_dir):
                    try: shutil.rmtree(out_dir)
                    except OSError: pass
        except Exception as e:
            import logging as _logging
            _logging.getLogger(__name__).debug(f"OpenDataLoader failed, falling back: {e}")
            uploaded_file.seek(0)

    # [MarkItDown] 우선 시도 (타임아웃 60초)
    # PPT 모드에서 Word는 별도 변환 로직을 사용한다.
    if MARKITDOWN_AVAILABLE and not (template_option == "presentation" and file_type in ["docx", "doc"]):
        try:
            suffix = os.path.splitext(uploaded_file.name)[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                uploaded_file.seek(0)
                tmp.write(uploaded_file.read())
                tmp_path = tmp.name
            uploaded_file.seek(0)

            try:
                import concurrent.futures
                def _convert_markitdown(path):
                    md = MarkItDown()
                    return md.convert(path)
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(_convert_markitdown, tmp_path)
                    result = future.result(timeout=60)
                if result and result.text_content:
                    return f"### [파일명: {uploaded_file.name} (MarkItDown)]\n{result.text_content}\n\n"
            finally:
                if os.path.exists(tmp_path):
                    try: os.unlink(tmp_path)
                    except: pass
        except Exception:
            uploaded_file.seek(0)

    file_type = uploaded_file.name.split('.')[-1].lower()
    text_content = ""

    try:
        # [PDF] PyMuPDF + Gemini Vision OCR (페이지 수 제한 적용)
        if file_type == 'pdf':
            with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
                total_pages = len(doc)
                if total_pages > MAX_PDF_PAGES:
                    text_content = (
                        f"### [파일명: {uploaded_file.name} - TRUNCATED]\n"
                        f"⚠️ PDF 페이지 초과: {total_pages}페이지 (제한: {MAX_PDF_PAGES}페이지)\n"
                        f"처음 {MAX_PDF_PAGES}페이지만 처리합니다.\n\n"
                    )
                    doc_to_process = fitz.open()
                    for i in range(MAX_PDF_PAGES):
                        doc_to_process.insert_pdf(doc, from_page=i, to_page=i)
                else:
                    text_content = ""
                    doc_to_process = doc

                if api_key:
                    text_content += ocr.extract_pdf_with_gemini_ocr(
                        doc_to_process, api_key, max_ocr_pages=MAX_OCR_PAGES,
                    )
                else:
                    text_content += ocr.extract_pdf_with_ocr(doc_to_process)

                if doc_to_process is not doc:
                    doc_to_process.close()
        
        # [Word] python-docx
        elif file_type in ['docx', 'doc']:
            doc = Document(uploaded_file)
            if template_option == "presentation":
                text_content = _docx_to_ppt_markdown(doc, uploaded_file.name)
            else:
                for para in doc.paragraphs:
                    text_content += para.text + "\n"

        # [PPT] python-pptx
        elif file_type in ['pptx', 'ppt']:
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

        # [Excel] pandas (전체 시트 파싱 적용)
        elif file_type in ['xlsx', 'xls', 'csv']:
            try:
                text_content = f"### [파일명: {uploaded_file.name}]\n"

                # 1. 파일 읽기 (CSV vs Excel)
                if file_type == 'csv':
                    df = pd.read_csv(uploaded_file)
                    df = df.fillna("")
                    try:
                        table_text = df.to_markdown(index=False)
                    except ImportError:
                        table_text = df.to_string(index=False)
                    text_content += f"\n{table_text}\n"
                else:
                    # [특별 변경] sheet_name=None으로 설정하여 모든 시트를 OrderedDict로 읽어옴
                    # 명시적 openpyxl 명시적으로 사용 (안정성)
                    xls_dict = pd.read_excel(uploaded_file, sheet_name=None, engine='openpyxl')

                    # 모든 시트 조회
                    for sheet_name, df in xls_dict.items():
                        df = df.fillna("") # 빈값 처리

                        # 시트별 헤더 추가
                        text_content += f"\n#### [Sheet: {sheet_name}]\n"

                        # 변환 (tabulate가 없으면 to_string으로 대체)
                        try:
                            table_text = df.to_markdown(index=False)
                        except ImportError:
                            table_text = df.to_string(index=False)

                        text_content += f"{table_text}\n"

            except Exception as e:
                text_content = f"[엑셀 파싱 오류: {str(e)}]\n(Tip: 암호 걸린 파일인 아닌지, 형식에 맞는지 확인해주세요)"

        # [Text]
        elif file_type in ['txt', 'md']:
            stringio = io.StringIO(uploaded_file.getvalue().decode("utf-8"))
            text_content = stringio.read()

        else:
            text_content = f"[지원하지 않는 파일 형식입니다: {uploaded_file.name}]"

    except Exception as e:
        return f"[파일 읽기 시도 중 오류: {uploaded_file.name} - {str(e)}]"

    # 파일 포인터 초기화
    if hasattr(uploaded_file, 'seek'):
        uploaded_file.seek(0)

    return f"### [파일명: {uploaded_file.name}]\n{text_content}\n\n"

def extract_title_from_markdown(text):
    """마크다운 텍스트에서 첫 번째 헤딩 또는 첫 줄을 제목으로 추출."""
    if not text:
        return None
    for line in text.split('\n'):
        line = line.strip()
        if line.startswith('#'):
            title = line.lstrip('#').strip()
            title = re.sub(r'[\\/*?:"<>|]', '', title).strip()
            if title:
                return title[:50].rstrip()
    # 헤딩이 없으면 첫 번째 비어있지 않은 줄 사용
    for line in text.split('\n'):
        line = line.strip()
        if line:
            title = re.sub(r'[\\/*?:"<>|]', '', line).strip()
            if title:
                return title[:50].rstrip()
    return None


def generate_filename(uploaded_files, template_option, generated_text=None):
    template_map = {
        "simple_review": "simple_review",
        "rfi": "RFI",
        "investment": "investment_report",
        "im": "IM",
        "management": "management_report",
        "presentation": "presentation",
        "custom": "report",
    }
    suffix = template_map.get(template_option, "report")

    # 생성된 텍스트에서 제목 추출 시도
    if generated_text:
        title = extract_title_from_markdown(generated_text)
        if title:
            return f"{title}.docx"

    project_name = "Investment_Report"
    if uploaded_files:
        first_file = uploaded_files[0].name
        base_name = os.path.splitext(first_file)[0]
        project_name = re.sub(r'[\\/*?:"<>|]', "", base_name).strip()
    return f"{project_name}_{suffix}.docx"

def add_list_paragraph(doc, content, level, is_bullet=True):
    """들여쓰기가 적용된 리스트 아이템 추가

    Args:
        doc: Document 객체
        content: 텍스트 내용
        level: 들여쓰기 레벨 (0부터 시작)
        is_bullet: True면 불릿, False면 번호
    """

    # Bullet characters by level (fallback to simple ASCII bullets)
    bullet_chars = ["-", "*", "+"]
    bullet_char = bullet_chars[level % len(bullet_chars)]

    p = doc.add_paragraph()

    # 들여쓰기 설정 (레벨당 0.1인치)
    indent = Inches(0.1 * (level + 1))
    p.paragraph_format.left_indent = indent
    p.paragraph_format.first_line_indent = Inches(-0.15)  # 불릿/번호 hanging indent

    # 불릿 문자 추가
    if is_bullet:
        p.add_run(f"{bullet_char} ")
    else:
        p.add_run(f"• ")  # 번호 리스트도 일단 불릿으로

    # 내용 추가 (볼드 처리 포함)
    parts = re.split(r'(\*\*.*?\*\*)', content)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = p.add_run(part[2:-2])
            run.bold = True
        else:
            p.add_run(part)

    return p

def _md_add_hyperlink(paragraph, text, url, font_name='맑은 고딕'):
    from docx.oxml import OxmlElement
    part = paragraph.part
    r_id = part.relate_to(
        url,
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
        is_external=True,
    )
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), r_id)
    new_run = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    rFonts = OxmlElement('w:rFonts')
    rFonts.set(qn('w:ascii'), font_name)
    rFonts.set(qn('w:hAnsi'), font_name)
    rFonts.set(qn('w:eastAsia'), font_name)
    rPr.append(rFonts)
    color = OxmlElement('w:color'); color.set(qn('w:val'), '0563C1'); rPr.append(color)
    u = OxmlElement('w:u'); u.set(qn('w:val'), 'single'); rPr.append(u)
    new_run.append(rPr)
    t = OxmlElement('w:t'); t.text = text; t.set(qn('xml:space'), 'preserve')
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


def _md_render_inline(paragraph, inline_token, *, font_name='맑은 고딕', mono_font='Consolas'):
    """inline 토큰의 children을 단락에 렌더링. bold/italic/strike/code/link/softbreak/hardbreak 지원."""
    marks = {'bold': False, 'italic': False, 'strike': False, 'code': False}
    link_stack = []  # [(text_buffer, url)]
    children = getattr(inline_token, 'children', None) or []

    def emit_run(text):
        if not text:
            return
        if link_stack:
            link_stack[-1][0].append(text)
            return
        if marks['code']:
            run = paragraph.add_run(text)
            run.font.name = mono_font
            run._element.rPr.rFonts.set(qn('w:eastAsia'), mono_font)
        else:
            run = paragraph.add_run(text)
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
        if marks['bold']:
            run.bold = True
        if marks['italic']:
            run.italic = True
        if marks['strike']:
            run.font.strike = True

    for tok in children:
        t = tok.type
        if t == 'text':
            emit_run(tok.content)
        elif t == 'strong_open':
            marks['bold'] = True
        elif t == 'strong_close':
            marks['bold'] = False
        elif t == 'em_open':
            marks['italic'] = True
        elif t == 'em_close':
            marks['italic'] = False
        elif t == 's_open':
            marks['strike'] = True
        elif t == 's_close':
            marks['strike'] = False
        elif t == 'code_inline':
            marks['code'] = True
            emit_run(tok.content)
            marks['code'] = False
        elif t == 'link_open':
            url = tok.attrGet('href') or ''
            link_stack.append(([], url))
        elif t == 'link_close':
            buf, url = link_stack.pop()
            text = ''.join(buf) or url
            if url:
                _md_add_hyperlink(paragraph, text, url, font_name=font_name)
            else:
                emit_run(text)
        elif t == 'softbreak':
            emit_run(' ')
        elif t == 'hardbreak':
            paragraph.add_run().add_break()
        elif t == 'image':
            alt = tok.attrGet('alt') or tok.content or ''
            src = tok.attrGet('src') or ''
            emit_run(f"[이미지: {alt or src}]")
        else:
            inner = getattr(tok, 'content', '')
            if inner:
                emit_run(inner)


def _md_set_paragraph_shading(paragraph, fill_hex):
    from docx.oxml import OxmlElement
    pPr = paragraph._p.get_or_add_pPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), fill_hex)
    pPr.append(shd)


def _md_set_left_border(paragraph, color_hex='808080', size=24):
    from docx.oxml import OxmlElement
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    left = OxmlElement('w:left')
    left.set(qn('w:val'), 'single')
    left.set(qn('w:sz'), str(size))
    left.set(qn('w:space'), '4')
    left.set(qn('w:color'), color_hex)
    pBdr.append(left)
    pPr.append(pBdr)


def _md_add_hr(doc):
    from docx.oxml import OxmlElement
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '6')
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), 'auto')
    pBdr.append(bottom)
    pPr.append(pBdr)


def create_docx(markdown_text):
    """markdown-it-py AST 기반 Markdown → DOCX 변환.

    지원: heading 1~6, paragraph(bold/italic/strike/code/link),
          bullet/ordered list(nested), table(GFM), blockquote, code block(fence),
          hr, softbreak/hardbreak.
    """
    if os.environ.get('MD_TO_DOCX_LEGACY') == '1':
        return _create_docx_legacy(markdown_text)

    try:
        from markdown_it import MarkdownIt
    except Exception:
        return _create_docx_legacy(markdown_text)

    font_name = '맑은 고딕'
    mono_font = 'Consolas'

    template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'template', 'Normal.docx')
    if os.path.exists(template_path):
        try:
            doc = Document(template_path)
            for el in list(doc.element.body):
                if el.tag.endswith('}p') or el.tag.endswith('}tbl'):
                    doc.element.body.remove(el)
        except Exception:
            doc = Document()
    else:
        doc = Document()

    md = MarkdownIt('commonmark', {'breaks': False, 'html': False})
    for plugin in ('table', 'strikethrough'):
        try:
            md.enable(plugin)
        except Exception:
            pass

    tokens = md.parse(markdown_text or '')

    # Token stream walker with stacks
    list_stack = []  # list of dicts: {'type': 'ul'|'ol', 'level': int, 'index': int}
    blockquote_depth = 0

    def current_list_level():
        return len(list_stack) - 1

    i = 0
    n = len(tokens)
    while i < n:
        tok = tokens[i]
        t = tok.type

        if t == 'heading_open':
            level = int(tok.tag[1])  # h1 → 1
            inline = tokens[i + 1] if i + 1 < n else None
            heading = doc.add_heading('', level=min(level, 9))
            if inline is not None and inline.type == 'inline':
                _md_render_inline(heading, inline, font_name=font_name, mono_font=mono_font)
            i += 3  # heading_open, inline, heading_close
            continue

        if t == 'paragraph_open':
            inline = tokens[i + 1] if i + 1 < n else None
            if list_stack:
                lvl = current_list_level()
                is_bullet = list_stack[-1]['type'] == 'ul'
                p = doc.add_paragraph()
                indent = Inches(0.25 * (lvl + 1))
                p.paragraph_format.left_indent = indent
                p.paragraph_format.first_line_indent = Inches(-0.2)
                if is_bullet:
                    bullets = ['•', '◦', '▪']
                    prefix = bullets[lvl % len(bullets)] + ' '
                else:
                    list_stack[-1]['index'] += 1
                    prefix = f"{list_stack[-1]['index']}. "
                run = p.add_run(prefix)
                run.font.name = font_name
                run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                if inline is not None and inline.type == 'inline':
                    _md_render_inline(p, inline, font_name=font_name, mono_font=mono_font)
            elif blockquote_depth > 0:
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Inches(0.3 * blockquote_depth)
                _md_set_left_border(p)
                if inline is not None and inline.type == 'inline':
                    _md_render_inline(p, inline, font_name=font_name, mono_font=mono_font)
            else:
                p = doc.add_paragraph()
                if inline is not None and inline.type == 'inline':
                    _md_render_inline(p, inline, font_name=font_name, mono_font=mono_font)
            i += 3
            continue

        if t == 'bullet_list_open':
            list_stack.append({'type': 'ul', 'level': current_list_level() + 1, 'index': 0})
            i += 1
            continue
        if t == 'ordered_list_open':
            start = 1
            try:
                start = int(tok.attrGet('start') or 1)
            except Exception:
                start = 1
            list_stack.append({'type': 'ol', 'level': current_list_level() + 1, 'index': start - 1})
            i += 1
            continue
        if t in ('bullet_list_close', 'ordered_list_close'):
            if list_stack:
                list_stack.pop()
            i += 1
            continue

        if t in ('list_item_open', 'list_item_close'):
            i += 1
            continue

        if t == 'blockquote_open':
            blockquote_depth += 1
            i += 1
            continue
        if t == 'blockquote_close':
            blockquote_depth = max(0, blockquote_depth - 1)
            i += 1
            continue

        if t == 'fence' or t == 'code_block':
            code = tok.content or ''
            for line in code.rstrip('\n').split('\n'):
                p = doc.add_paragraph()
                _md_set_paragraph_shading(p, 'F2F2F2')
                p.paragraph_format.left_indent = Inches(0.1)
                run = p.add_run(line if line else ' ')
                run.font.name = mono_font
                run._element.rPr.rFonts.set(qn('w:eastAsia'), mono_font)
                run.font.size = Pt(9.5)
            i += 1
            continue

        if t == 'hr':
            _md_add_hr(doc)
            i += 1
            continue

        if t == 'table_open':
            # Collect table tokens
            j = i
            depth = 1
            while j + 1 < n and depth > 0:
                j += 1
                if tokens[j].type == 'table_open':
                    depth += 1
                elif tokens[j].type == 'table_close':
                    depth -= 1
            table_tokens = tokens[i:j + 1]
            # Parse rows
            rows = []  # list of list of (inline_token, is_header)
            cur_row = None
            in_header = False
            for tt in table_tokens:
                tp = tt.type
                if tp == 'thead_open':
                    in_header = True
                elif tp == 'thead_close':
                    in_header = False
                elif tp == 'tr_open':
                    cur_row = []
                elif tp == 'tr_close':
                    if cur_row is not None:
                        rows.append(cur_row)
                    cur_row = None
                elif tp in ('th_open', 'td_open'):
                    # next token should be inline
                    pass
                elif tp == 'inline':
                    if cur_row is not None:
                        cur_row.append((tt, in_header))
            if rows:
                cols = max(len(r) for r in rows)
                table = doc.add_table(rows=len(rows), cols=cols)
                table.style = 'Table Grid'
                for r_idx, row in enumerate(rows):
                    for c_idx in range(cols):
                        cell = table.rows[r_idx].cells[c_idx]
                        cell.text = ''  # clear default paragraph content
                        para = cell.paragraphs[0]
                        if c_idx < len(row):
                            inline_tok, is_header = row[c_idx]
                            _md_render_inline(para, inline_tok, font_name=font_name, mono_font=mono_font)
                            if is_header:
                                for run in para.runs:
                                    run.bold = True
                                _md_set_paragraph_shading(para, 'F2F2F2')
            i = j + 1
            continue

        # Skip unhandled
        i += 1

    # --- 후처리: 폰트/간격 일괄 ---
    normal = doc.styles['Normal']
    normal.font.name = font_name
    normal.font.size = Pt(10)
    normal.element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
    normal.paragraph_format.space_after = Pt(0)

    heading_sizes = {1: Pt(16), 2: Pt(13), 3: Pt(11), 4: Pt(10.5), 5: Pt(10), 6: Pt(10)}
    for level, size in heading_sizes.items():
        style_name = f'Heading {level}'
        if style_name in doc.styles:
            hs = doc.styles[style_name]
            hs.font.name = font_name
            hs.font.size = size
            hs.element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
            hs.paragraph_format.space_after = Pt(0)

    for paragraph in doc.paragraphs:
        paragraph.paragraph_format.space_after = Pt(0)
        for run in paragraph.runs:
            if not run.font.name:
                run.font.name = font_name
                run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    paragraph.paragraph_format.space_after = Pt(0)
                    for run in paragraph.runs:
                        if not run.font.name:
                            run.font.name = font_name
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _create_docx_legacy(markdown_text):
    """레거시 regex 기반 MD→DOCX (fallback). MD_TO_DOCX_LEGACY=1 환경변수로 강제 사용 가능."""
    doc = Document()
    lines = markdown_text.split('\n')
    i = 0

    roman_header_pattern = re.compile(r'^(I{1,3}|IV|VI{0,3}|V|IX|X)\.\s+(.+)$')
    indent_stack = [0]

    while i < len(lines):
        raw_line = lines[i]
        line = raw_line.strip()

        if line.startswith('##### '):
            doc.add_heading(line.replace('##### ', ''), level=5)
            indent_stack = [0]
            i += 1
        elif line.startswith('#### '):
            doc.add_heading(line.replace('#### ', ''), level=4)
            indent_stack = [0]
            i += 1
        elif line.startswith('### '):
            doc.add_heading(line.replace('### ', ''), level=3)
            indent_stack = [0]
            i += 1
        elif line.startswith('## '):
            doc.add_heading(line.replace('## ', ''), level=2)
            indent_stack = [0]
            i += 1
        elif line.startswith('# '):
            doc.add_heading(line.replace('# ', ''), level=1)
            indent_stack = [0]
            i += 1
        elif roman_header_pattern.match(line):
            doc.add_heading(line, level=1)
            indent_stack = [0]
            i += 1
        elif line.startswith('|'):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                table_lines.append(lines[i].strip())
                i += 1
            if len(table_lines) >= 2:
                headers = [c.strip() for c in table_lines[0].split('|') if c.strip()]
                data_rows = []
                for row_line in table_lines[1:]:
                    if '---' in row_line: continue
                    parts = row_line.split('|')
                    if len(parts) >= 2: data_rows.append([c.strip() for c in parts[1:-1]])
                if headers:
                    table = doc.add_table(rows=1, cols=len(headers))
                    table.style = 'Table Grid'
                    for idx, text in enumerate(headers):
                        if idx < len(table.rows[0].cells):
                            cell = table.rows[0].cells[idx]
                            cell.text = text.replace('**', '')
                            if cell.paragraphs[0].runs:
                                cell.paragraphs[0].runs[0].bold = True
                    for row_data in data_rows:
                        row_cells = table.add_row().cells
                        for idx, text in enumerate(row_data):
                            if idx < len(row_cells): row_cells[idx].text = text.replace('**', '')
            indent_stack = [0]
        elif re.match(r'^\s*([-*•]|\d+\.)\s', raw_line):
            match = re.match(r'^(\s*)([-*•]|\d+\.)\s+(.*)', raw_line)
            if match:
                indent_str, marker, content = match.groups()
                spaces = indent_str.replace('\t', '    ')
                indent_len = len(spaces)
                if indent_len == 0:
                    level = 0
                    indent_stack = [0]
                elif indent_len > indent_stack[-1]:
                    level = len(indent_stack)
                    indent_stack.append(indent_len)
                else:
                    while len(indent_stack) > 1 and indent_stack[-1] > indent_len:
                        indent_stack.pop()
                    level = len(indent_stack) - 1
                if level > 8: level = 8
                is_bullet = marker in ['-', '*', '•']
                add_list_paragraph(doc, content, level, is_bullet)
            i += 1
        else:
            if line:
                p = doc.add_paragraph()
                parts = re.split(r'(\*\*.*?\*\*)', line)
                for part in parts:
                    if part.startswith('**') and part.endswith('**'):
                        run = p.add_run(part[2:-2])
                        run.bold = True
                    else: p.add_run(part)
            i += 1

    font_name = '맑은 고딕'
    normal = doc.styles['Normal']
    normal.font.name = font_name
    normal.font.size = Pt(10)
    normal.element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
    normal.paragraph_format.space_after = Pt(0)

    heading_sizes = {1: Pt(16), 2: Pt(13), 3: Pt(11), 4: Pt(10.5), 5: Pt(10)}
    for level, size in heading_sizes.items():
        style_name = f'Heading {level}'
        if style_name in doc.styles:
            hs = doc.styles[style_name]
            hs.font.name = font_name
            hs.font.size = size
            hs.element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
            hs.paragraph_format.space_after = Pt(0)

    for paragraph in doc.paragraphs:
        paragraph.paragraph_format.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    paragraph.paragraph_format.space_after = Pt(0)
                    for run in paragraph.runs:
                        run.font.name = font_name
                        run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def create_excel(markdown_text):
    data = []
    lines = markdown_text.split('\n')
    for line in lines:
        line = line.strip()
        if line.startswith('|') and '---' not in line:
            row = [c.strip().replace('**', '') for c in line.split('|')[1:-1]]
            if row: data.append(row)
    bio = io.BytesIO()
    if data:
        df = pd.DataFrame(data[1:], columns=data[0])
        with pd.ExcelWriter(bio, engine='xlsxwriter') as writer:
            df.to_excel(writer, index=False, sheet_name='RFI_List')
    return bio.getvalue()

# ========================================
# Local Storage Management (RAG)
# ========================================
SAVED_DOCS_DIR = "saved_documents"

def ensure_saved_docs_dir():
    if not os.path.exists(SAVED_DOCS_DIR):
        os.makedirs(SAVED_DOCS_DIR)

def save_to_local_storage(filename, content):
    """파싱된 텍스트를 로컬 MD 파일로 저장"""
    ensure_saved_docs_dir()
    safe_name = os.path.basename(filename)
    base, _ = os.path.splitext(safe_name)
    save_path = os.path.join(SAVED_DOCS_DIR, f"{base}.md")
    
    with open(save_path, "w", encoding="utf-8") as f:
        f.write(content)
    return save_path

def list_saved_docs():
    ensure_saved_docs_dir()
    files = [f for f in os.listdir(SAVED_DOCS_DIR) if f.endswith(".md")]
    return sorted(files)

def load_saved_doc(filename):
    ensure_saved_docs_dir()
    path = os.path.join(SAVED_DOCS_DIR, filename)
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    return ""

# ========================================
# Project Management (프로젝트별 자료 관리)
# ========================================
PROJECTS_DIR = os.path.join(SAVED_DOCS_DIR, "projects")

def _ensure_projects_dir():
    if not os.path.exists(PROJECTS_DIR):
        os.makedirs(PROJECTS_DIR)

def list_projects():
    """저장된 프로젝트 목록 반환"""
    _ensure_projects_dir()
    return sorted([
        d for d in os.listdir(PROJECTS_DIR)
        if os.path.isdir(os.path.join(PROJECTS_DIR, d))
    ])

def create_project(name):
    """새 프로젝트 디렉토리 생성"""
    _ensure_projects_dir()
    safe_name = re.sub(r'[\\/*?:"<>|]', "", name).strip()
    if not safe_name:
        return None
    path = os.path.join(PROJECTS_DIR, safe_name)
    if not os.path.exists(path):
        os.makedirs(path)
    return safe_name

def delete_project(name):
    """프로젝트 삭제"""
    import shutil
    path = os.path.join(PROJECTS_DIR, name)
    if os.path.exists(path):
        shutil.rmtree(path)

def add_doc_to_project(project_name, filename, content):
    """파싱된 문서를 프로젝트에 저장"""
    path = os.path.join(PROJECTS_DIR, project_name)
    if not os.path.exists(path):
        os.makedirs(path)
    safe_name = os.path.basename(filename)
    base, _ = os.path.splitext(safe_name)
    save_path = os.path.join(path, f"{base}.md")
    with open(save_path, "w", encoding="utf-8") as f:
        f.write(content)
    return save_path

def remove_doc_from_project(project_name, filename):
    """프로젝트에서 문서 제거"""
    path = os.path.join(PROJECTS_DIR, project_name, filename)
    if os.path.exists(path):
        os.remove(path)

def list_project_docs(project_name):
    """프로젝트 내 문서 목록"""
    path = os.path.join(PROJECTS_DIR, project_name)
    if not os.path.exists(path):
        return []
    return sorted([f for f in os.listdir(path) if f.endswith(".md")])

def load_all_project_docs(project_name):
    """프로젝트 전체 문서 텍스트 반환"""
    docs = list_project_docs(project_name)
    all_text = ""
    for doc_name in docs:
        path = os.path.join(PROJECTS_DIR, project_name, doc_name)
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                all_text += f.read() + "\n\n"
    return all_text


def unlock_pdf(pdf_bytes: bytes, password: str) -> bytes:
    """비밀번호로 보호된 PDF의 잠금을 해제하여 새 PDF 바이트를 반환한다.

    Args:
        pdf_bytes: 원본 PDF 파일의 바이트 데이터
        password: PDF 비밀번호

    Returns:
        잠금 해제된 PDF의 바이트 데이터

    Raises:
        ValueError: 비밀번호가 틀리거나 PDF를 열 수 없는 경우
    """
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as e:
        raise ValueError(f"PDF 파일을 열 수 없습니다: {e}")

    if not doc.is_encrypted:
        # 이미 잠금 해제된 PDF
        result = doc.tobytes()
        doc.close()
        return result

    if not doc.authenticate(password):
        doc.close()
        raise ValueError("비밀번호가 올바르지 않습니다.")

    # 잠금 해제된 PDF를 새로 저장 (암호화 없이)
    result = doc.tobytes(garbage=3, deflate=True)
    doc.close()
    return result
