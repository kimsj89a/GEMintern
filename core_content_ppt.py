"""
Content-first 4:3 PPT 빌더.

설계 원칙:
- LLM 은 슬라이드별 콘텐츠 블록(sentences/tables/charts/sources) 만 생성
- Layout 결정은 코드 (블록 수·종류 → 7개 layout 패턴 중 자동 매칭)
- 4:3 캔버스 (10" × 7.5"), 노앤 PE 톤 (네이비 + 그레이)
- 차트는 matplotlib → PNG 임베드 (네이티브 차트 대신 품질 ↑)

슬라이드 입력 schema (LLM 출력):
{
  "title": "한 줄 슬라이드 제목",
  "summary": "1~2줄 핵심 메시지",          (선택)
  "sentences": ["문장1", "문장2", ...],    (선택, bullet 으로 배치)
  "tables": [{
      "title": "표 제목",                   (선택)
      "headers": ["연도", "매출", "EBITDA"],
      "rows": [["2024", "100", "20"], ...]
  }],                                       (선택, 0~2개)
  "charts": [{
      "title": "차트 제목",                 (선택)
      "chart_type": "bar"|"line"|"pie",
      "categories": ["2023", "2024", "2025"],
      "series": [{"name": "매출", "values": [80, 100, 120]}],
      "y_label": "억원"                     (선택)
  }],                                       (선택, 0~2개)
  "sources": ["출처1", "출처2"]             (선택, footer)
}
"""
from __future__ import annotations

import io
import logging
import os
import tempfile
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# ── 양식 (노앤 PE 톤, 4:3) ──
SLIDE_W_IN = 10.0
SLIDE_H_IN = 7.5

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_LIGHT = RGBColor(0x2D, 0x4A, 0x7C)
GREY_900 = RGBColor(0x1F, 0x29, 0x37)
GREY_700 = RGBColor(0x4B, 0x55, 0x63)
GREY_500 = RGBColor(0x6B, 0x72, 0x80)
GREY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GREY_100 = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_ACCENT = RGBColor(0x2B, 0x6C, 0xB0)

FONT_FAMILY = "Pretendard, Malgun Gothic, Arial"

# ── matplotlib (lazy import) ──

def _make_chart_png(chart: Dict[str, Any], width_in: float, height_in: float) -> Optional[bytes]:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        logger.warning("matplotlib not available, skipping chart")
        return None

    ctype = (chart.get("chart_type") or "bar").lower()
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    if not cats or not series:
        return None

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    palette = ["#1B2A4A", "#2B6CB0", "#6B7280", "#94A3B8", "#CBD5E1"]

    try:
        if ctype == "pie":
            s0 = series[0]
            vals = [_to_num(v) for v in (s0.get("values") or [])]
            ax.pie(vals, labels=cats, colors=palette[:len(vals)],
                   autopct="%1.0f%%", textprops={"fontsize": 9})
        elif ctype == "line":
            for i, s in enumerate(series):
                vals = [_to_num(v) for v in (s.get("values") or [])]
                ax.plot(cats, vals, marker="o", linewidth=2,
                        color=palette[i % len(palette)], label=s.get("name", f"S{i+1}"))
            if len(series) > 1:
                ax.legend(loc="best", fontsize=8, frameon=False)
        else:  # bar (grouped if multiple series)
            n = len(series)
            x = list(range(len(cats)))
            width = 0.8 / max(n, 1)
            for i, s in enumerate(series):
                vals = [_to_num(v) for v in (s.get("values") or [])]
                offsets = [xi + (i - (n - 1) / 2) * width for xi in x]
                ax.bar(offsets, vals, width=width, color=palette[i % len(palette)],
                       label=s.get("name", f"S{i+1}"))
            ax.set_xticks(x)
            ax.set_xticklabels(cats, fontsize=8)
            if n > 1:
                ax.legend(loc="best", fontsize=8, frameon=False)

        if ctype != "pie":
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            ax.tick_params(labelsize=8, colors="#4B5563")
            for spine in ("left", "bottom"):
                ax.spines[spine].set_color("#D1D5DB")
            if chart.get("y_label"):
                ax.set_ylabel(chart["y_label"], fontsize=8, color="#6B7280")
            ax.grid(axis="y", linestyle=":", linewidth=0.5, color="#E5E7EB")

        plt.tight_layout(pad=0.5)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        plt.close(fig)
        return buf.getvalue()
    except Exception as e:
        logger.warning(f"chart render failed: {e}")
        try: plt.close(fig)
        except Exception: pass
        return None


def _to_num(v) -> float:
    if isinstance(v, (int, float)):
        return float(v)
    if not v:
        return 0.0
    s = str(v).strip().replace(",", "").replace("%", "")
    # "100억" → 100, "1.2조" → 12000
    mult = 1.0
    if s.endswith("조"): mult, s = 10000.0, s[:-1]
    elif s.endswith("억"): mult, s = 1.0, s[:-1]
    elif s.endswith("천만"): mult, s = 0.1, s[:-2]
    elif s.endswith("만"): mult, s = 0.0001, s[:-1]
    try:
        return float(s) * mult
    except ValueError:
        return 0.0


# ── 텍스트 헬퍼 ──

def _add_text(slide, x, y, w, h, text, *,
              font_size=12, bold=False, color=GREY_900,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    f = run.font
    f.name = FONT_FAMILY
    f.size = Pt(font_size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tb


def _add_rect(slide, x, y, w, h, fill=None, line=None, line_width=0.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width or 0.75)
    return shape


def _draw_chrome(slide, page_no: int, total: int, deck_title: str = ""):
    """슬라이드 공통 chrome — 상단 타이틀 underline, 하단 footer."""
    # 상단 네이비 바 (좌측 짧은 강조)
    _add_rect(slide, 0.5, 0.55, 0.4, 0.05, fill=NAVY)
    # 하단 footer 라인
    _add_rect(slide, 0.5, 7.05, 9.0, 0.01, fill=GREY_300)
    # footer: 좌측 deck title, 우측 page
    if deck_title:
        _add_text(slide, 0.5, 7.12, 6.0, 0.3, deck_title,
                  font_size=8, color=GREY_500, align=PP_ALIGN.LEFT)
    _add_text(slide, 7.5, 7.12, 2.0, 0.3, f"{page_no} / {total}",
              font_size=8, color=GREY_500, align=PP_ALIGN.RIGHT)


def _draw_title(slide, title: str, subtitle: Optional[str] = None):
    """상단 타이틀 영역 (y=0.4 ~ 1.05)."""
    _add_text(slide, 0.5, 0.2, 9.0, 0.4, title,
              font_size=20, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_text(slide, 0.5, 0.65, 9.0, 0.3, subtitle,
                  font_size=10, color=GREY_500, align=PP_ALIGN.LEFT)


def _draw_summary(slide, summary: str, x=0.5, y=1.05, w=9.0, h=0.55):
    """핵심 메시지 박스 — 좌측 네이비 강조 바 + 회색 배경."""
    _add_rect(slide, x, y, 0.06, h, fill=NAVY)
    _add_rect(slide, x + 0.06, y, w - 0.06, h, fill=GREY_100)
    _add_text(slide, x + 0.2, y + 0.05, w - 0.3, h - 0.1, summary,
              font_size=11, bold=True, color=GREY_900,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def _draw_sources(slide, sources: List[str]):
    if not sources:
        return
    text = "출처: " + " · ".join(s.strip() for s in sources if s and s.strip())[:200]
    _add_text(slide, 0.5, 6.75, 9.0, 0.25, text,
              font_size=8, italic=True, color=GREY_500, align=PP_ALIGN.LEFT)


# ── 표 ──

def _draw_table(slide, table_data: Dict[str, Any], x: float, y: float, w: float, h: float):
    headers = table_data.get("headers") or []
    rows = table_data.get("rows") or []
    title = table_data.get("title")

    cur_y = y
    if title:
        _add_text(slide, x, cur_y, w, 0.25, title,
                  font_size=10, bold=True, color=GREY_700, align=PP_ALIGN.LEFT)
        cur_y += 0.28
        h -= 0.28

    if not rows:
        _add_text(slide, x, cur_y, w, 0.4, "(데이터 없음)",
                  font_size=9, italic=True, color=GREY_500, align=PP_ALIGN.CENTER)
        return

    n_cols = max(len(headers), max((len(r) for r in rows), default=1))
    n_rows = len(rows) + (1 if headers else 0)
    if n_cols == 0 or n_rows == 0:
        return

    # python-pptx native table
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(cur_y), Inches(w), Inches(h))
    tbl = tbl_shape.table

    # 컬럼 폭 균등 (첫 컬럼만 좀 넓게)
    if n_cols > 1:
        first = w * 0.35
        rest = (w - first) / (n_cols - 1)
        tbl.columns[0].width = Inches(first)
        for i in range(1, n_cols):
            tbl.columns[i].width = Inches(rest)

    row_offset = 0
    if headers:
        for j in range(n_cols):
            cell = tbl.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            txt = headers[j] if j < len(headers) else ""
            _set_cell_text(cell, str(txt), font_size=9, bold=True, color=WHITE,
                           align=PP_ALIGN.CENTER)
        row_offset = 1

    for ri, row in enumerate(rows):
        for j in range(n_cols):
            cell = tbl.cell(ri + row_offset, j)
            val = row[j] if j < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else GREY_100
            _set_cell_text(cell, str(val), font_size=9, color=GREY_900,
                           align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)


def _set_cell_text(cell, text, *, font_size=9, bold=False, color=GREY_900, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = FONT_FAMILY
    f.size = Pt(font_size)
    f.bold = bold
    f.color.rgb = color


# ── Bullet ──

def _draw_bullets(slide, sentences: List[str], x: float, y: float, w: float, h: float,
                  *, font_size=11, max_lines=10):
    items = [s for s in (sentences or []) if s and s.strip()][:max_lines]
    if not items:
        return
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, s in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.text = ""
        # bullet circle
        b = p.add_run()
        b.text = "■  "
        b.font.name = FONT_FAMILY
        b.font.size = Pt(font_size - 1)
        b.font.color.rgb = NAVY
        # body
        run = p.add_run()
        run.text = str(s).strip()
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.color.rgb = GREY_900
        p.space_after = Pt(4)


# ── 차트 ──

def _draw_chart(slide, chart: Dict[str, Any], x: float, y: float, w: float, h: float):
    title = chart.get("title")
    cur_y = y
    if title:
        _add_text(slide, x, cur_y, w, 0.25, title,
                  font_size=10, bold=True, color=GREY_700, align=PP_ALIGN.LEFT)
        cur_y += 0.28
        h -= 0.28

    png = _make_chart_png(chart, width_in=w, height_in=max(h, 1.0))
    if png:
        slide.shapes.add_picture(io.BytesIO(png), Inches(x), Inches(cur_y),
                                 width=Inches(w), height=Inches(h))
    else:
        _add_text(slide, x, cur_y, w, h, "(차트 데이터 부족)",
                  font_size=9, italic=True, color=GREY_500,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── Layout matcher ──
# 콘텐츠 영역: y = 1.7 ~ 6.65 (h = 4.95), x = 0.5 ~ 9.5 (w = 9.0)
CONTENT_Y = 1.7
CONTENT_H = 4.95
CONTENT_X = 0.5
CONTENT_W = 9.0


def _layout_and_draw(slide, blocks: Dict[str, Any]):
    """
    blocks: {sentences, tables, charts}
    Layout 패턴 (블록 수 기준):
      L1. text only (sentences) — full bullets
      L2. table only — full table
      L3. chart only — full chart
      L4. text + table — left bullets / right table
      L5. text + chart — left bullets / right chart
      L6. table + chart — left chart / right table
      L7. text + table + chart — top text / bottom 2-col (chart + table)
      L8. 2 charts — left / right
      L9. 2 tables — top / bottom
      L10. fallback — 모든 것을 위→아래로 쌓기
    """
    sentences = blocks.get("sentences") or []
    tables = blocks.get("tables") or []
    charts = blocks.get("charts") or []
    n_t = len(tables)
    n_c = len(charts)
    has_s = bool(sentences)

    # L1 / L2 / L3
    if has_s and n_t == 0 and n_c == 0:
        _draw_bullets(slide, sentences, CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H, max_lines=12)
        return
    if not has_s and n_t == 1 and n_c == 0:
        _draw_table(slide, tables[0], CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H)
        return
    if not has_s and n_t == 0 and n_c == 1:
        _draw_chart(slide, charts[0], CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H)
        return

    # L4: text + table (1+1)
    if has_s and n_t == 1 and n_c == 0:
        col_w = (CONTENT_W - 0.3) / 2
        _draw_bullets(slide, sentences, CONTENT_X, CONTENT_Y, col_w, CONTENT_H)
        _draw_table(slide, tables[0], CONTENT_X + col_w + 0.3, CONTENT_Y, col_w, CONTENT_H)
        return
    # L5: text + chart
    if has_s and n_c == 1 and n_t == 0:
        col_w = (CONTENT_W - 0.3) / 2
        _draw_bullets(slide, sentences, CONTENT_X, CONTENT_Y, col_w, CONTENT_H)
        _draw_chart(slide, charts[0], CONTENT_X + col_w + 0.3, CONTENT_Y, col_w, CONTENT_H)
        return
    # L6: chart + table
    if not has_s and n_c == 1 and n_t == 1:
        col_w = (CONTENT_W - 0.3) / 2
        _draw_chart(slide, charts[0], CONTENT_X, CONTENT_Y, col_w, CONTENT_H)
        _draw_table(slide, tables[0], CONTENT_X + col_w + 0.3, CONTENT_Y, col_w, CONTENT_H)
        return
    # L7: text + table + chart (3 blocks)
    if has_s and n_c == 1 and n_t == 1:
        top_h = 1.5
        _draw_bullets(slide, sentences, CONTENT_X, CONTENT_Y, CONTENT_W, top_h, max_lines=4)
        bot_y = CONTENT_Y + top_h + 0.2
        bot_h = CONTENT_H - top_h - 0.2
        col_w = (CONTENT_W - 0.3) / 2
        _draw_chart(slide, charts[0], CONTENT_X, bot_y, col_w, bot_h)
        _draw_table(slide, tables[0], CONTENT_X + col_w + 0.3, bot_y, col_w, bot_h)
        return
    # L8: 2 charts
    if n_c == 2 and not has_s and n_t == 0:
        col_w = (CONTENT_W - 0.3) / 2
        _draw_chart(slide, charts[0], CONTENT_X, CONTENT_Y, col_w, CONTENT_H)
        _draw_chart(slide, charts[1], CONTENT_X + col_w + 0.3, CONTENT_Y, col_w, CONTENT_H)
        return
    # L9: 2 tables
    if n_t == 2 and not has_s and n_c == 0:
        row_h = (CONTENT_H - 0.2) / 2
        _draw_table(slide, tables[0], CONTENT_X, CONTENT_Y, CONTENT_W, row_h)
        _draw_table(slide, tables[1], CONTENT_X, CONTENT_Y + row_h + 0.2, CONTENT_W, row_h)
        return

    # L10: fallback — 알찬 stacking
    cur_y = CONTENT_Y
    remaining = CONTENT_H
    for i, ch in enumerate(charts):
        h = max(2.0, remaining / (len(charts) + len(tables) + (1 if has_s else 0)))
        _draw_chart(slide, ch, CONTENT_X, cur_y, CONTENT_W, h)
        cur_y += h + 0.15
        remaining -= h + 0.15
    if has_s and remaining > 0.5:
        bh = min(2.0, remaining * 0.4)
        _draw_bullets(slide, sentences, CONTENT_X, cur_y, CONTENT_W, bh, max_lines=4)
        cur_y += bh + 0.15
        remaining -= bh + 0.15
    for i, tb in enumerate(tables):
        if remaining <= 0.3:
            break
        h = remaining if i == len(tables) - 1 else max(1.5, remaining / (len(tables) - i))
        _draw_table(slide, tb, CONTENT_X, cur_y, CONTENT_W, h)
        cur_y += h + 0.15
        remaining -= h + 0.15


# ── 슬라이드 단위 빌더 ──

def _build_slide(prs, slide_data: Dict[str, Any], page_no: int, total: int, deck_title: str):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)

    title = (slide_data.get("title") or "").strip() or f"Slide {page_no}"
    subtitle = slide_data.get("subtitle")
    summary = slide_data.get("summary") or slide_data.get("key_message")

    _draw_chrome(slide, page_no, total, deck_title)
    _draw_title(slide, title, subtitle)

    blocks = {
        "sentences": slide_data.get("sentences") or [],
        "tables": slide_data.get("tables") or [],
        "charts": slide_data.get("charts") or [],
    }

    # summary 유무에 따라 콘텐츠 영역 보정
    global CONTENT_Y, CONTENT_H
    prev_y, prev_h = CONTENT_Y, CONTENT_H
    try:
        if summary:
            _draw_summary(slide, summary)
            CONTENT_Y, CONTENT_H = 1.75, 4.9
        else:
            CONTENT_Y, CONTENT_H = 1.1, 5.55
        _layout_and_draw(slide, blocks)
    finally:
        CONTENT_Y, CONTENT_H = prev_y, prev_h

    _draw_sources(slide, slide_data.get("sources") or [])


# ── 표지 / 섹션 구분 ──

def _build_cover(prs, deck_title: str, subtitle: Optional[str] = None,
                 page_no: int = 1, total: int = 1):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # 좌측 네이비 풀 사이드바
    _add_rect(slide, 0, 0, 1.0, 7.5, fill=NAVY)
    # 본문
    _add_text(slide, 1.5, 2.8, 8.0, 1.0, deck_title,
              font_size=32, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_text(slide, 1.5, 3.9, 8.0, 0.5, subtitle,
                  font_size=14, color=GREY_500, align=PP_ALIGN.LEFT)
    # 하단 라인
    _add_rect(slide, 1.5, 5.0, 2.5, 0.05, fill=NAVY_LIGHT)


def _build_divider(prs, section_no: str, section_title: str,
                   page_no: int, total: int, deck_title: str):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_rect(slide, 0, 0, 10.0, 7.5, fill=NAVY)
    _add_text(slide, 0.8, 2.8, 8.5, 1.0, section_no or "",
              font_size=72, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1),
              align=PP_ALIGN.LEFT)
    _add_text(slide, 0.8, 4.0, 8.5, 1.0, section_title or "",
              font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_text(slide, 0.8, 7.05, 8.5, 0.3, f"{page_no} / {total}  ·  {deck_title}",
              font_size=8, color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.LEFT)


# ── Public API ──

def build_pptx(slides: List[Dict[str, Any]], output_path: Optional[str] = None,
               *, deck_title: str = "", deck_subtitle: str = "") -> str:
    """슬라이드 리스트 → 4:3 노앤 PE 톤 .pptx 생성. output_path 반환.

    슬라이드 형식:
    - 'kind': 'cover' | 'divider' | 'content' (기본 content)
    - cover: { kind: 'cover', title, subtitle? }
    - divider: { kind: 'divider', section_number, title }
    - content: 위 모듈 docstring 참조
    """
    prs = Presentation()
    # 4:3 강제
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    if output_path is None:
        fd, output_path = tempfile.mkstemp(suffix=".pptx", prefix="content_")
        os.close(fd)

    total = len(slides)
    rendered = 0
    for i, sl in enumerate(slides):
        page_no = i + 1
        kind = (sl.get("kind") or sl.get("slide_type") or "").lower()
        try:
            if kind in ("cover", "title"):
                _build_cover(prs, sl.get("title") or deck_title, sl.get("subtitle") or deck_subtitle,
                             page_no=page_no, total=total)
            elif kind in ("divider", "section_divider"):
                _build_divider(prs,
                               sl.get("section_number") or sl.get("section_no") or "",
                               sl.get("title") or sl.get("section_title") or "",
                               page_no, total, deck_title)
            else:
                _build_slide(prs, sl, page_no, total, deck_title)
            rendered += 1
        except Exception as e:
            logger.exception(f"[content_ppt] slide {i} failed: {e}")
            # 최후 fallback — 제목만이라도 살림
            try:
                blank = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank)
                _draw_chrome(slide, page_no, total, deck_title)
                _draw_title(slide, sl.get("title") or f"Slide {page_no}")
                _add_text(slide, 0.5, 3.0, 9.0, 1.0,
                          f"[렌더링 실패] {str(e)[:100]}",
                          font_size=11, italic=True, color=GREY_500, align=PP_ALIGN.CENTER)
                rendered += 1
            except Exception:
                pass

    if rendered == 0:
        raise RuntimeError("렌더링된 슬라이드 0장")

    prs.save(output_path)
    logger.info(f"[content_ppt] {rendered}/{len(slides)} 슬라이드 → {output_path}")
    return output_path


def smoke_test() -> str:
    slides = [
        {"kind": "cover", "title": "Bestech IM", "subtitle": "Project Bestech · 2026"},
        {"kind": "divider", "section_number": "I", "title": "Executive Summary"},
        {
            "title": "재무 현황",
            "summary": "최근 3년 매출 CAGR 22%, EBITDA 마진 18% 유지.",
            "sentences": [
                "2024년 매출 100억, EBITDA 18억 (마진 18.0%)",
                "2025년 매출 120억으로 20% 성장, EBITDA 22억",
                "수익성은 동종 PE 포트폴리오 대비 +5pp 우위",
            ],
            "tables": [{
                "title": "Key Financials (단위: 억원)",
                "headers": ["연도", "매출", "EBITDA", "마진"],
                "rows": [
                    ["2023", "82", "13", "15.9%"],
                    ["2024", "100", "18", "18.0%"],
                    ["2025E", "120", "22", "18.3%"],
                ],
            }],
            "charts": [{
                "title": "매출 추이",
                "chart_type": "bar",
                "categories": ["2023", "2024", "2025E"],
                "series": [{"name": "매출", "values": [82, 100, 120]}],
                "y_label": "억원",
            }],
            "sources": ["감사보고서 2024", "Management Projection v3"],
        },
        {
            "title": "시장 점유율",
            "summary": "동남아 시장 진출로 글로벌 점유 8% 달성 가능.",
            "charts": [{
                "title": "지역별 매출 비중 (2025E)",
                "chart_type": "pie",
                "categories": ["국내", "동남아", "북미", "기타"],
                "series": [{"name": "매출", "values": [55, 25, 15, 5]}],
            }],
        },
    ]
    out = build_pptx(slides, deck_title="Bestech IM", deck_subtitle="2026")
    print(f"OK: {out}")
    return out


if __name__ == "__main__":
    import sys as _sys
    if len(_sys.argv) > 1 and _sys.argv[1] == "--smoke":
        smoke_test()
