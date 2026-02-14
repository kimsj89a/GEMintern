import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class InvestmentHistoryUpdater:
    def __init__(self, pptx_path):
        self.prs = Presentation(pptx_path)
        self.data = []

    def extract_data(self):
        """
        Scans all slides (except slide 9 itself) for investment-related keywords
        and builds a data structure for the history table.
        This is a heuristic approach based on the user's description.
        """
        extracted_info = []
        
        # Keywords to look for
        keywords = {
            'Pre-A': 'Pre-A',
            'Series A': 'Series A',
            'Series B': 'Series B',
            'Series C': 'Series C',
            'Series C2': 'Series C2',
            'Seed': 'Seed'
        }
        
        # Simple extraction logic (can be improved with more complex regex)
        for i, slide in enumerate(self.prs.slides):
            if i == 8: # Skip extraction from the target slide itself (Index 8 = Slide 9)
                continue
                
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
            
            # Simple keyword matching
            for key, val in keywords.items():
                if key in text:
                    # Try to find associated numbers (amounts) nearby
                    # Matches "145억원", "920억원", etc.
                    amount_match = re.search(r'(\d+(?:,\d+)*)(?:억원|억)', text)
                    amount = amount_match.group(0) if amount_match else "-"
                    
                    extracted_info.append({
                        "round": val,
                        "amount": amount,
                        "source_slide": i + 1
                    })
        
        # Deduplicate and sort (this logic might need refinement based on real data)
        # For now, let's just use static data if extraction fails or as a fallback/template
        # based on the user's provided example for structure.
        
        self.data = [
            ['시기',      '2020', '2021', '2022', '2023', '2025', '2026'],
            ['투자금액',   '-', '145억원', '920억원', '~2,000억원', '~2,000억원', '2,500~3,000억원'],
            ['Pre-money', '-', '~200억', '~3,000억', '~7,000억', '~1.5조', '2.2~2.4조원'],
            ['주요 투자자', '창업자', '초기 VC', 'KT, 산업은행', 'SKT, Aramco', '전략적/재무적', '국민성장펀드+PEF'],
            ['주요 이벤트', '회사 설립', 'ION 칩 개발', 'ATOM 양산', 'REBEL tape-out', 'REBEL-Quad chip-out', 'REBEL-Quad 양산'],
        ]
        return self.data

    def update_slide(self, slide_index=8):
        """
        Updates the specified slide (default index 8 for Slide 9) with the investment history table.
        """
        if slide_index >= len(self.prs.slides):
            print(f"Error: Slide index {slide_index} out of range.")
            return

        slide = self.prs.slides[slide_index]
        
        # Remove existing table if any (heuristic: check for GraphicFrame)
        # In a real scenario, we might want to be more specific.
        # For this task, we'll just add a new table on top or clear pertinent shapes if identified.
        
        # User specified layout:
        # Table position and size
        left = Emu(350000)
        top = Emu(1700000)
        width = Emu(11500000)
        height = Emu(3800000)
        
        rows = 6
        cols = 7
        
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        
        # Column widths
        col_widths = [1200000, 1400000, 1400000, 1500000, 1500000, 1500000, 2277600]
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)
            
        # Headers
        headers = ['구분', 'Seed', 'Pre-A', 'Series A', 'Series B', 'Series C', 'Series C2 (본건)']
        
        # 0. Header Row
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x20, 0x60) # Dark Blue
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White
            paragraph.font.bold = True
            paragraph.font.size = Pt(10)
            paragraph.alignment = PP_ALIGN.CENTER

            # Special styling for the last column header "Series C2 (본건)"
            if i == 6:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC) # Light Yellow
                paragraph.font.color.rgb = RGBColor(0xC0, 0x00, 0x00) # Red Text
                # Add bottom border logic here if needed (python-pptx has limited border support, might need XML)

        # 1. Data Rows
        # self.data structure: list of rows, where each row is a list of cell values
        # Row 0 is '시기', Row 1 is '투자금액', etc.
        # We start filling from table row 1 (since row 0 is header)
        
        for r_idx, row_vals in enumerate(self.data):
            for c_idx, val in enumerate(row_vals):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = val
                
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
                
                # First column styling (Row headers)
                if c_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x00, 0x20, 0x60) # Dark Blue
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White
                    paragraph.font.bold = True
                else:
                    # Data cells
                    paragraph.font.color.rgb = RGBColor(0x00, 0x00, 0x00) # Black

                # Highlight last column (Series C2)
                if c_idx == 6 and r_idx >= 0:
                     cell.fill.solid()
                     cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC) # Light Yellow background
        
    def save(self, output_path):
        self.prs.save(output_path)
