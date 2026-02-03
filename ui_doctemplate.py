"""
문서양식 복사기 모듈
기존 문서 양식 업로드 -> AI가 형식 분석 -> 콘텐츠 파일 분석 -> 새 문서 생성
"""
import streamlit as st
import io
import os
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from google import genai

# 기본 템플릿 경로 (.docx만 지원, .dotm은 python-docx 미지원)
TEMPLATE_PATH_DOCX = os.path.join(os.path.dirname(__file__), "template", "Normal.docx")

def get_template_path():
    """사용 가능한 템플릿 경로 반환 (.docx 우선)"""
    if os.path.exists(TEMPLATE_PATH_DOCX):
        return TEMPLATE_PATH_DOCX
    return None  # .dotm은 python-docx 미지원

# 지원 파일 형식
SUPPORTED_FILE_TYPES = ['docx', 'pdf', 'pptx', 'xlsx', 'xls', 'txt', 'md']


def extract_text_from_file(uploaded_file) -> str:
    """업로드된 파일에서 텍스트 추출"""
    filename = uploaded_file.name.lower()
    content = ""

    try:
        uploaded_file.seek(0)  # 파일 포인터 리셋

        if filename.endswith('.txt') or filename.endswith('.md'):
            raw_bytes = uploaded_file.read()
            for encoding in ['utf-8-sig', 'utf-8', 'cp949', 'euc-kr']:
                try:
                    content = raw_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            if not content:
                content = raw_bytes.decode('utf-8', errors='replace')

        elif filename.endswith('.docx'):
            doc = Document(io.BytesIO(uploaded_file.read()))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = '\n\n'.join(paragraphs)

        elif filename.endswith('.pdf'):
            try:
                import fitz
                pdf_bytes = uploaded_file.read()
                pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                text_parts = []
                for page in pdf_doc:
                    text_parts.append(page.get_text())
                content = '\n'.join(text_parts)
                pdf_doc.close()
            except ImportError:
                content = "[PDF 읽기 실패: PyMuPDF 필요]"

        elif filename.endswith('.pptx'):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(uploaded_file.read()))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                content = '\n\n'.join(text_parts)
            except ImportError:
                content = "[PPT 읽기 실패: python-pptx 필요]"

        elif filename.endswith('.xlsx') or filename.endswith('.xls'):
            try:
                import pandas as pd
                df = pd.read_excel(io.BytesIO(uploaded_file.read()), sheet_name=None)
                text_parts = []
                for sheet_name, sheet_df in df.items():
                    text_parts.append(f"[{sheet_name}]\n{sheet_df.to_string()}")
                content = '\n\n'.join(text_parts)
            except ImportError:
                content = "[Excel 읽기 실패: pandas 필요]"

    except Exception as e:
        content = f"[파일 읽기 오류: {str(e)}]"

    return content


def analyze_document_format(raw_text: str, api_key: str, model: str) -> str:
    """AI로 문서 형식을 마크다운으로 분석"""
    client = genai.Client(api_key=api_key)

    prompt = """다음 문서의 **형식과 구조**를 분석해서 마크다운 템플릿으로 변환해주세요.

[분석 원칙]
1. 문서의 전체적인 구조(섹션, 제목, 소제목 등)를 파악
2. 각 섹션에 어떤 종류의 내용이 들어가는지 파악 (예: 회사 개요, 재무 현황, 결론 등)
3. 반복되는 패턴이나 표 형식이 있다면 구조 파악
4. 실제 내용은 {{섹션명}} 형태의 플레이스홀더로 대체

[출력 형식]
- 마크다운 형식으로 출력
- 각 섹션의 제목은 그대로 유지
- 내용이 들어갈 부분은 {{해당_섹션_설명}} 형태로 표시
- 문서의 톤앤매너, 형식적 특징도 주석으로 메모

[문서 내용]
""" + raw_text[:15000]  # 토큰 제한

    try:
        response = client.models.generate_content(model=model, contents=prompt)
        return response.text.strip() if response.text else ""
    except Exception as e:
        return f"[형식 분석 오류: {str(e)}]"


def extract_content_as_markdown(raw_text: str, api_key: str, model: str) -> str:
    """AI로 콘텐츠를 마크다운으로 정리"""
    client = genai.Client(api_key=api_key)

    prompt = """다음 문서의 **내용**을 마크다운 형식으로 깔끔하게 정리해주세요.

[정리 원칙]
1. 핵심 정보와 데이터를 빠짐없이 추출
2. 논리적 구조로 재정리 (제목, 소제목, 불릿 포인트 활용)
3. 숫자, 고유명사, 날짜 등은 정확히 유지
4. 불필요한 반복이나 형식적 문구는 제거
5. 표 형식 데이터는 마크다운 테이블로 변환

[문서 내용]
""" + raw_text[:15000]

    try:
        response = client.models.generate_content(model=model, contents=prompt)
        return response.text.strip() if response.text else ""
    except Exception as e:
        return f"[콘텐츠 추출 오류: {str(e)}]"


def generate_final_document(format_md: str, content_md: str, api_key: str, model: str) -> str:
    """형식 템플릿 + 콘텐츠 = 최종 문서 생성"""
    client = genai.Client(api_key=api_key)

    prompt = f"""다음 두 가지를 결합하여 완성된 문서를 마크다운으로 작성해주세요.

[문서 형식 템플릿]
{format_md}

---

[채워 넣을 콘텐츠]
{content_md}

---

[작성 원칙]
1. 형식 템플릿의 구조와 스타일을 따름
2. 콘텐츠의 정보를 적절한 섹션에 배치
3. 플레이스홀더({{...}})를 실제 내용으로 채움
4. 콘텐츠에 없는 정보는 추가하지 않음
5. 자연스럽고 전문적인 문체 유지
6. 최종 결과는 마크다운 형식으로 출력
"""

    try:
        response = client.models.generate_content(model=model, contents=prompt)
        return response.text.strip() if response.text else ""
    except Exception as e:
        return f"[문서 생성 오류: {str(e)}]"


def markdown_to_docx(markdown_text: str, use_template: bool = True) -> bytes:
    """마크다운을 Word 문서로 변환"""
    import re

    template_path = get_template_path()
    if use_template and template_path:
        doc = Document(template_path)
        for element in doc.element.body[:]:
            doc.element.body.remove(element)
    else:
        doc = Document()
        style = doc.styles['Normal']
        style.font.name = '맑은 고딕'
        style.font.size = Pt(11)

    lines = markdown_text.split('\n')
    i = 0
    in_code_block = False
    code_content = []

    while i < len(lines):
        line = lines[i]

        # 코드 블록
        if line.strip().startswith('```'):
            if in_code_block:
                code_text = '\n'.join(code_content)
                p = doc.add_paragraph()
                run = p.add_run(code_text)
                run.font.name = 'Consolas'
                run.font.size = Pt(10)
                code_content = []
                in_code_block = False
            else:
                in_code_block = True
            i += 1
            continue

        if in_code_block:
            code_content.append(line)
            i += 1
            continue

        # 빈 줄
        if not line.strip():
            i += 1
            continue

        # 헤더
        if line.startswith('### '):
            doc.add_paragraph(line[4:].strip(), style='Heading 3')
        elif line.startswith('## '):
            doc.add_paragraph(line[3:].strip(), style='Heading 2')
        elif line.startswith('# '):
            doc.add_paragraph(line[2:].strip(), style='Heading 1')
        # 리스트
        elif line.strip().startswith('- ') or line.strip().startswith('* '):
            text = line.strip()[2:]
            doc.add_paragraph(text, style='List Bullet')
        elif re.match(r'^\s*\d+\.\s', line):
            text = re.sub(r'^\s*\d+\.\s', '', line)
            doc.add_paragraph(text, style='List Number')
        # 수평선
        elif line.strip() in ['---', '***', '___']:
            p = doc.add_paragraph('_' * 50)
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        # 일반 텍스트
        else:
            # 마크다운 포맷 제거
            clean_text = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
            clean_text = re.sub(r'\*(.+?)\*', r'\1', clean_text)
            clean_text = re.sub(r'`(.+?)`', r'\1', clean_text)
            doc.add_paragraph(clean_text)

        i += 1

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def render_doctemplate_panel(settings):
    """문서양식 복사기 UI 패널"""
    st.markdown("### 📋 문서양식 복사기")

    st.markdown("""
        <div style='background-color: #f0f2f6; padding: 15px; border-radius: 8px; margin-bottom: 20px; border-left: 4px solid #0068c9;'>
        <h4 style='margin-top: 0; color: #0068c9;'>📋 기능 안내</h4>
        <b>기존 문서의 형식을 분석하여 새로운 내용으로 문서를 생성합니다.</b><br/><br/>
        <b>사용 방법:</b><br/>
        1️⃣ 양식 문서 업로드 (형식/구조 참고용)<br/>
        2️⃣ AI가 문서 형식을 마크다운으로 분석<br/>
        3️⃣ 콘텐츠 파일 업로드 (새로 넣을 내용)<br/>
        4️⃣ AI가 형식 + 콘텐츠를 결합하여 새 문서 생성
        </div>
    """, unsafe_allow_html=True)

    # API Key 확인
    api_key = settings.get('api_key', '') if settings else ''
    if not api_key:
        st.warning("⚠️ 상단 설정에서 Google API Key를 입력해주세요.")
        return

    # 모델 선택
    model = st.selectbox(
        "🤖 AI 모델",
        ["gemini-3-flash-preview", "gemini-3-pro-preview"],
        key="doctemplate_model"
    )

    # =========================================
    # 1단계: 양식 문서 업로드
    # =========================================
    st.markdown("---")
    st.markdown("## 1️⃣ 양식 문서 업로드")
    st.caption("형식과 구조를 참고할 문서를 업로드하세요.")

    format_file = st.file_uploader(
        "양식 문서",
        type=SUPPORTED_FILE_TYPES,
        key="doctemplate_format_file"
    )

    if format_file:
        format_raw = extract_text_from_file(format_file)
        st.success(f"✅ 양식 파일 로드: {format_file.name} ({len(format_raw):,}자)")

        with st.expander("📄 원본 내용 미리보기", expanded=False):
            st.text(format_raw[:2000] + ("..." if len(format_raw) > 2000 else ""))

        # =========================================
        # 2단계: AI 형식 분석
        # =========================================
        st.markdown("---")
        st.markdown("## 2️⃣ 문서 형식 분석")

        if st.button("🔍 AI로 형식 분석", type="secondary", key="btn_analyze_format"):
            with st.spinner("문서 형식 분석 중..."):
                format_md = analyze_document_format(format_raw, api_key, model)
                st.session_state['doctemplate_format_md'] = format_md
                st.success("✅ 형식 분석 완료!")

        if st.session_state.get('doctemplate_format_md'):
            st.markdown("#### 📝 분석된 문서 형식 (마크다운)")
            format_md = st.text_area(
                "형식 템플릿 (편집 가능)",
                value=st.session_state['doctemplate_format_md'],
                height=300,
                key="doctemplate_format_edit"
            )
            st.session_state['doctemplate_format_md'] = format_md

            # =========================================
            # 3단계: 콘텐츠 파일 업로드
            # =========================================
            st.markdown("---")
            st.markdown("## 3️⃣ 콘텐츠 파일 업로드")
            st.caption("새 문서에 들어갈 내용이 담긴 파일을 업로드하세요.")

            content_file = st.file_uploader(
                "콘텐츠 파일",
                type=SUPPORTED_FILE_TYPES,
                key="doctemplate_content_file"
            )

            if content_file:
                content_raw = extract_text_from_file(content_file)
                st.success(f"✅ 콘텐츠 파일 로드: {content_file.name} ({len(content_raw):,}자)")
                st.session_state['doctemplate_content_filename'] = content_file.name

                with st.expander("📄 콘텐츠 미리보기", expanded=False):
                    st.text(content_raw[:2000] + ("..." if len(content_raw) > 2000 else ""))

                if st.button("🔍 AI로 콘텐츠 정리", type="secondary", key="btn_analyze_content"):
                    with st.spinner("콘텐츠 정리 중..."):
                        content_md = extract_content_as_markdown(content_raw, api_key, model)
                        st.session_state['doctemplate_content_md'] = content_md
                        st.success("✅ 콘텐츠 정리 완료!")

                if st.session_state.get('doctemplate_content_md'):
                    st.markdown("#### 📝 정리된 콘텐츠 (마크다운)")
                    content_md = st.text_area(
                        "콘텐츠 (편집 가능)",
                        value=st.session_state['doctemplate_content_md'],
                        height=300,
                        key="doctemplate_content_edit"
                    )
                    st.session_state['doctemplate_content_md'] = content_md

                    # =========================================
                    # 4단계: 문서 생성
                    # =========================================
                    st.markdown("---")
                    st.markdown("## 4️⃣ 새 문서 생성")

                    # 출력 파일명 (콘텐츠 파일명 기반)
                    default_filename = os.path.splitext(
                        st.session_state.get('doctemplate_content_filename', 'output')
                    )[0]

                    col1, col2 = st.columns(2)
                    with col1:
                        output_filename = st.text_input(
                            "출력 파일명",
                            value=default_filename,
                            key="doctemplate_output_filename"
                        )
                    with col2:
                        template_path = get_template_path()
                        template_exists = template_path is not None
                        use_template = st.checkbox(
                            "기본 템플릿 스타일 적용",
                            value=template_exists,
                            disabled=not template_exists,
                            help="template/Normal.docx 파일이 필요합니다." if not template_exists else None,
                            key="doctemplate_use_template"
                        )

                    if st.button("🚀 문서 생성", type="primary", use_container_width=True, key="btn_generate"):
                        with st.spinner("AI가 문서를 생성하는 중..."):
                            final_md = generate_final_document(
                                st.session_state['doctemplate_format_md'],
                                st.session_state['doctemplate_content_md'],
                                api_key,
                                model
                            )
                            st.session_state['doctemplate_final_md'] = final_md
                            st.success("✅ 문서 생성 완료!")

                    if st.session_state.get('doctemplate_final_md'):
                        st.markdown("#### 📄 생성된 문서")

                        with st.container(border=True):
                            st.markdown(st.session_state['doctemplate_final_md'])

                        # 다운로드 옵션
                        st.markdown("#### 💾 다운로드")
                        col1, col2 = st.columns(2)

                        with col1:
                            st.download_button(
                                "📥 마크다운 (.md)",
                                st.session_state['doctemplate_final_md'],
                                f"{output_filename}.md",
                                "text/markdown",
                                use_container_width=True
                            )

                        with col2:
                            try:
                                docx_bytes = markdown_to_docx(
                                    st.session_state['doctemplate_final_md'],
                                    use_template=use_template
                                )
                                st.download_button(
                                    "📥 Word (.docx)",
                                    docx_bytes,
                                    f"{output_filename}.docx",
                                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                    use_container_width=True,
                                    type="primary"
                                )
                            except Exception as e:
                                st.error(f"Word 변환 오류: {e}")

    else:
        st.info("양식 문서를 업로드하면 AI가 형식을 분석합니다.")
