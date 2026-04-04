"""FastAPI REST API routes for GEMintern."""
import io
import json
import logging
import os
import uuid
import tempfile
from typing import List, Dict

from fastapi import APIRouter, UploadFile, File, Form, Depends, HTTPException
from fastapi.responses import Response

from pydantic import BaseModel
from backend.api_models import (
    ProjectCreate, ProjectRename, FolderCreate, FolderRename, DocMoveRequest, NoteCreate, NoteUpdate,
    GenerateRequest, QaRequest, AnalysisRequest,
    CreatePptxRequest, SlideRegenerateRequest, SlidesFromOutlineRequest,
    SaveResearchRequest,
    FolderScanRequest, FolderScanFileInfo, FolderScanPreviewResponse,
    FolderIngestRequest,
    WikiSectionUpdate, WikiSectionCreate,
)
from backend.api_ws import create_task, run_generate_task, run_analysis_task, get_task
from backend.auth import get_current_user
from backend.database import log_usage, save_generation
from ai_client import AIClient

logger = logging.getLogger(__name__)

router = APIRouter()


def _save_inline_task_history(task: dict, result_text: str):
    """inline _run() 완료 시 generation_history에 저장."""
    user_id = task.get("_user_id")
    if not user_id:
        return
    try:
        save_generation(
            user_id=user_id,
            endpoint=task.get("_endpoint", ""),
            title=task.get("_title", ""),
            model=task.get("_model"),
            inputs=task.get("_inputs"),
            result_text=result_text,
        )
    except Exception as e:
        print(f"[history] save error: {e}")


def _get_config(settings_key: str, env_var: str = None) -> str:
    """Generic config getter: env var first, then settings.json fallback."""
    if env_var:
        env_val = os.environ.get(env_var, "")
        if env_val:
            return env_val
    return _load_settings().get(settings_key, "")


def _get_api_key() -> str:
    """Return Gemini API key: env var first, then settings.json fallback."""
    return _get_config("api_key", "GEMINI_API_KEY")


def _get_model_name() -> str:
    """Return model name: env var first, then settings.json fallback."""
    return _get_config("model_name", "MODEL_NAME") or "gemini-2.5-flash"


def _get_anthropic_api_key() -> str:
    """Return Anthropic API key: env var first, then settings.json fallback."""
    return _get_config("anthropic_api_key", "ANTHROPIC_API_KEY")

# Max context size per model family
MAX_CONTEXT_CHARS_GEMINI = 800_000   # ~200K tokens
MAX_CONTEXT_CHARS_CLAUDE = 50_000    # ~25K tokens (Claude 분당 30K 토큰 제한 대응)
MAX_CONTEXT_CHARS = MAX_CONTEXT_CHARS_GEMINI  # default


def _max_chars_for_model(model: str) -> int:
    """모델에 따라 최대 컨텍스트 크기 반환."""
    if model.startswith("claude-"):
        return MAX_CONTEXT_CHARS_CLAUDE
    return MAX_CONTEXT_CHARS_GEMINI


def _strip_doc_stem(filename: str) -> str:
    """Strip .md wrapper and original extension to get bare stem.
    e.g. '신한벤치PE.txt.md' -> '신한벤치PE', 'report.pdf' -> 'report'
    """
    import os as _os
    name = filename
    if name.endswith('.md'):
        name = name[:-3]
    return _os.path.splitext(name)[0]


def _load_context_with_budget(project_name: str, selected_docs: list = None,
                               max_chars: int = MAX_CONTEXT_CHARS, owner_id: int | None = None) -> str:
    """Load project docs with per-document truncation and clear headers."""
    import core_rag
    docs_dict = core_rag.load_project_docs_dict(project_name, owner_id=owner_id)

    if selected_docs:
        sel_stems = {_strip_doc_stem(s) for s in selected_docs}
        docs_dict = {k: v for k, v in docs_dict.items()
                     if _strip_doc_stem(k) in sel_stems or k in selected_docs}

    if not docs_dict:
        return ""

    n = len(docs_dict)
    per_doc = max_chars // n
    parts = []
    for name, content in sorted(docs_dict.items()):
        doc_name = name.replace('.md', '')
        if len(content) > per_doc:
            content = content[:per_doc] + f"\n\n[... '{doc_name}' 문서 일부 생략 ({len(content):,}자 중 {per_doc:,}자)]"
        parts.append(f"===== 문서: {doc_name} =====\n{content}")

    return '\n\n'.join(parts)


def _truncate_context(text: str, max_chars: int = MAX_CONTEXT_CHARS) -> str:
    """Simple truncation fallback for pre-loaded text."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n\n[... 이하 생략 — 전체 자료의 일부만 제공됨]"


def _select_relevant_docs(project_name: str, query: str, model: str = "",
                           selected_docs: list = None, owner_id: int | None = None) -> str:
    """질문과 관련된 문서만 선별하여 컨텍스트 구성.
    1순위: 하이브리드 검색 (BM25 + 벡터 + RRF + 리랭킹)
    2순위: 기존 벡터 검색 폴백
    3순위: 키워드 매칭 + 예산 분배
    4순위: 전체 문서 로드
    """
    max_chars = _max_chars_for_model(model)
    api_key = _get_api_key()

    # 1. 하이브리드 검색
    if query:
        try:
            import core_rag_hybrid
            hybrid_ctx = core_rag_hybrid.search_and_build_context(
                api_key, project_name, query,
                top_k=15, max_chars=max_chars,
                selected_docs=selected_docs,
            )
            if hybrid_ctx and len(hybrid_ctx.strip()) > 100:
                logger.info(f"RAG: hybrid search OK for '{project_name}' (selected={len(selected_docs) if selected_docs else 'all'})")
                return _truncate_context(hybrid_ctx, max_chars)
        except Exception as e:
            logger.debug(f"RAG: hybrid search failed, falling back: {e}")

    # 2. 기존 벡터 검색 폴백
    if query:
        vector_ctx = _get_vector_context(api_key, project_name, query, selected_docs, owner_id=owner_id)
        if vector_ctx and len(vector_ctx.strip()) > 100:
            return _truncate_context(vector_ctx, max_chars)

    # 3. 문서 로드 + 키워드 스코어링
    import core_rag
    docs_dict = core_rag.load_project_docs_dict(project_name, owner_id=owner_id)
    if not docs_dict:
        return ""

    # 선택된 문서만 필터링
    if selected_docs:
        sel_stems = {_strip_doc_stem(s) for s in selected_docs}
        docs_dict = {k: v for k, v in docs_dict.items()
                     if _strip_doc_stem(k) in sel_stems or k in selected_docs}

    if not docs_dict:
        return ""

    if query and len(docs_dict) > 3:
        import re as _re
        query_lower = query.lower()
        keywords = set(_re.findall(r'[\w가-힣]{2,}', query_lower))

        scored = []
        for name, content in docs_dict.items():
            doc_lower = (name + " " + content[:2000]).lower()
            score = sum(1 for kw in keywords if kw in doc_lower)
            name_lower = name.lower()
            score += sum(3 for kw in keywords if kw in name_lower)
            scored.append((name, content, score))

        scored.sort(key=lambda x: -x[2])
        selected = []
        total = 0
        for name, content, score in scored:
            if total + min(len(content), max_chars // max(len(scored), 1)) > max_chars and len(selected) >= 1:
                break
            selected.append((name, content))
            total += len(content)
        docs_dict = dict(selected)

    return _load_context_with_budget_from_dict(docs_dict, max_chars)


def _load_context_with_budget_from_dict(docs_dict: dict, max_chars: int) -> str:
    """문서 dict에서 예산 분배하여 컨텍스트 구성."""
    if not docs_dict:
        return ""
    n = len(docs_dict)
    per_doc = max_chars // n
    parts = []
    for name, content in sorted(docs_dict.items()):
        doc_name = name.replace('.md', '')
        if len(content) > per_doc:
            content = content[:per_doc] + f"\n\n[... '{doc_name}' 일부 생략]"
        parts.append(f"===== 문서: {doc_name} =====\n{content}")
    return '\n\n'.join(parts)


def _get_vector_context(api_key: str, project_name: str, query: str,
                        selected_docs: list = None, owner_id: int | None = None) -> str:
    """Vector 검색으로 관련 청크 추출. ChromaDB segfault 방지를 위해 서브프로세스에서 실행."""
    import subprocess, sys, json as _json
    try:
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        vdb_dir = os.path.join(base, "rag_storage", "_vectordb")
        if not os.path.isdir(vdb_dir):
            raise FileNotFoundError("no vectordb")

        # 서브프로세스에서 ChromaDB 검색 실행 (segfault 격리)
        script = f"""
import sys, json
sys.path.insert(0, {repr(base)})
import core_rag_vector
result = core_rag_vector.build_context_from_search(
    {repr(api_key)}, {repr(project_name)}, {repr(query)},
    selected_docs={repr(selected_docs)})
if result:
    print("__VECTOR_OK__")
    print(result)
else:
    print("__VECTOR_EMPTY__")
"""
        proc = subprocess.run(
            [sys.executable, "-c", script],
            capture_output=True, text=True, timeout=10, cwd=base,
        )
        if proc.returncode == 0 and "__VECTOR_OK__" in proc.stdout:
            context = proc.stdout.split("__VECTOR_OK__\n", 1)[1]
            if context.strip():
                return context
    except Exception:
        logger.debug("Vector search subprocess failed, falling back to full doc load")
    # Fallback: 전체 문서 로드 + 예산 분배
    return _load_context_with_budget(project_name, selected_docs, owner_id=owner_id)

# Settings file path
SETTINGS_FILE = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "settings.json"
)


def _load_settings() -> dict:
    if os.path.exists(SETTINGS_FILE):
        with open(SETTINGS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def _load_settings_for_user(user_id: int) -> dict:
    """settings.json + 유저별 DB 설정 오버레이. 유저가 바꾼 model_name 등이 우선."""
    from backend.database import get_user_settings
    base = _load_settings()
    user_settings = get_user_settings(user_id)
    if user_settings:
        base.update(user_settings)
    return base


def _save_settings(data: dict):
    with open(SETTINGS_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)




def _verify_project_ownership(project_name: str, user_id: int):
    """Raise 403 if user doesn't own the project. Check SQLite first, fallback to rag_storage."""
    from backend.database import get_db
    with get_db() as conn:
        project = conn.execute(
            "SELECT id, owner_id FROM projects WHERE name = ? AND owner_id = ?",
            (project_name, user_id)
        ).fetchone()
    if project:
        return  # OK
    # Fallback: check rag_storage for legacy compatibility
    import core_rag
    projects = core_rag._load_projects()
    p = next((p for p in projects if p["name"] == project_name), None)
    if not p:
        raise HTTPException(status_code=404, detail="프로젝트를 찾을 수 없습니다.")
    if p.get("owner_id") is not None and p.get("owner_id") != user_id:
        raise HTTPException(status_code=403, detail="접근 권한이 없습니다.")

# ========================================
# Health
# ========================================

@router.get("/health")
def health_check():
    return {"status": "ok", "version": "2026.03.21"}


# ========================================
# Settings
# ========================================

@router.get("/settings")
def get_settings(user: dict = Depends(get_current_user)):
    from backend.database import get_user_settings
    data = _load_settings()
    masked = {**data}
    # Never expose full API keys
    masked.pop("api_key", None)
    masked.pop("anthropic_api_key", None)
    masked["api_key_configured"] = bool(_get_api_key())
    masked["anthropic_api_key_configured"] = bool(_get_anthropic_api_key())
    # Overlay per-user settings
    user_settings = get_user_settings(user["id"])
    masked.update(user_settings)
    return masked


@router.put("/settings")
def update_settings(settings: dict, user: dict = Depends(get_current_user)):
    from backend.database import save_user_settings, get_user_settings
    settings.pop("api_key", None)  # API key managed via env var
    settings.pop("anthropic_api_key", None)  # Anthropic key managed via env var
    # Save per-user settings to DB
    current_user_settings = get_user_settings(user["id"])
    current_user_settings.update(settings)
    save_user_settings(user["id"], current_user_settings)
    # Also sync model_name/thinking_level to settings.json so all endpoints pick it up
    global_keys = ("model_name", "thinking_level")
    global_settings = _load_settings()
    changed = False
    for k in global_keys:
        if k in settings and settings[k] != global_settings.get(k):
            global_settings[k] = settings[k]
            changed = True
    if changed:
        _save_settings(global_settings)
    return {"success": True}


@router.post("/settings/apply")
def apply_settings(user: dict = Depends(get_current_user)):
    settings = _load_settings_for_user(user["id"])
    model = settings.get("model_name", "")

    # Claude 모델 선택 시 Anthropic key 검증
    if model.startswith("claude-"):
        anthropic_key = _get_anthropic_api_key()
        if not anthropic_key:
            return {"success": False, "error": "Anthropic API Key가 설정되지 않았습니다. (.env 파일의 ANTHROPIC_API_KEY 확인)"}
        try:
            import anthropic
            client = anthropic.Anthropic(api_key=anthropic_key)
            # 간단한 검증 호출
            client.messages.create(
                model=model, max_tokens=10,
                messages=[{"role": "user", "content": "hi"}],
            )
            return {"success": True}
        except Exception as e:
            return {"success": False, "error": f"Anthropic API 검증 실패: {str(e)}"}

    # Gemini 모델
    api_key = _get_api_key()
    if not api_key:
        return {"success": False, "error": "API Key가 설정되지 않았습니다. (GEMINI_API_KEY 환경변수 확인)"}
    try:
        import core_logic
        core_logic.get_client(api_key)
        return {"success": True}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ========================================
# Projects
# ========================================

@router.get("/projects")
def list_projects(user: dict = Depends(get_current_user)):
    from backend.database import get_db
    import core_rag
    with get_db() as conn:
        rows = conn.execute(
            "SELECT id, name, storage_name, created_at FROM projects WHERE owner_id = ? ORDER BY created_at DESC",
            (user["id"],)
        ).fetchall()
    result = []
    for r in rows:
        doc_count = 0
        try:
            doc_count = len(core_rag.get_indexed_doc_names(r["name"], owner_id=user["id"]))
        except Exception:
            logger.debug("Failed to get indexed doc count for project %s", r["name"])
        # Fallback: check SQLite documents table (for Railway ephemeral FS)
        if doc_count == 0:
            with get_db() as conn2:
                cnt = conn2.execute(
                    "SELECT COUNT(*) as cnt FROM documents WHERE project_id = ?",
                    (r["id"],)
                ).fetchone()
                if cnt:
                    doc_count = cnt["cnt"]
        result.append({
            "id": r["id"], "name": r["name"], "storage_name": r["storage_name"],
            "created_at": r["created_at"], "doc_count": doc_count,
        })
    return result


@router.post("/projects")
def create_project(req: ProjectCreate, user: dict = Depends(get_current_user)):
    import core_rag
    from backend.database import get_db
    # Create in rag_storage (files + index)
    rag_result = core_rag.create_project(req.name, owner_id=user["id"])
    if not rag_result.get("success"):
        return rag_result
    # Also insert into SQLite
    project_info = rag_result.get("project", {})
    with get_db() as conn:
        existing = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (req.name.strip(), user["id"])
        ).fetchone()
        if not existing:
            conn.execute(
                "INSERT INTO projects (name, owner_id, storage_name) VALUES (?, ?, ?)",
                (project_info.get("name", req.name.strip()), user["id"],
                 project_info.get("storage_name", req.name.strip()))
            )
    return rag_result


@router.patch("/projects/{name}")
def rename_project(name: str, req: ProjectRename, user: dict = Depends(get_current_user)):
    """Rename a project (display name only, storage directory unchanged)."""
    import re
    import core_rag
    from backend.database import get_db
    new_name = req.new_name.strip()
    if not new_name:
        raise HTTPException(status_code=400, detail="프로젝트명을 입력해주세요.")
    safe_name = re.sub(r'[\\/*?:"<>|]', "", new_name).strip()
    if not safe_name:
        raise HTTPException(status_code=400, detail="유효하지 않은 프로젝트명입니다.")

    with get_db() as conn:
        # Check ownership
        project = conn.execute(
            "SELECT id, storage_name FROM projects WHERE name = ? AND owner_id = ?",
            (name, user["id"])
        ).fetchone()
        if not project:
            raise HTTPException(status_code=404, detail="프로젝트를 찾을 수 없습니다.")
        # Check duplicate
        dup = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (safe_name, user["id"])
        ).fetchone()
        if dup and dup["id"] != project["id"]:
            raise HTTPException(status_code=409, detail=f"'{safe_name}' 프로젝트가 이미 존재합니다.")
        # Update DB
        conn.execute(
            "UPDATE projects SET name = ? WHERE id = ?",
            (safe_name, project["id"])
        )
    # Update _projects.json for legacy compatibility
    projects = core_rag._load_projects()
    for p in projects:
        if p["name"] == name and (p.get("owner_id") is None or p.get("owner_id") == user["id"]):
            p["name"] = safe_name
            break
    core_rag._save_projects(projects)
    return {"success": True, "name": safe_name}


@router.delete("/projects/{name}")
def delete_project(name: str, user: dict = Depends(get_current_user)):
    import core_rag
    from backend.database import get_db
    # Delete from SQLite
    with get_db() as conn:
        conn.execute(
            "DELETE FROM projects WHERE name = ? AND owner_id = ?",
            (name, user["id"])
        )
    # Delete from rag_storage
    return core_rag.delete_project(name, owner_id=user["id"])


# ========================================
# Documents & Folders
# ========================================

@router.get("/projects/{name}/docs")
def get_project_docs(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_rag
    from backend.database import get_db

    # 1. 파일시스템에서 문서 목록
    fs_names = set(core_rag.get_indexed_doc_names(name, owner_id=user["id"]) or [])
    fs_tree = core_rag.get_folder_tree(name, owner_id=user["id"]) or {}

    # 2. SQLite에서 문서 목록 (항상 확인 — 파일시스템과 합침)
    # Skip SQLite entries whose stem already exists in fs (avoids stem/filename duplicates)
    db_tree: dict = {}
    db_names: set = set()
    with get_db() as conn:
        project = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (name, user["id"])
        ).fetchone()
        if project:
            rows = conn.execute(
                "SELECT DISTINCT filename, folder FROM documents WHERE project_id = ?",
                (project["id"],)
            ).fetchall()
            for r in rows:
                fname = r["filename"]
                stem = os.path.splitext(fname)[0]
                # If filesystem already tracks this doc (as stem), skip the SQLite entry
                if stem in fs_names or fname in fs_names:
                    continue
                folder = r["folder"] or core_rag.ROOT_FOLDER
                db_tree.setdefault(folder, []).append(fname)
                db_names.add(fname)

    # 3. 합치기 (파일시스템 우선, SQLite는 보완)
    all_names = fs_names | db_names
    merged_tree: dict = {}
    for folder, docs in {**db_tree, **fs_tree}.items():
        merged_tree.setdefault(folder, [])
        for doc in docs:
            if doc not in merged_tree[folder]:
                merged_tree[folder].append(doc)

    # 4. 어떤 폴더에도 없는 문서를 __root__에 추가 (orphan 방지)
    assigned = set()
    for docs in merged_tree.values():
        assigned.update(docs)
    orphans = all_names - assigned
    if orphans:
        merged_tree.setdefault(core_rag.ROOT_FOLDER, [])
        merged_tree[core_rag.ROOT_FOLDER].extend(sorted(orphans))

    return {"folder_tree": merged_tree, "doc_names": list(all_names), "count": len(all_names)}


@router.get("/projects/{name}/documents")
def list_documents(name: str, user: dict = Depends(get_current_user)):
    """Return all documents for a project from SQLite."""
    _verify_project_ownership(name, user["id"])
    from backend.database import get_db
    with get_db() as conn:
        project = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (name, user["id"])
        ).fetchone()
        if not project:
            raise HTTPException(status_code=404, detail="프로젝트를 찾을 수 없습니다.")
        docs = conn.execute(
            "SELECT id, folder, filename, size, uploaded_at FROM documents WHERE project_id = ? ORDER BY folder, filename",
            (project["id"],)
        ).fetchall()
    return [dict(d) for d in docs]


@router.post("/projects/{name}/folders")
def create_folder(name: str, req: FolderCreate, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_rag
    return core_rag.create_folder(name, req.name, owner_id=user["id"])


@router.post("/projects/{name}/folders/rename")
def rename_folder_ep(name: str, req: FolderRename, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_rag
    # Preserve parent path: "finance/reports" renamed to "Q1" → "finance/Q1"
    parts = req.old_name.rsplit("/", 1)
    new_full = f"{parts[0]}/{req.new_name}" if len(parts) > 1 else req.new_name
    return core_rag.rename_folder(name, req.old_name, new_full, owner_id=user["id"])


@router.delete("/projects/{name}/folders/{folder}")
def delete_folder(name: str, folder: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_rag
    return core_rag.delete_folder(name, folder, owner_id=user["id"])


@router.post("/projects/{name}/docs/{doc}/move")
def move_doc(name: str, doc: str, req: DocMoveRequest, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_rag
    result = core_rag.move_doc_to_folder(name, doc, req.target_folder, owner_id=user["id"])
    # Also sync SQLite documents table
    if result.get("success"):
        from backend.database import get_db
        with get_db() as conn:
            project = conn.execute(
                "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
                (name, user["id"])
            ).fetchone()
            if project:
                conn.execute(
                    "UPDATE documents SET folder = ? WHERE project_id = ? AND (filename = ? OR filename LIKE ?)",
                    (req.target_folder, project["id"], doc, f"{doc}.%"),
                )
    return result


@router.delete("/projects/{name}/docs/{doc}")
def trash_doc(name: str, doc: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_rag
    from backend.database import get_db
    result = core_rag.trash_document(name, doc, owner_id=user["id"])

    # Also delete from SQLite documents table (Railway ephemeral FS 대응)
    try:
        with get_db() as conn:
            project = conn.execute(
                "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
                (name, user["id"])
            ).fetchone()
            if project:
                doc_stem = os.path.splitext(doc)[0] if '.' in doc else doc
                conn.execute(
                    "DELETE FROM documents WHERE project_id = ? AND (filename = ? OR filename = ?)",
                    (project["id"], doc, doc_stem)
                )
    except Exception:
        pass

    return result


@router.post("/projects/{name}/sync-texts")
def sync_texts_to_server(name: str, payload: dict, user: dict = Depends(get_current_user)):
    """IndexedDB에서 파싱된 텍스트를 서버 RAG 저장소로 동기화.
    payload: { docs: [{ filename, parsedText, folder? }] }
    프로젝트가 서버에 없으면 자동 생성.
    """
    import core_rag
    from backend.database import get_db

    # 서버에 프로젝트가 없으면 자동 생성
    projects = core_rag.list_projects(owner_id=user["id"])
    if not any(p["name"] == name for p in projects):
        core_rag.create_project(name, owner_id=user["id"])

    docs = payload.get("docs", [])
    if not docs:
        return {"success": True, "indexed": [], "message": "동기화할 문서 없음"}

    texts = {d["filename"]: d["parsedText"] for d in docs if d.get("parsedText")}
    result = core_rag.index_texts("", texts, name, owner_id=user["id"])

    # Store in SQLite documents table
    with get_db() as conn:
        project = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (name, user["id"])
        ).fetchone()
        if project:
            for d in docs:
                fn = d.get("filename", "")
                folder = d.get("folder", "__root__")
                parsed = d.get("parsedText", "")
                size = len(parsed.encode("utf-8")) if parsed else 0
                conn.execute(
                    """INSERT INTO documents (project_id, folder, filename, parsed_text, size)
                       VALUES (?, ?, ?, ?, ?)
                       ON CONFLICT(project_id, folder, filename) DO UPDATE SET
                         parsed_text = excluded.parsed_text, size = excluded.size""",
                    (project["id"], folder, fn, parsed, size),
                )

    return result


@router.post("/projects/{name}/upload")
async def upload_files(name: str, files: List[UploadFile] = File(...), folder: str = Form("__root__"), user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_rag
    from backend.database import get_db
    api_key = _get_api_key()
    texts = {}
    parse_errors = []
    for f in files:
        content_bytes = await f.read()
        parsed = _parse_file_bytes(f.filename, content_bytes, api_key)
        if parsed:
            texts[f.filename] = parsed
        else:
            parse_errors.append(f.filename)
    result = core_rag.index_texts(api_key, texts, name, owner_id=user["id"])
    if parse_errors:
        result["parse_errors"] = parse_errors
    result["parsed_texts"] = texts

    # Move files to specified folder in folder tree
    target_folder = folder or "__root__"
    if target_folder != "__root__" and texts:
        for fn in texts:
            try:
                # Use stem (without extension) since _folders.json stores stems
                stem = os.path.splitext(fn)[0]
                core_rag.move_doc_to_folder(name, stem, target_folder, owner_id=user["id"])
            except Exception:
                pass

    # Store in SQLite documents table
    with get_db() as conn:
        project = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (name, user["id"])
        ).fetchone()
        if project:
            for fn, parsed_text in texts.items():
                size = len(parsed_text.encode("utf-8")) if parsed_text else 0
                conn.execute(
                    """INSERT INTO documents (project_id, folder, filename, parsed_text, size)
                       VALUES (?, ?, ?, ?, ?)
                       ON CONFLICT(project_id, folder, filename) DO UPDATE SET
                         parsed_text = excluded.parsed_text, size = excluded.size""",
                    (project["id"], folder or "__root__", fn, parsed_text, size),
                )

    # 자동 BM25 인덱싱 (백그라운드)
    if texts:
        import threading
        def _auto_index():
            try:
                all_docs = core_rag.load_project_docs_dict(name, owner_id=user["id"])
                if all_docs:
                    import core_rag_bm25
                    core_rag_bm25.build_index(name, all_docs)
                    logger.info(f"Auto BM25 index built for '{name}'")
            except Exception as e:
                logger.warning(f"Auto BM25 index failed for '{name}': {e}")
        threading.Thread(target=_auto_index, daemon=True).start()

    return result


# ========================================
# Wiki
# ========================================

@router.get("/projects/{name}/wiki")
def get_wiki(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_wiki
    wiki = core_wiki.load_wiki(name, owner_id=user["id"])
    return wiki or {"sections": [], "citations": [], "generated_at": None}


@router.post("/projects/{name}/wiki/generate")
def generate_wiki_endpoint(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    owner_id = user["id"]

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/wiki/generate")
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_wiki
            result = core_wiki.generate_wiki(api_key, name, owner_id=owner_id)
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.post("/projects/{name}/wiki/update")
def update_wiki_endpoint(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    owner_id = user["id"]

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/wiki/update")
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_wiki
            result = core_wiki.update_wiki(api_key, name, owner_id=owner_id)
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.post("/projects/{name}/wiki/sections/{section_id}/revise")
def revise_wiki_section(name: str, section_id: str, body: dict, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    owner_id = user["id"]
    instruction = body.get("instruction", "")
    selected_docs = body.get("selected_docs")

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/wiki/revise")
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_wiki
            result = core_wiki.revise_section(
                api_key, name, section_id, instruction,
                owner_id=owner_id, selected_docs=selected_docs,
            )
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.post("/projects/{name}/wiki/suggest-sections")
def suggest_wiki_sections(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_wiki
    api_key = _get_api_key()
    return core_wiki.suggest_sections(api_key, name, owner_id=user["id"])


@router.patch("/projects/{name}/wiki/sections/{section_id}")
def patch_wiki_section(name: str, section_id: str, body: WikiSectionUpdate, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_wiki
    return core_wiki.update_section(name, section_id, body.dict(exclude_none=True), owner_id=user["id"])


@router.post("/projects/{name}/wiki/sections")
def add_wiki_section(name: str, body: WikiSectionCreate, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_wiki
    return core_wiki.add_section(name, body.id, body.title, body.content, owner_id=user["id"])


@router.delete("/projects/{name}/wiki/sections/{section_id}")
def delete_wiki_section(name: str, section_id: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_wiki
    return core_wiki.delete_section(name, section_id, owner_id=user["id"])


@router.post("/projects/{name}/wiki/citation-context")
def citation_context(name: str, body: dict, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_wiki
    return core_wiki.get_citation_context(
        name, body.get("source_doc", ""), body.get("excerpt", ""),
        owner_id=user["id"],
    )


# ========================================
# Review Workflow — RFI Extraction & Cross-check
# ========================================

@router.post("/projects/{name}/review/extract-rfi")
def extract_rfi_endpoint(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    owner_id = user["id"]

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/review/extract-rfi", model=model)
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_review_workflow
            result = core_review_workflow.extract_rfi_from_wiki(
                api_key, model, name, owner_id=owner_id
            )
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.post("/projects/{name}/review/crosscheck")
def crosscheck_rfi_endpoint(name: str, body: dict, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    owner_id = user["id"]
    items = body.get("items", [])

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/review/crosscheck", model=model)
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_review_workflow
            result = core_review_workflow.crosscheck_rfi(
                api_key, model, name, items, owner_id=owner_id
            )
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.post("/projects/{name}/review/generate-external-rfi")
def generate_external_rfi_endpoint(name: str, body: dict, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    owner_id = user["id"]
    gap_items = body.get("gap_items", [])

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/review/generate-external-rfi", model=model)
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_review_workflow
            result = core_review_workflow.generate_external_rfi(
                api_key, model, gap_items, project_name=name, owner_id=owner_id
            )
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


# ========================================
# Excel Financial Model
# ========================================

@router.post("/projects/{name}/excel-model")
def generate_excel_model_endpoint(name: str, user: dict = Depends(get_current_user)):
    """Extract deal structure from wiki+docs and generate PEF cash flow Excel model."""
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    owner_id = user["id"]

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/excel-model", model=model)
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_excel_model
            result = core_excel_model.generate_excel_model(
                api_key, model, name, owner_id=owner_id
            )
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.post("/projects/{name}/excel-model/download")
def download_excel_model(name: str, body: dict, user: dict = Depends(get_current_user)):
    """Download generated Excel model from base64 data."""
    import base64
    excel_b64 = body.get("excel_b64", "")
    filename = body.get("filename", f"{name}_Model.xlsx")
    excel_bytes = base64.b64decode(excel_b64)
    return Response(
        content=excel_bytes,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ========================================
# Contract Comparison (신구조문 비교)
# ========================================

@router.post("/projects/{name}/contract-compare")
def contract_compare_endpoint(name: str, body: dict = None, user: dict = Depends(get_current_user)):
    """Compare termsheet vs contract drafts and generate comparison report."""
    _verify_project_ownership(name, user["id"])
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    owner_id = user["id"]
    selected_docs = (body or {}).get("selected_docs")

    from backend.api_ws import _tasks
    task_id = create_task(user_id=user["id"], endpoint="/contract-compare", model=model)
    task = _tasks[task_id]

    import threading
    def _run():
        try:
            import core_contract_compare
            result = core_contract_compare.compare_contracts(
                api_key, model, name, owner_id=owner_id, selected_docs=selected_docs
            )
            task["status"] = "complete"
            task["result"] = result
        except Exception as e:
            task["status"] = "error"
            task["error"] = str(e)
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


# ========================================
# Local Folder Scan — Preview & Ingest
# ========================================

SUPPORTED_SCAN_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.pptx', '.ppt',
    '.xlsx', '.xls', '.csv',
    '.txt', '.md', '.json', '.xml', '.html', '.htm',
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',
}

MAX_SCAN_FILES = 500


def _validate_scan_path(folder_path: str) -> str:
    """폴더 경로 유효성 검증. 안전한 실제 경로를 반환."""
    real_path = os.path.realpath(folder_path)
    if ".." in os.path.normpath(folder_path).split(os.sep):
        raise HTTPException(status_code=400, detail="경로에 '..'은 사용할 수 없습니다.")
    if not os.path.exists(real_path):
        raise HTTPException(status_code=400, detail=f"경로가 존재하지 않습니다: {folder_path}")
    if not os.path.isdir(real_path):
        raise HTTPException(status_code=400, detail=f"디렉토리가 아닙니다: {folder_path}")
    # Check allowed_scan_dirs from settings
    settings = _load_settings()
    allowed_dirs = settings.get("allowed_scan_dirs", [])
    if allowed_dirs:
        real_lower = real_path.lower().replace("\\", "/")
        allowed = False
        for d in allowed_dirs:
            d_real = os.path.realpath(d).lower().replace("\\", "/")
            if real_lower.startswith(d_real):
                allowed = True
                break
        if not allowed:
            raise HTTPException(status_code=403, detail="허용되지 않은 경로입니다.")
    return real_path


@router.post("/scan-folder/preview")
async def scan_folder_preview(req: FolderScanRequest, user: dict = Depends(get_current_user)):
    """로컬 폴더를 스캔하여 지원 파일 목록을 미리보기로 반환."""
    real_path = _validate_scan_path(req.folder_path)
    ext_filter = set(e.lower() if e.startswith('.') else f'.{e.lower()}' for e in req.file_extensions) if req.file_extensions else None

    files: List[FolderScanFileInfo] = []
    if req.recursive:
        walker = os.walk(real_path)
    else:
        # Single level: just the top directory
        try:
            entries = os.listdir(real_path)
        except PermissionError:
            raise HTTPException(status_code=403, detail="폴더 접근 권한이 없습니다.")
        walker = [(real_path, [], [e for e in entries if os.path.isfile(os.path.join(real_path, e))])]

    for dirpath, _dirnames, filenames in walker:
        for fname in filenames:
            if len(files) >= MAX_SCAN_FILES:
                break
            # Windows surrogateescape 방지: 잘못된 유니코드 파일명 건너뛰기
            try:
                fname.encode("utf-8")
            except (UnicodeEncodeError, UnicodeDecodeError):
                continue
            ext = os.path.splitext(fname)[1].lower()
            if ext not in SUPPORTED_SCAN_EXTENSIONS:
                continue
            if ext_filter and ext not in ext_filter:
                continue
            # ~$ 임시 파일 건너뛰기 (Word/Excel lock 파일)
            if fname.startswith("~$"):
                continue
            full_path = os.path.join(dirpath, fname)
            try:
                size = os.path.getsize(full_path)
            except OSError:
                continue
            rel_path = os.path.relpath(full_path, real_path).replace("\\", "/")
            files.append(FolderScanFileInfo(
                path=full_path.replace("\\", "/"),
                name=fname,
                size=size,
                ext=ext,
                relative_path=rel_path,
            ))
        if len(files) >= MAX_SCAN_FILES:
            break

    total_size = sum(f.size for f in files)
    return FolderScanPreviewResponse(files=files, total_size=total_size, file_count=len(files))


@router.post("/projects/{name}/scan-folder/ingest")
async def scan_folder_ingest(name: str, req: FolderIngestRequest, user: dict = Depends(get_current_user)):
    """선택된 로컬 파일들을 파싱하여 프로젝트에 인덱싱."""
    _verify_project_ownership(name, user["id"])
    real_folder = _validate_scan_path(req.folder_path)
    api_key = _get_api_key()
    import core_rag
    from backend.database import get_db

    indexed = []
    skipped = []
    errors = []

    # Group files by folder for preserve_structure
    folder_groups: Dict[str, Dict[str, str]] = {}  # folder_name -> {filename: parsed_text}

    for file_path in req.selected_files:
        # Security: ensure file is under the scanned folder
        file_real = os.path.realpath(file_path)
        if not file_real.lower().startswith(real_folder.lower()):
            errors.append({"file": file_path, "error": "허용된 폴더 밖의 파일입니다."})
            continue
        if not os.path.isfile(file_real):
            errors.append({"file": file_path, "error": "파일이 존재하지 않습니다."})
            continue

        filename = os.path.basename(file_real)
        try:
            with open(file_real, "rb") as fh:
                data = fh.read()
            parsed = _parse_file_bytes(filename, data, api_key)
            if not parsed:
                skipped.append(filename)
                continue

            if req.preserve_structure:
                rel = os.path.relpath(file_real, real_folder)
                rel_dir = os.path.dirname(rel).replace("\\", "/")
                folder_key = rel_dir if rel_dir and rel_dir != "." else "__root__"
            else:
                folder_key = "__root__"

            if folder_key not in folder_groups:
                folder_groups[folder_key] = {}
            folder_groups[folder_key][filename] = parsed
        except Exception as e:
            errors.append({"file": filename, "error": str(e)})

    # Index each folder group
    for folder_key, texts in folder_groups.items():
        try:
            if folder_key == "__root__":
                result = core_rag.index_texts(api_key, texts, name, owner_id=user["id"])
            else:
                result = core_rag.index_texts_to_folder(
                    api_key, texts, name, folder=folder_key, owner_id=user["id"]
                )
            indexed.extend(result.get("indexed", []))
            skipped.extend(result.get("skipped", []))
            if result.get("errors"):
                errors.extend(result["errors"])

            # Store in SQLite documents table
            with get_db() as conn:
                project = conn.execute(
                    "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
                    (name, user["id"])
                ).fetchone()
                if project:
                    db_folder = folder_key if folder_key != "__root__" else "__root__"
                    for fn, parsed_text in texts.items():
                        size = len(parsed_text.encode("utf-8")) if parsed_text else 0
                        conn.execute(
                            """INSERT INTO documents (project_id, folder, filename, parsed_text, size)
                               VALUES (?, ?, ?, ?, ?)
                               ON CONFLICT(project_id, folder, filename) DO UPDATE SET
                                 parsed_text = excluded.parsed_text, size = excluded.size""",
                            (project["id"], db_folder, fn, parsed_text, size),
                        )
        except Exception as e:
            errors.append({"folder": folder_key, "error": str(e)})

    return {
        "success": True,
        "indexed": indexed,
        "skipped": skipped,
        "errors": errors,
        "indexed_count": len(indexed),
        "skipped_count": len(skipped),
        "error_count": len(errors),
    }


@router.post("/parse-files")
async def parse_files_only(files: List[UploadFile] = File(...), user: dict = Depends(get_current_user)):
    """Parse uploaded files and return text content without server-side storage.
    Used by local-storage mode where files are stored in browser IndexedDB."""
    api_key = _get_api_key()
    parsed = {}
    errors = []
    for f in files:
        content_bytes = await f.read()
        text = _parse_file_bytes(f.filename, content_bytes, api_key)
        if text:
            parsed[f.filename] = text
        else:
            errors.append(f.filename)
    return {"parsed_texts": parsed, "errors": errors, "count": len(parsed)}


@router.post("/extract-excel-cells")
async def extract_excel_cells(files: List[UploadFile] = File(...), user: dict = Depends(get_current_user)):
    """엑셀 파일에서 행 단위로 텍스트 추출 (LP Q&A 질문 목록용).
    동일 행의 모든 셀을 ' | '로 합치고, 셀 내 줄바꿈은 유지하여 하나의 질문으로 취급.
    """
    rows_out: list[str] = []
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        if ext not in ('.xlsx', '.xls', '.csv'):
            data = await f.read()
            lines = data.decode("utf-8", errors="replace").split("\n")
            rows_out.extend(l.strip() for l in lines if l.strip())
            continue
        try:
            import pandas as pd
            data = await f.read()
            if ext == '.csv':
                df_dict = {"Sheet1": pd.read_csv(io.BytesIO(data), header=None)}
            else:
                df_dict = pd.read_excel(io.BytesIO(data), sheet_name=None, header=None)
            for _sheet_name, df in df_dict.items():
                for _, row in df.iterrows():
                    parts = []
                    for val in row:
                        if pd.notna(val):
                            text = str(val).strip()
                            if text:
                                parts.append(text)
                    if parts:
                        combined = " | ".join(parts)
                        # 너무 짧거나 숫자만이면 제목/번호행 → 건너뜀
                        if len(combined) > 5:
                            rows_out.append(combined)
        except Exception as e:
            rows_out.append(f"파싱 오류: {e}")
    return {"cells": rows_out, "count": len(rows_out)}


def _parse_file_bytes(filename: str, data: bytes, api_key: str = "") -> str:
    """파일 바이트를 텍스트로 변환. PDF/DOCX/PPTX/XLSX 파서 지원."""
    import tempfile
    ext = os.path.splitext(filename)[1].lower()

    # 텍스트 파일은 직접 디코드
    if ext in ('.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm'):
        return data.decode("utf-8", errors="replace")

    # MarkItDown 우선 시도 (PDF, DOCX, PPTX, XLSX 등 다양한 포맷 지원)
    try:
        from markitdown import MarkItDown
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(data)
            tmp_path = tmp.name
        try:
            import concurrent.futures
            def _convert(path):
                md = MarkItDown()
                return md.convert(path)
            with concurrent.futures.ThreadPoolExecutor() as executor:
                future = executor.submit(_convert, tmp_path)
                result = future.result(timeout=120)
            text_content = result.text_content if result else ""
            # PDF: 스캔본 감지 — 의미 있는 텍스트인지 검증
            if ext == '.pdf' and text_content:
                # 한글/영문 비율이 낮으면 바이너리/인코딩 쓰레기
                import re as _re
                readable = _re.findall(r'[가-힣a-zA-Z0-9]', text_content)
                ratio = len(readable) / max(len(text_content), 1)
                if ratio < 0.3 or len(text_content.strip()) < 200:
                    text_content = ""  # OCR로 폴백
            min_len = 200 if ext == '.pdf' else 50
            if text_content and len(text_content.strip()) > min_len:
                return f"### [파일명: {filename}]\n{text_content}"
        finally:
            if os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except Exception:
                    logger.debug("Failed to remove temp file %s", tmp_path)
    except Exception:
        logger.debug("MarkItDown conversion failed for %s", filename)

    # PDF fallback: PyMuPDF → OCR if scanned
    if ext == '.pdf':
        try:
            import fitz
            with fitz.open(stream=data, filetype="pdf") as doc:
                pages = []
                for page in doc:
                    pages.append(page.get_text())
                text = "\n\n".join(pages)
                # 읽을 수 있는 텍스트 비율 검증
                import re as _re2
                readable2 = _re2.findall(r'[가-힣a-zA-Z0-9]', text)
                ratio2 = len(readable2) / max(len(text), 1)
                if text.strip() and len(text.strip()) > 100 and ratio2 > 0.3:
                    return f"### [파일명: {filename}]\n{text}"
                # 텍스트가 거의 없거나 바이너리면 스캔 PDF → Gemini OCR 시도
                if api_key:
                    try:
                        import ocr as ocr_module
                        ocr_text = ocr_module.extract_pdf_with_gemini_ocr(doc, api_key)
                        if ocr_text and len(ocr_text.strip()) > 50:
                            return f"### [파일명: {filename} (OCR)]\n{ocr_text}"
                    except Exception as ocr_err:
                        print(f"[parse] OCR failed for {filename}: {ocr_err}")
        except Exception:
            logger.debug("PDF fallback parsing failed for %s", filename)

    # DOCX fallback: python-docx
    if ext in ('.docx', '.doc'):
        try:
            from docx import Document
            doc = Document(io.BytesIO(data))
            text = "\n".join(p.text for p in doc.paragraphs)
            if text.strip():
                return f"### [파일명: {filename}]\n{text}"
        except Exception:
            logger.debug("DOCX fallback parsing failed for %s", filename)

    # PPTX fallback: python-pptx
    if ext in ('.pptx', '.ppt'):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            text = "\n".join(text_parts)
            if text.strip():
                return f"### [파일명: {filename}]\n{text}"
        except Exception:
            logger.debug("PPTX fallback parsing failed for %s", filename)

    # XLSX fallback: openpyxl
    if ext in ('.xlsx', '.xls'):
        try:
            import pandas as pd
            df = pd.read_excel(io.BytesIO(data), sheet_name=None)
            parts = []
            for sheet_name, sheet_df in df.items():
                parts.append(f"## 시트: {sheet_name}\n{sheet_df.to_markdown()}")
            text = "\n\n".join(parts)
            if text.strip():
                return f"### [파일명: {filename}]\n{text}"
        except Exception:
            logger.debug("Excel fallback parsing failed for %s", filename)

    # 최종 fallback: UTF-8 디코드
    return data.decode("utf-8", errors="replace")


# ========================================
# AI Generation
# ========================================

@router.post("/generate")
def start_generate(req: GenerateRequest, user: dict = Depends(get_current_user)):
    if req.project_name and not req.file_context.strip():
        _verify_project_ownership(req.project_name, user["id"])
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")

    user_context = req.file_context.strip()
    selected_docs = req.inputs.get("selected_docs", [])

    inputs = dict(req.inputs)
    inputs.setdefault("template_option", req.template_option)
    inputs["context_text"] = user_context

    # 문서 선택을 빠르게 (try/except로 502 방지)
    max_chars = _max_chars_for_model(model)
    file_context = ""
    try:
        if req.project_name:
            query = inputs.get("context_text", "") or req.template_option
            file_context = _select_relevant_docs(req.project_name, query, model, selected_docs if selected_docs else None)
        if not file_context and user_context:
            file_context = _truncate_context(user_context, max_chars)
    except Exception as e:
        print(f"[generate] doc selection error: {e}")
        if user_context:
            file_context = _truncate_context(user_context, max_chars)

    title = req.template_option or req.inputs.get("template_option", "")
    task_id = create_task(
        user_id=user["id"], endpoint="/generate", model=model,
        title=title, inputs={"template_option": req.template_option, "mode": req.mode},
    )
    run_generate_task(
        task_id, api_key, model,
        inputs, req.thinking_level, file_context,
        mode=req.mode
    )
    log_usage(user["id"], "/generate", model)
    return {"task_id": task_id}


@router.post("/qa")
def start_qa(req: QaRequest, user: dict = Depends(get_current_user)):
    if req.project_name and not req.file_context.strip():
        _verify_project_ownership(req.project_name, user["id"])
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")

    # 문서 선택 (try/except로 502 방지)
    max_chars = _max_chars_for_model(model)
    context = ""
    try:
        if req.project_name:
            context = _select_relevant_docs(req.project_name, req.question, model, req.selected_docs, owner_id=user["id"])
        if not context and req.file_context:
            context = _truncate_context(req.file_context.strip(), max_chars)
    except Exception as e:
        print(f"[qa] doc selection error: {e}")
        if req.file_context:
            context = _truncate_context(req.file_context.strip(), max_chars)

    task_id = create_task(
        user_id=user["id"], endpoint="/qa", model=model,
        title=req.question[:100], inputs={"question": req.question},
    )
    run_analysis_task(
        task_id, "qa_answer", api_key, model,
        file_context=context, question=req.question
    )
    log_usage(user["id"], "/qa", model)
    return {"task_id": task_id}


@router.post("/analyze")
def start_analysis(req: AnalysisRequest, user: dict = Depends(get_current_user)):
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")

    kwargs = dict(req.kwargs)
    max_chars = _max_chars_for_model(model)

    # Inline docs from local folder (client-side, not stored on server)
    inline_docs = kwargs.pop("inline_docs", None) or {}

    # Batch mode: pass individual docs_dict instead of combined file_context
    if req.task_type == 'material_summary_batch' and "project_name" in kwargs:
        import core_rag
        pname = kwargs.pop("project_name")
        sel_docs = kwargs.pop("selected_docs", [])
        docs_dict = core_rag.load_project_docs_dict(pname, owner_id=user["id"])
        if sel_docs and docs_dict:
            sel_stems = {_strip_doc_stem(s) for s in sel_docs}
            docs_dict = {k: v for k, v in docs_dict.items()
                         if _strip_doc_stem(k) in sel_stems or k in sel_docs}
        kwargs.pop("file_context", None)
        kwargs["docs_dict"] = docs_dict or {}
    elif "project_name" in kwargs:
        pname = kwargs.pop("project_name")
        sel_docs = kwargs.pop("selected_docs", [])
        query = kwargs.get("question", req.task_type)
        try:
            server_ctx = _select_relevant_docs(pname, query, model, sel_docs, owner_id=user["id"])
            if server_ctx:
                kwargs["file_context"] = server_ctx
            elif kwargs.get("file_context"):
                kwargs["file_context"] = _truncate_context(kwargs["file_context"], max_chars)
        except Exception as e:
            print(f"[analyze] doc selection error: {e}")
            if kwargs.get("file_context"):
                kwargs["file_context"] = _truncate_context(kwargs["file_context"], max_chars)

    # Append inline local docs to file_context
    if inline_docs and isinstance(inline_docs, dict):
        inline_text = "\n\n".join(
            f"[로컬 파일: {k}]\n{v}" for k, v in inline_docs.items()
        )
        if kwargs.get("file_context"):
            kwargs["file_context"] += "\n\n" + inline_text
        else:
            kwargs["file_context"] = inline_text

    task_id = create_task(
        user_id=user["id"], endpoint=f"/analyze/{req.task_type}", model=model,
        title=req.task_type, inputs={"task_type": req.task_type},
    )
    run_analysis_task(task_id, req.task_type, api_key, model, **kwargs)
    log_usage(user["id"], f"/analyze/{req.task_type}", model)
    return {"task_id": task_id}


# ========================================
# Document Download
# ========================================

@router.get("/projects/{name}/docs/{doc_name}/download")
def download_doc(name: str, doc_name: str, user: dict = Depends(get_current_user)):
    """프로젝트 문서를 마크다운 파일로 다운로드."""
    import core_rag
    from backend.database import get_db
    storage = core_rag._get_storage_name(name, owner_id=user["id"])
    content = core_rag._load_doc_file(storage, doc_name)

    # 파일시스템에 없으면 SQLite에서 로드 (Railway 폴백)
    if not content:
        with get_db() as conn:
            project = conn.execute(
                "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
                (name, user["id"])
            ).fetchone()
            if project:
                row = conn.execute(
                    "SELECT parsed_text FROM documents WHERE project_id = ? AND filename = ?",
                    (project["id"], doc_name)
                ).fetchone()
                if row:
                    content = row["parsed_text"]

    if not content:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다")
    filename = f"{doc_name}.md" if not doc_name.endswith('.md') else doc_name
    return Response(
        content=content.encode("utf-8"),
        media_type="text/markdown; charset=utf-8",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ========================================
# Web Research — Save to Project
# ========================================

@router.post("/save-research")
def save_research(req: SaveResearchRequest, user: dict = Depends(get_current_user)):
    """웹 리서치 결과를 프로젝트 문서로 저장."""
    import core_rag
    from backend.database import get_db
    api_key = _get_api_key()
    texts = {req.doc_name: req.content}
    result = core_rag.index_texts(api_key, texts, req.project_name, owner_id=user["id"])

    # SQLite에도 저장 (Railway 에포메럴 FS 대비)
    with get_db() as conn:
        project = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (req.project_name, user["id"])
        ).fetchone()
        if project and req.content:
            size = len(req.content.encode("utf-8"))
            conn.execute(
                """INSERT INTO documents (project_id, folder, filename, parsed_text, size)
                   VALUES (?, ?, ?, ?, ?)
                   ON CONFLICT(project_id, folder, filename) DO UPDATE SET
                     parsed_text = excluded.parsed_text, size = excluded.size""",
                (project["id"], "__root__", req.doc_name, req.content, size),
            )

    return result


# ========================================
# Research Notes (Obsidian-like)
# ========================================

@router.get("/projects/{name}/notes")
def list_notes(name: str, tag: str = None, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_notes
    return core_notes.list_notes(name, user["id"], tag=tag)


@router.post("/projects/{name}/notes")
def create_note(name: str, req: NoteCreate, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_notes
    return core_notes.create_note(name, user["id"], req.title, req.content, req.tags)


@router.get("/projects/{name}/notes/tags")
def get_note_tags(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_notes
    return core_notes.get_tags(name, user["id"])


@router.get("/projects/{name}/notes/{slug}")
def get_note(name: str, slug: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_notes
    note = core_notes.get_note(name, user["id"], slug)
    if not note:
        raise HTTPException(status_code=404, detail="노트를 찾을 수 없습니다.")
    return note


@router.patch("/projects/{name}/notes/{slug}")
def update_note(name: str, slug: str, req: NoteUpdate, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_notes
    return core_notes.update_note(name, user["id"], slug, req.title, req.content, req.tags)


@router.delete("/projects/{name}/notes/{slug}")
def delete_note(name: str, slug: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_notes
    return core_notes.delete_note(name, user["id"], slug)


@router.get("/projects/{name}/notes/{slug}/backlinks")
def get_note_backlinks(name: str, slug: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    import core_notes
    return core_notes.get_backlinks(name, user["id"], slug)


# ========================================
# PPT Generation
# ========================================

@router.post("/create-pptx")
def create_pptx(req: CreatePptxRequest, user: dict = Depends(get_current_user)):
    """JSON slide data → PPTX 파일 생성 및 다운로드."""
    import utils_ppt

    slide_json = req.slide_json
    template_path = getattr(req, 'template_path', None)
    pptx_bytes = utils_ppt.create_deck_from_json(slide_json, template_path=template_path)
    if not pptx_bytes:
        raise HTTPException(status_code=400, detail="PPTX 생성 실패: 유효하지 않은 슬라이드 데이터")

    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
    )


@router.post("/create-ib-pptx")
def create_ib_pptx(req: CreatePptxRequest, user: dict = Depends(get_current_user)):
    """좌표 기반 동적 PPT 생성. pptxgenjs 우선, 실패 시 python-pptx 폴백."""
    import subprocess, sys, tempfile

    slide_json = req.slide_json
    if isinstance(slide_json, str):
        slide_json = json.loads(slide_json)

    base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    # 1차: pptxgenjs (Node.js) 시도
    pptx_bytes = None
    json_path = None
    pptx_path = None
    try:
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".json", dir=base, delete=False, encoding="utf-8"
        ) as f:
            json.dump(slide_json, f, ensure_ascii=False, indent=2)
            json_path = f.name
        pptx_path = json_path.replace(".json", ".pptx")

        result = subprocess.run(
            ["node", os.path.join(base, "generate_pptx_dynamic.js"), json_path, pptx_path],
            capture_output=True, text=True, timeout=30, cwd=base,
        )
        if result.returncode == 0 and os.path.exists(pptx_path):
            with open(pptx_path, "rb") as f:
                pptx_bytes = f.read()
            logger.info("PPT generated via pptxgenjs (Node.js)")
    except Exception as e:
        logger.warning(f"pptxgenjs failed, falling back to python-pptx: {e}")
    finally:
        for p in [json_path, pptx_path]:
            if p and os.path.exists(p):
                try: os.remove(p)
                except OSError: pass

    # 2차: python-pptx 폴백
    if not pptx_bytes:
        try:
            pptx_bytes = _render_dynamic_pptx_python(slide_json)
            logger.info("PPT generated via python-pptx fallback")
        except Exception as e:
            logger.error(f"python-pptx fallback also failed: {e}")
            raise HTTPException(status_code=500, detail=f"PPT 생성 실패: {e}")

    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
    )


def _render_dynamic_pptx_python(deck_json: dict) -> bytes:
    """좌표 기반 elements[]를 python-pptx로 렌더링 (Node.js 없을 때 폴백)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    import io

    # --- Design system colors ---
    CLR_PRIMARY = "1B2A4A"
    CLR_ACCENT = "86BC25"
    CLR_SECONDARY = "0076A8"
    CLR_ALERT = "C4262E"
    CLR_BG = "F6F6F6"
    CLR_TEXT = "2D2D2D"
    CLR_TABLE_HDR = "1B2A4A"
    CLR_ALT_ROW = "F2F5F7"

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    def _add_text_paragraphs(tf, text_str, font_size, bold, color_str, align=None):
        """Split text by newlines, handle bullet points, add paragraphs."""
        lines = str(text_str).split("\n")
        for li, line in enumerate(lines):
            if li == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            is_bullet = line.lstrip().startswith("•") or line.lstrip().startswith("- ")
            display_text = line.lstrip("•- ").strip() if is_bullet else line
            if is_bullet:
                p.level = 1
                pPr = p._pPr
                if pPr is None:
                    pPr = p._p.get_or_add_pPr()
                buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
                pPr.append(buChar)
            run = p.add_run()
            run.text = display_text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color_str)
            if align:
                p.alignment = align

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)
    blank_layout = prs.slide_layouts[6]  # blank

    slides = deck_json.get("slides", [])

    for slide_idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        bg_color = slide_data.get("background", "FFFFFF")
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_color)

        for el in slide_data.get("elements", []):
            el_type = el.get("type", "")
            x = Inches(el.get("x", 0))
            y = Inches(el.get("y", 0))
            w = Inches(el.get("w", 1))
            h = Inches(el.get("h", 0.5))

            # ===== TEXT =====
            if el_type == "text":
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = True
                text_content = el.get("text", "")
                el_align = align_map.get(el.get("align", "left"), PP_ALIGN.LEFT)

                if isinstance(text_content, list):
                    # Rich text runs
                    for ri, run_data in enumerate(text_content):
                        if ri == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        if isinstance(run_data, str):
                            run_text = run_data
                            r_size = el.get("fontSize", 11)
                            r_bold = el.get("bold", False)
                            r_italic = el.get("italic", False)
                            r_color = el.get("color", CLR_TEXT)
                        else:
                            run_text = run_data.get("text", "")
                            r_size = run_data.get("fontSize", el.get("fontSize", 11))
                            r_bold = run_data.get("bold", el.get("bold", False))
                            r_italic = run_data.get("italic", el.get("italic", False))
                            r_color = run_data.get("color") or el.get("color", CLR_TEXT)

                        # Handle line breaks within a single run
                        sub_lines = run_text.split("\n")
                        for sli, sub_line in enumerate(sub_lines):
                            if sli > 0:
                                p = tf.add_paragraph()
                            run = p.add_run()
                            run.text = sub_line
                            run.font.size = Pt(r_size)
                            run.font.bold = r_bold
                            run.font.italic = r_italic
                            run.font.color.rgb = RGBColor.from_string(r_color)
                        p.alignment = el_align
                else:
                    _add_text_paragraphs(
                        tf, text_content,
                        font_size=el.get("fontSize", 11),
                        bold=el.get("bold", False),
                        color_str=el.get("color", CLR_TEXT),
                        align=el_align,
                    )

                # Background fill
                if el.get("fill"):
                    txBox.fill.solid()
                    txBox.fill.fore_color.rgb = RGBColor.from_string(el["fill"])

                # Shadow (drop shadow via XML)
                if el.get("shadow"):
                    spPr = txBox._element.spPr
                    effectLst = spPr.makeelement(qn("a:effectLst"), {})
                    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
                        "blurRad": "50800", "dist": "38100", "dir": "2700000",
                        "algn": "tl", "rotWithShape": "0",
                    })
                    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
                    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": "40000"})
                    srgbClr.append(alpha)
                    outerShdw.append(srgbClr)
                    effectLst.append(outerShdw)
                    spPr.append(effectLst)

            # ===== SHAPE =====
            elif el_type == "shape":
                shape_type_name = el.get("shapeType", "RECTANGLE").upper()
                shape_enum = getattr(MSO_SHAPE, shape_type_name, MSO_SHAPE.RECTANGLE)
                shape = slide.shapes.add_shape(shape_enum, x, y, w, h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(el.get("fill", CLR_BG))
                shape.line.fill.background()
                # Optional text inside shape
                if el.get("text"):
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = str(el["text"])
                    run.font.size = Pt(el.get("fontSize", 10))
                    run.font.color.rgb = RGBColor.from_string(el.get("color", CLR_TEXT))
                    p.alignment = align_map.get(el.get("align", "center"), PP_ALIGN.CENTER)

            # ===== DIVIDER =====
            elif el_type == "divider":
                div_h = Inches(el.get("h", 0.02))
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, div_h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(el.get("fill", "CCCCCC"))
                shape.line.fill.background()

            # ===== ICON_TEXT =====
            elif el_type == "icon_text":
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                icon_char = el.get("icon", "\u25CF")  # default: filled circle
                icon_run = p.add_run()
                icon_run.text = icon_char + "  "
                icon_run.font.size = Pt(el.get("iconSize", el.get("fontSize", 14)))
                icon_run.font.color.rgb = RGBColor.from_string(el.get("iconColor", CLR_ACCENT))
                text_run = p.add_run()
                text_run.text = el.get("text", "")
                text_run.font.size = Pt(el.get("fontSize", 11))
                text_run.font.bold = el.get("bold", False)
                text_run.font.color.rgb = RGBColor.from_string(el.get("color", CLR_TEXT))
                p.alignment = align_map.get(el.get("align", "left"), PP_ALIGN.LEFT)
                if el.get("fill"):
                    txBox.fill.solid()
                    txBox.fill.fore_color.rgb = RGBColor.from_string(el["fill"])

            # ===== TABLE =====
            elif el_type == "table":
                rows = el.get("rows", [])
                if not rows:
                    continue
                n_rows = len(rows)
                n_cols = max(len(r) for r in rows) if rows else 1
                table_h = Inches(el.get("h", n_rows * 0.3))
                table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, table_h)
                tbl = table_shape.table
                has_total = el.get("totalRow", False)
                font_size = el.get("fontSize", 9)

                for ri, row in enumerate(rows):
                    is_header = (ri == 0)
                    is_total = has_total and (ri == n_rows - 1)
                    is_alt = (not is_header) and (ri % 2 == 0) and (not is_total)

                    for ci, cell_data in enumerate(row):
                        if ci >= n_cols:
                            break
                        cell = tbl.cell(ri, ci)

                        # Extract cell text and optional cell-level props
                        if isinstance(cell_data, dict):
                            cell_text = str(cell_data.get("text", ""))
                            cell_fill = cell_data.get("fill")
                            cell_bold = cell_data.get("bold")
                            cell_color = cell_data.get("color")
                            cell_align = cell_data.get("align")
                        else:
                            cell_text = str(cell_data) if cell_data is not None else ""
                            cell_fill = None
                            cell_bold = None
                            cell_color = None
                            cell_align = None

                        cell.text = cell_text

                        # --- Cell styling ---
                        for p in cell.text_frame.paragraphs:
                            if cell_align:
                                p.alignment = align_map.get(cell_align, PP_ALIGN.LEFT)
                            for run in p.runs:
                                run.font.size = Pt(font_size)

                        if is_header:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor.from_string(el.get("headerFill", CLR_TABLE_HDR))
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    run.font.color.rgb = RGBColor.from_string("FFFFFF")
                                    run.font.bold = True
                        elif is_total:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor.from_string(CLR_PRIMARY)
                            # Top border for total row
                            tc = cell._tc
                            tcPr = tc.get_or_add_tcPr()
                            for border_tag in ["a:lnT"]:
                                ln = tcPr.makeelement(qn(border_tag), {"w": "12700", "cmpd": "sng"})
                                solidFill = ln.makeelement(qn("a:solidFill"), {})
                                srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": CLR_PRIMARY})
                                solidFill.append(srgb)
                                ln.append(solidFill)
                                tcPr.append(ln)
                        elif is_alt:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor.from_string(CLR_ALT_ROW)

                        # Cell-level fill override (conditional formatting)
                        if cell_fill:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor.from_string(cell_fill)
                        if cell_bold is not None:
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    run.font.bold = cell_bold
                        if cell_color:
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    run.font.color.rgb = RGBColor.from_string(cell_color)

            # ===== CHART =====
            elif el_type == "chart":
                try:
                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE

                    chart_type_str = el.get("chartType", "bar").lower()
                    chart_type_map = {
                        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                        "line": XL_CHART_TYPE.LINE,
                        "pie": XL_CHART_TYPE.PIE,
                        "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
                        "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
                    }
                    xl_chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

                    chart_data_list = el.get("chart_data") or el.get("chartData") or []
                    # chart_data_list: [{name, labels, values}]

                    cd = CategoryChartData()
                    if chart_data_list:
                        # Use labels from the first series
                        labels = chart_data_list[0].get("labels", [])
                        cd.categories = labels
                        for series_info in chart_data_list:
                            s_name = series_info.get("name", "Series")
                            s_values = series_info.get("values", [])
                            # Pad or truncate values to match labels length
                            padded = list(s_values) + [0] * max(0, len(labels) - len(s_values))
                            cd.add_series(s_name, padded[:len(labels)])

                    chart_shape = slide.shapes.add_chart(xl_chart_type, x, y, w, h, cd)
                    chart = chart_shape.chart
                    chart.has_legend = len(chart_data_list) > 1

                    # Apply design colors to series
                    series_colors = [CLR_PRIMARY, CLR_ACCENT, CLR_SECONDARY, CLR_ALERT, "4472C4", "ED7D31"]
                    for si, series in enumerate(chart.series):
                        color_hex = series_colors[si % len(series_colors)]
                        fill = series.format.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor.from_string(color_hex)
                        if chart_type_str == "line":
                            series.format.line.color.rgb = RGBColor.from_string(color_hex)
                            series.format.line.width = Pt(2)
                            series.smooth = False

                    # Chart title
                    if el.get("title"):
                        chart.has_title = True
                        chart.chart_title.text_frame.paragraphs[0].text = el["title"]
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
                        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(CLR_PRIMARY)
                except Exception:
                    # Fallback: render a placeholder shape if chart fails
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(CLR_BG)
                    shape.line.fill.background()
                    tf = shape.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = f"[Chart: {el.get('title', el.get('chartType', 'chart'))}]"
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor.from_string(CLR_TEXT)
                    p.alignment = PP_ALIGN.CENTER

            # ===== KPI_CARD / CALLOUT =====
            elif el_type in ("kpi_card", "callout"):
                txBox = slide.shapes.add_textbox(x, y, w, h)
                txBox.fill.solid()
                txBox.fill.fore_color.rgb = RGBColor.from_string(el.get("fill", CLR_BG))
                tf = txBox.text_frame
                tf.word_wrap = True
                if el.get("label"):
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = el["label"]
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor.from_string("6B6B6B")
                    p.alignment = align_map.get(el.get("align", "center"), PP_ALIGN.CENTER)
                if el.get("value"):
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = el["value"]
                    run.font.size = Pt(20)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor.from_string(CLR_PRIMARY)
                    p.alignment = align_map.get(el.get("align", "center"), PP_ALIGN.CENTER)
                if el.get("change"):
                    p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = el["change"]
                    run.font.size = Pt(9)
                    is_pos = "+" in el["change"] or "\u25B2" in el["change"]
                    run.font.color.rgb = RGBColor.from_string(CLR_ACCENT if is_pos else CLR_ALERT)
                    p.alignment = align_map.get(el.get("align", "center"), PP_ALIGN.CENTER)

                # Shadow for cards
                if el.get("shadow", el_type == "kpi_card"):
                    spPr = txBox._element.spPr
                    effectLst = spPr.makeelement(qn("a:effectLst"), {})
                    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
                        "blurRad": "40000", "dist": "25400", "dir": "2700000",
                        "algn": "tl", "rotWithShape": "0",
                    })
                    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
                    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": "25000"})
                    srgbClr.append(alpha)
                    outerShdw.append(srgbClr)
                    effectLst.append(outerShdw)
                    spPr.append(effectLst)

        # Speaker notes
        if slide_data.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

        # Slide number (bottom-right)
        sn_left = Inches(9.1)
        sn_top = Inches(5.3)
        sn_w = Inches(0.7)
        sn_h = Inches(0.25)
        sn_box = slide.shapes.add_textbox(sn_left, sn_top, sn_w, sn_h)
        sn_tf = sn_box.text_frame
        sn_p = sn_tf.paragraphs[0]
        sn_run = sn_p.add_run()
        sn_run.text = str(slide_idx + 1)
        sn_run.font.size = Pt(8)
        sn_run.font.color.rgb = RGBColor.from_string("999999")
        sn_p.alignment = PP_ALIGN.RIGHT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@router.post("/slide-regenerate")
def slide_regenerate(req: SlideRegenerateRequest, user: dict = Depends(get_current_user)):
    """단일 슬라이드 재생성 (비동기 task)."""
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")

    current_str = json.dumps(req.current_slide, ensure_ascii=False) if not isinstance(req.current_slide, str) else req.current_slide
    prev_str = json.dumps(req.prev_slide, ensure_ascii=False) if req.prev_slide and not isinstance(req.prev_slide, str) else (req.prev_slide or "null")
    next_str = json.dumps(req.next_slide, ensure_ascii=False) if req.next_slide and not isinstance(req.next_slide, str) else (req.next_slide or "null")

    task_id = create_task(
        user_id=user["id"], endpoint="/slide-regenerate", model=model,
        title="슬라이드 재생성",
    )
    run_analysis_task(task_id, "slide_regenerate", api_key, model,
                      current_slide=current_str, prev_slide=prev_str,
                      next_slide=next_str, instruction=req.instruction)
    log_usage(user["id"], "/slide-regenerate", model)
    return {"task_id": task_id}


@router.post("/slide-outline")
def slide_outline(req: AnalysisRequest, user: dict = Depends(get_current_user)):
    """PPT 아웃라인만 생성 (Phase 1). 사용자가 편집 후 /slides-from-outline으로 상세 생성."""
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")

    task_id = create_task(
        user_id=user["id"], endpoint="/slide-outline", model=model,
        title="PPT 아웃라인 생성",
    )
    run_analysis_task(task_id, "slide_outline", api_key, model, **req.kwargs)
    log_usage(user["id"], "/slide-outline", model)
    return {"task_id": task_id}


@router.post("/slides-from-outline")
def slides_from_outline(req: SlidesFromOutlineRequest, user: dict = Depends(get_current_user)):
    """편집된 아웃라인에서 상세 슬라이드 생성 (Phase 2)."""
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")

    # Load project docs for file_context
    file_context = ""
    if req.project_name:
        try:
            import core_logic
            file_context, _ = core_logic.parse_all_files(
                [], saved_files=req.selected_docs,
                api_key=api_key
            ) if req.selected_docs else ("", "")
        except Exception:
            pass

    task_id = create_task(
        user_id=user["id"], endpoint="/slides-from-outline", model=model,
        title="슬라이드 상세 생성",
    )
    run_analysis_task(task_id, "slides_from_outline", api_key, model,
                      outline=req.outline,
                      file_context=file_context,
                      context_text=req.context_text)
    log_usage(user["id"], "/slides-from-outline", model)
    return {"task_id": task_id}


@router.post("/update-pptx-history")
async def update_pptx_history(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    """기존 PPTX 파일의 투자이력(날짜/수치) 업데이트."""
    import core_im

    content = await file.read()
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")

    try:
        updated_bytes = core_im.update_pptx_history(content, api_key, model)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"업데이트 실패: {str(e)}")

    if not updated_bytes:
        raise HTTPException(status_code=400, detail="업데이트할 내용이 없습니다.")

    return Response(
        content=updated_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=updated_{file.filename}"},
    )


@router.get("/tasks/{task_id}")
def get_task_status(task_id: str, user: dict = Depends(get_current_user)):
    task = get_task(task_id)
    if not task:
        return {"task_id": task_id, "status": "error", "error": "Task not found", "result": None}
    resp = {
        "task_id": task_id,
        "status": task["status"],
        "result": task["result"],
        "error": task["error"],
    }
    if task.get("batch_progress"):
        resp["batch_progress"] = task["batch_progress"]
    return resp


# ========================================
# Generation History
# ========================================

@router.get("/history")
def list_history(limit: int = 50, offset: int = 0, user: dict = Depends(get_current_user)):
    """사용자의 생성 이력 목록을 반환한다."""
    from backend.database import list_generations
    rows = list_generations(user["id"], limit, offset)
    return {"items": rows, "limit": limit, "offset": offset}


@router.get("/history/{gen_id}")
def get_history_detail(gen_id: int, user: dict = Depends(get_current_user)):
    """특정 생성 이력의 상세를 반환한다."""
    from backend.database import get_generation
    item = get_generation(gen_id, user["id"])
    if not item:
        raise HTTPException(status_code=404, detail="이력을 찾을 수 없습니다.")
    return item


@router.delete("/history/{gen_id}")
def delete_history(gen_id: int, user: dict = Depends(get_current_user)):
    """생성 이력을 삭제한다."""
    from backend.database import delete_generation
    deleted = delete_generation(gen_id, user["id"])
    if not deleted:
        raise HTTPException(status_code=404, detail="이력을 찾을 수 없습니다.")
    return {"ok": True}


# ========================================
# Vector RAG
# ========================================

@router.post("/projects/{name}/reindex")
def reindex_project(name: str, force: bool = False, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    """프로젝트 문서를 인덱싱 (BM25 + 벡터DB + Gemini Files).

    force=True면 전체 재인덱싱, 기본은 변경분만 처리.
    """
    import subprocess, sys, json as _json
    import core_rag
    api_key = _get_api_key()
    if not api_key:
        return {"success": False, "error": "API Key가 설정되지 않았습니다."}

    results = {}

    # 1. BM25 인덱싱 (인프로세스, 빠름)
    try:
        import core_rag_bm25
        docs_dict = core_rag.load_project_docs_dict(name, owner_id=user["id"])
        if docs_dict:
            bm25_result = core_rag_bm25.build_index(name, docs_dict)
            results["bm25"] = {"success": True, **bm25_result}
        else:
            results["bm25"] = {"success": False, "error": "문서 없음"}
    except Exception as e:
        results["bm25"] = {"success": False, "error": str(e)}

    # 2. Gemini File Search 업로드 (선택적)
    try:
        import core_rag_gemini
        if docs_dict:
            gemini_result = core_rag_gemini.upload_project_docs(api_key, name, docs_dict)
            results["gemini_files"] = {"success": True, **gemini_result}
    except Exception as e:
        results["gemini_files"] = {"success": False, "error": str(e)}

    # 3. 벡터DB 인덱싱 (서브프로세스, segfault 격리)
    try:
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        script = f"""
import sys, json
sys.path.insert(0, {repr(base)})
import core_rag_vector
result = core_rag_vector.index_project_all({repr(api_key)}, {repr(name)}, force={repr(force)})
print(json.dumps(result, ensure_ascii=False))
"""
        proc = subprocess.run(
            [sys.executable, "-c", script],
            capture_output=True, text=True, timeout=300, cwd=base,
        )
        if proc.returncode == 0 and proc.stdout.strip():
            results["vector"] = _json.loads(proc.stdout.strip().split("\n")[-1])
        else:
            error_msg = proc.stderr.strip() if proc.stderr else f"exit code {proc.returncode}"
            results["vector"] = {"success": False, "error": error_msg}
    except Exception as e:
        results["vector"] = {"success": False, "error": str(e)}

    return {"success": True, "results": results}


@router.get("/projects/{name}/vector-stats")
def vector_stats(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    """프로젝트 벡터 인덱스 통계 (파일시스템 기반, ChromaDB 직접 열지 않음)."""
    try:
        local_data = os.path.join(os.environ.get("LOCALAPPDATA", os.path.expanduser("~")), "GEMintern")
        vdb_dir = os.path.join(local_data, "vectordb")
        # ChromaDB 디렉토리 존재 여부로 인덱싱 상태 판단
        if os.path.isdir(vdb_dir) and any(
            f for f in os.listdir(vdb_dir) if not f.startswith(".")
        ):
            return {"total_chunks": -1, "documents": -1, "indexed": True,
                    "note": "벡터 DB 존재 (상세 통계는 재인덱싱 시 확인)"}
        return {"total_chunks": 0, "documents": 0, "indexed": False}
    except Exception as e:
        return {"total_chunks": 0, "documents": 0, "indexed": False, "error": str(e)}


@router.get("/projects/{name}/sync-status")
def doc_sync_status(name: str, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    """파일 목록과 RAG DB 동기화 상태 비교."""
    import core_rag
    try:
        # RAG storage에 저장된 문서 목록 (indexed .md files)
        indexed_names = set(core_rag.get_indexed_doc_names(name, owner_id=user["id"]))

        # 실제 docs 디렉토리의 .md 파일
        storage = core_rag._get_storage_name(name, owner_id=user["id"])
        docs_dir = core_rag._get_project_docs_dir(storage)
        disk_files = {}
        if os.path.isdir(docs_dir):
            for f in sorted(os.listdir(docs_dir)):
                if f.endswith(".md"):
                    fpath = os.path.join(docs_dir, f)
                    stat = os.stat(fpath)
                    doc_name = f[:-3]  # remove .md
                    disk_files[doc_name] = {
                        "size": stat.st_size,
                        "modified": stat.st_mtime,
                    }

        disk_names = set(disk_files.keys())

        docs = []
        for doc_name in sorted(disk_names | indexed_names):
            on_disk = doc_name in disk_names
            in_index = doc_name in indexed_names
            info = disk_files.get(doc_name, {})
            status = "synced" if on_disk and in_index else "disk_only" if on_disk else "index_only"
            docs.append({
                "name": doc_name,
                "status": status,
                "size": info.get("size", 0),
                "modified": info.get("modified", 0),
            })

        return {
            "total_disk": len(disk_names),
            "total_indexed": len(indexed_names),
            "synced": len(disk_names & indexed_names),
            "disk_only": len(disk_names - indexed_names),
            "index_only": len(indexed_names - disk_names),
            "docs": docs,
        }
    except Exception as e:
        return {"error": str(e), "docs": []}


class SyncDocsRequest(BaseModel):
    add: List[str] = []       # disk_only 파일을 인덱스에 추가
    remove: List[str] = []    # index_only 항목을 인덱스에서 제거


@router.post("/projects/{name}/sync-docs")
def sync_selected_docs(name: str, req: SyncDocsRequest, user: dict = Depends(get_current_user)):
    _verify_project_ownership(name, user["id"])
    """선택한 파일들의 동기화 상태를 변경."""
    import core_rag
    try:
        storage = core_rag._get_storage_name(name, owner_id=user["id"])
        indexed = set(core_rag.get_indexed_doc_names(name, owner_id=user["id"]))
        folders = core_rag._load_folders(storage)

        # disk_only → 인덱스에 추가
        for doc_name in req.add:
            indexed.add(doc_name)
            # 폴더 구조에도 추가 (어디에도 없으면 root에)
            in_any = any(doc_name in docs for docs in folders.values())
            if not in_any:
                root = folders.setdefault(core_rag.ROOT_FOLDER, [])
                root.append(doc_name)

        # index_only → 인덱스에서 제거
        for doc_name in req.remove:
            indexed.discard(doc_name)
            # 폴더 구조에서도 제거
            for folder_docs in folders.values():
                if doc_name in folder_docs:
                    folder_docs.remove(doc_name)

        core_rag._save_indexed_docs(storage, list(indexed))
        core_rag._save_folders(storage, folders)

        return {
            "success": True,
            "added": len(req.add),
            "removed": len(req.remove),
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


# ========================================
# OCR
# ========================================

@router.post("/ocr")
async def ocr_files(files: List[UploadFile] = File(...), engine: str = "gemini", user: dict = Depends(get_current_user)):
    """이미지/PDF에서 텍스트 추출 (Gemini Vision / Claude Vision / Document AI)."""
    import fitz
    api_key = _get_api_key()
    log_usage(user["id"], "/ocr")
    results = []

    for f in files:
        data = await f.read()
        ext = os.path.splitext(f.filename)[1].lower()
        text = ""

        try:
            if ext == '.pdf':
                doc = fitz.open(stream=data, filetype="pdf")
                if engine == "claude":
                    text = _ocr_pdf_with_claude(doc)
                elif engine == "gemini" and api_key:
                    import ocr as ocr_module
                    text = ocr_module.extract_pdf_with_gemini_ocr(doc, api_key)
                elif engine == "docai":
                    try:
                        import utils_docai
                        text = utils_docai.process_document(data)
                    except Exception:
                        import ocr as ocr_module
                        text = ocr_module.extract_pdf_with_ocr(doc)
                else:
                    import ocr as ocr_module
                    text = ocr_module.extract_pdf_with_ocr(doc)
                doc.close()
            elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.webp'):
                if engine == "claude":
                    text = _ocr_image_with_claude(data, ext)
                elif engine == "gemini" and api_key:
                    from google.genai import types
                    client = AIClient(api_key=api_key)
                    mime_map = {'.png': 'image/png', '.jpg': 'image/jpeg',
                                '.jpeg': 'image/jpeg', '.tiff': 'image/tiff',
                                '.bmp': 'image/bmp', '.webp': 'image/webp'}
                    mime = mime_map.get(ext, 'image/png')
                    response = client.models.generate_content(
                        model="gemini-3-pro-preview",
                        contents=[
                            types.Part.from_bytes(data=data, mime_type=mime),
                            "이 이미지의 내용을 텍스트로 추출해줘. 표는 Markdown 표 문법으로 변환해줘. 서론 없이 결과만 출력해."
                        ],
                        config=types.GenerateContentConfig(temperature=0.0)
                    )
                    text = response.text.strip() if response.text else ""
                elif engine == "docai":
                    try:
                        import utils_docai
                        text = utils_docai.process_document(data)
                    except Exception as e:
                        text = f"Document AI 오류: {e}"
                else:
                    text = "API 키가 필요합니다."
            else:
                text = f"지원하지 않는 형식: {ext}"
        except Exception as e:
            text = f"처리 오류: {e}"

        results.append({"filename": f.filename, "text": text})

    return {"results": results}


def _ocr_image_with_claude(data: bytes, ext: str) -> str:
    """Claude Vision으로 이미지 OCR."""
    import anthropic
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        # Fallback to Gemini
        return _ocr_image_with_gemini_fallback(data, ext)
    client = anthropic.Anthropic(api_key=api_key)
    mime_map = {'.png': 'image/png', '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg', '.tiff': 'image/png',
                '.bmp': 'image/png', '.webp': 'image/webp'}
    import base64
    b64 = base64.standard_b64encode(data).decode("utf-8")
    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=16384,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image", "source": {"type": "base64", "media_type": mime_map.get(ext, "image/png"), "data": b64}},
                {"type": "text", "text": "이 이미지의 내용을 텍스트로 정확히 추출해줘. 표는 Markdown 표 문법으로 변환해줘. 레이아웃 구조를 유지해줘. 서론 없이 결과만 출력해."},
            ],
        }],
    )
    return response.content[0].text.strip() if response.content else ""


def _ocr_image_with_gemini_fallback(data: bytes, ext: str) -> str:
    """Claude 키 없을 때 Gemini로 fallback."""
    api_key = _get_api_key()
    if not api_key:
        return "API 키가 설정되지 않았습니다. (ANTHROPIC_API_KEY 또는 GOOGLE_API_KEY)"
    from google.genai import types
    client = AIClient(api_key=api_key)
    mime_map = {'.png': 'image/png', '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg', '.tiff': 'image/tiff',
                '.bmp': 'image/bmp', '.webp': 'image/webp'}
    response = client.models.generate_content(
        model="gemini-3-pro-preview",
        contents=[
            types.Part.from_bytes(data=data, mime_type=mime_map.get(ext, 'image/png')),
            "이 이미지의 내용을 텍스트로 추출해줘. 표는 Markdown 표 문법으로 변환해줘. 서론 없이 결과만 출력해."
        ],
        config=types.GenerateContentConfig(temperature=0.0)
    )
    return response.text.strip() if response.text else ""


def _ocr_pdf_with_claude(doc) -> str:
    """Claude Vision으로 PDF OCR (페이지별 이미지 변환 후 처리)."""
    import anthropic
    import fitz
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        # Fallback to Gemini
        gemini_key = _get_api_key()
        if gemini_key:
            import ocr as ocr_module
            return ocr_module.extract_pdf_with_gemini_ocr(doc, gemini_key)
        return "API 키가 설정되지 않았습니다."

    client = anthropic.Anthropic(api_key=api_key)
    import base64
    page_results = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        page_text = page.get_text().strip()

        # 텍스트가 충분하면 OCR 건너뜀
        if len(page_text) >= 50:
            page_results.append(f"[Page {page_num + 1}]\n{page_text}")
            continue

        # 이미지로 변환 (2x 해상도)
        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
        img_bytes = pix.tobytes("png")
        b64 = base64.standard_b64encode(img_bytes).decode("utf-8")

        try:
            response = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=8192,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64}},
                        {"type": "text", "text": "이 페이지의 내용을 텍스트로 정확히 추출해줘. 표는 Markdown 표 문법으로 변환해줘. 서론 없이 결과만 출력해."},
                    ],
                }],
            )
            ocr_text = response.content[0].text.strip() if response.content else ""
            page_results.append(f"[Page {page_num + 1} - Claude Vision]\n{ocr_text}")
        except Exception as e:
            page_results.append(f"[Page {page_num + 1} - Error: {e}]\n{page_text}")

    return "\n\n".join(page_results)


# ========================================
# Markdown to Word
# ========================================

class FreeDocRequest(BaseModel):
    instruction: str
    file_text: str = ""
    paste_text: str = ""


class MarkdownToDocxRequest(BaseModel):
    markdown: str
    filename: str = "output.docx"


# ========================================
# PDF Unlock
# ========================================

@router.post("/unlock-pdf")
async def unlock_pdf(
    file: UploadFile = File(...),
    password: str = "",
    user: dict = Depends(get_current_user),
):
    """비밀번호로 보호된 PDF의 잠금을 해제하여 반환한다."""
    from fastapi.responses import Response
    import utils

    log_usage(user["id"], "/unlock-pdf")
    pdf_bytes = await file.read()
    try:
        unlocked = utils.unlock_pdf(pdf_bytes, password)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    orig_name = os.path.splitext(file.filename)[0]
    return Response(
        content=unlocked,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="unlocked_{orig_name}.pdf"'},
    )


# ========================================
# Free-form Document Writing
# ========================================

@router.post("/freedoc/upload")
async def freedoc_upload(files: List[UploadFile] = File(...), user: dict = Depends(get_current_user)):
    """자유양식 문서작성: 파일 업로드 → 파싱된 텍스트 반환."""
    api_key = _get_api_key()
    parsed_parts = []
    warnings = []
    for f in files:
        content_bytes = await f.read()
        parsed = _parse_file_bytes(f.filename, content_bytes, api_key)
        if parsed:
            parsed_parts.append(f"--- [{f.filename}] ---\n{parsed}")
        else:
            warnings.append(f"{f.filename}: 텍스트를 추출할 수 없습니다. 스캔 PDF인 경우 이미지 품질을 확인해주세요.")
    return {"file_text": "\n\n".join(parsed_parts), "count": len(parsed_parts), "warnings": warnings}


@router.post("/freedoc/generate")
def freedoc_generate(req: FreeDocRequest, user: dict = Depends(get_current_user)):
    """자유양식 문서작성: 파싱된 텍스트 + 지시사항 → AI 문서 생성 (task)."""
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    log_usage(user["id"], "/freedoc/generate", model)

    combined = req.file_text
    if req.paste_text.strip():
        combined += f"\n\n--- [직접 입력 텍스트] ---\n{req.paste_text.strip()}"

    if not combined.strip():
        return {"error": "자료가 비어있습니다."}

    task_id = create_task(
        user_id=user["id"], endpoint="/freedoc/generate", model=model,
        title=req.instruction[:100], inputs={"instruction": req.instruction},
    )

    def _run():
        task = get_task(task_id)
        task["status"] = "generating"
        try:
            from google.genai import types
            from ai_client import AIClient
            client = AIClient(api_key=api_key)

            system_prompt = (
                "당신은 전문 문서 작성 AI입니다.\n"
                "제공된 자료를 분석하여 사용자의 지시에 맞는 문서를 작성합니다.\n\n"
                "[핵심 규칙]\n"
                "1. 제공된 자료의 내용을 충실히 반영하세요. 자료에 없는 내용을 임의로 추가하지 마세요.\n"
                "2. 서문, 인트로, 설명 문장 없이 바로 마크다운 본문(# 헤딩)으로 시작하세요.\n"
                "3. 수치, 고유명사, 날짜 등 구체적 데이터는 절대 변경하지 마세요.\n"
                "4. 논리적이고 체계적인 구조로 작성하세요.\n"
                "5. 마크다운 형식으로 출력하세요.\n"
                "6. 한국어로 작성하세요.\n"
            )

            prompt = (
                f"[사용자 지시사항]\n{req.instruction}\n\n"
                f"[제공된 자료]\n{_truncate_context(combined, _max_chars_for_model(model))}"
            )

            config = types.GenerateContentConfig(
                max_output_tokens=65536,
                temperature=0.4,
                system_instruction=system_prompt,
            )

            import time as _time
            max_retries = 3
            last_err = None
            for attempt in range(max_retries):
                try:
                    stream = client.models.generate_content_stream(
                        model=model, contents=prompt, config=config
                    )
                    full_text = ""
                    for chunk in stream:
                        text = chunk.text or "" if hasattr(chunk, "text") else ""
                        if text:
                            full_text += text
                            task["chunks"].append(text)
                    task["result"] = full_text
                    task["status"] = "complete"
                    _save_inline_task_history(task, full_text)
                    last_err = None
                    break
                except Exception as api_err:
                    err_str = str(api_err)
                    if ('503' in err_str or 'overloaded' in err_str.lower() or 'unavailable' in err_str.lower()) and attempt < max_retries - 1:
                        wait = 5 * (attempt + 1)
                        logger.warning(f"기안문 생성 API 503 (attempt {attempt + 1}/{max_retries}), {wait}초 후 재시도")
                        task["chunks"].append(f"\n[서버 과부하로 {wait}초 후 재시도 중...]\n")
                        _time.sleep(wait)
                        task["chunks"] = []  # 이전 청크 초기화
                        last_err = api_err
                        continue
                    raise
            if last_err:
                raise last_err
        except Exception as e:
            task["error"] = str(e)
            task["status"] = "error"

    import threading
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


# ========================================
# Draft Document (기안문) Writing
# ========================================

class DraftDocRequest(BaseModel):
    file_text: str = ""
    paste_text: str = ""
    instruction: str = ""


@router.post("/draftdoc/generate")
def draftdoc_generate(req: DraftDocRequest, user: dict = Depends(get_current_user)):
    """기안문 작성: 파싱된 텍스트 → AI 기안문 생성 (task)."""
    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    log_usage(user["id"], "/draftdoc/generate", model)

    combined = req.file_text
    if req.paste_text.strip():
        combined += f"\n\n--- [직접 입력 텍스트] ---\n{req.paste_text.strip()}"

    if not combined.strip():
        return {"error": "자료가 비어있습니다."}

    task_id = create_task(
        user_id=user["id"], endpoint="/draftdoc/generate", model=model,
        title="기안문 작성", inputs={"instruction": req.instruction[:200] if req.instruction else ""},
    )

    def _run():
        task = get_task(task_id)
        task["status"] = "generating"
        try:
            from google.genai import types
            from ai_client import AIClient
            client = AIClient(api_key=api_key)

            system_prompt = (
                "당신은 '기안문 작성 도우미'입니다. 사용자가 제공하는 내용과 참고 자료를 기반으로 기안문을 작성합니다.\n\n"
                "[기안문 공통 형식]\n"
                "1. **제목:** [대상회사명] ~~~의 건(부연설명)\n\n"
                "2. **내용:** 공식 문체 사용 (예: \"~하고자 하오니 재가하여 주시기 바랍니다.\")\n\n"
                "[지원하는 기안문 종류 및 지침]\n"
                "**사전 동의:** 피투자회사의 요청사항에 대해 사전 동의 진행 시 작성.\n"
                "**의결권 행사:** 주주총회 또는 이사회 의결권 행사 시 작성.\n"
                "**NDA, LOI, LOC 등 문서 결재:** - NDA 작성 시 \"반드시\" 다음 형식을 따름 (이 경우 공통 형식의 '다음' 항목은 작성하지 않음):\n"
                "     신규 프로젝트([프로젝트명])와 관련하여 다음의 NDA를 체결하고자 하오니 승인하여 주시기 바랍니다.\n"
                "     \n"
                "     비밀유지확약서(NDA) 체결의 건\n"
                "     - 체결 대상자: [대상자 정보]\n"
                "     - 투자 검토대상: [대상 투자 또는 검토 정보]\n"
                "     - 내용: 본 건과 관련하여 비밀유지 의무\n"
                "     - 체결 기간 및 특이사항: [계약 기간 및 특이사항]\n"
                "     \n"
                "     **[NDA 체크리스트]**\n"
                "     | 체크리스트 항목 | 답변 |\n"
                "     |---|---|\n"
                "     | 1 비밀유지의무가 일방적인가, 아니면 상호적인가? | [일방 or 상호] |\n"
                "     | 2 비밀유지대상의 정의가 명확히 기재되어 있는가? | [Y/N] |\n"
                "     | 3 비밀정보의 범위 내 잔존정보가 포함되는가? | [Y/N] |\n"
                "     | 4 비밀정보의 사용 목적이 특정되어 있는가? | [Y/N] |\n"
                "     | 5 비밀유지 의무의 내용이 구체적으로 명시되어 있는가? | [Y/N] |\n"
                "     | 6 계약 종료 시, 비밀정보 반환/폐기 확인서를 제출해야 하는가? | [Y/N] |\n"
                "     | 7 계약 위반 시, 손해배상액의 규모가 구체적이지 않은가? | [Y/N] |\n"
                "     | 8 그 외 특이 조항이 존재하는가? | [Y/N] |\n\n"
                "**다음:** 주요 항목을 나열하며 **음슴체** 사용. (**단, NDA 작성 시에는 이 '다음' 항목을 절대 생성하지 말고 생략하시오.**)\n"
                "   - **목적:** 기안 목적 명확화.\n"
                "   - **주요 내용 요약:** 주요 내용 정리.\n"
                "   - **의견:** 분석적 의견과 향후 조치사항 등 능동적이고 확정적인 어조로 작성.\n\n"
                "**붙임:** 첨부 자료 명시. 양식은\n"
                "붙임1. \n"
                "붙임2. \n"
                "~\n"
                "붙임n.[첨부파일제목].끝.\n\n"
                "[주의사항]\n"
                "- 여러 파일이 제공된 경우 내용을 종합하여 하나의 완결된 기안문을 작성하십시오.\n"
                "- 이미지가 포함된 경우 이미지 내의 텍스트를 읽고 내용을 반영하십시오.\n"
                "- 마크다운 형식으로 출력하세요.\n"
                "- 서문이나 인트로 없이 바로 기안문 본문으로 시작하세요.\n"
            )

            user_instruction = req.instruction.strip() if req.instruction.strip() else "제공된 자료를 바탕으로 기안문을 작성해주세요."

            prompt = (
                f"[사용자 추가 요청]\n{user_instruction}\n\n"
                f"[제공된 자료]\n{_truncate_context(combined, _max_chars_for_model(model))}"
            )

            config = types.GenerateContentConfig(
                max_output_tokens=65536,
                temperature=0.3,
                system_instruction=system_prompt,
            )

            stream = client.models.generate_content_stream(
                model=model, contents=prompt, config=config
            )
            full_text = ""
            for chunk in stream:
                text = chunk.text or "" if hasattr(chunk, "text") else ""
                if text:
                    full_text += text
                    task["chunks"].append(text)
            task["result"] = full_text
            task["status"] = "complete"
            _save_inline_task_history(task, full_text)
        except Exception as e:
            task["error"] = str(e)
            task["status"] = "error"

    import threading
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.post("/markdown-to-docx")
def markdown_to_docx(req: MarkdownToDocxRequest, user: dict = Depends(get_current_user)):
    """마크다운 텍스트를 Word 문서로 변환."""
    from fastapi.responses import Response
    import utils
    try:
        docx_bytes = utils.create_docx(req.markdown)
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Word 변환 오류: {e}")
    return Response(
        content=docx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{req.filename}"'},
    )


# ========================================
# Document Updater
# ========================================

_doc_updater_sessions: Dict[str, str] = {}  # session_id -> temp_dir


class DocUpdaterRunRequest(BaseModel):
    session_id: str
    supplementary_text: str = ""
    instruction: str
    mode: str = "full"


@router.post("/doc-updater/upload-original")
async def doc_updater_upload_original(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    """원본 문서 업로드 → 인덱싱 프리뷰 반환."""
    import core_doc_updater

    session_id = str(uuid.uuid4())[:8]
    temp_dir = tempfile.mkdtemp(prefix="gemintern_docupd_")
    _doc_updater_sessions[session_id] = temp_dir

    file_path = os.path.join(temp_dir, file.filename)
    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)

    try:
        doc_type, indexed = core_doc_updater.index_document(file_path)
        preview = core_doc_updater.format_document_map(indexed)
    except Exception as e:
        return {"error": f"문서 인덱싱 오류: {e}"}

    return {
        "session_id": session_id,
        "filename": file.filename,
        "doc_type": doc_type,
        "paragraph_count": len(indexed),
        "preview": preview[:5000],
    }


@router.post("/doc-updater/{session_id}/supplementary")
async def doc_updater_upload_supplementary(
    session_id: str, files: List[UploadFile] = File(...), user: dict = Depends(get_current_user)
):
    """추가 자료 파일 업로드."""
    temp_dir = _doc_updater_sessions.get(session_id)
    if not temp_dir:
        return {"error": "세션이 없습니다. 원본 문서를 먼저 업로드해주세요."}

    sup_dir = os.path.join(temp_dir, "supplementary")
    os.makedirs(sup_dir, exist_ok=True)

    filenames = []
    for f in files:
        path = os.path.join(sup_dir, f.filename)
        content = await f.read()
        with open(path, "wb") as out:
            out.write(content)
        filenames.append(f.filename)

    return {"filenames": filenames, "count": len(filenames)}


@router.post("/doc-updater/run")
def doc_updater_run(req: DocUpdaterRunRequest, user: dict = Depends(get_current_user)):
    """문서 업데이트 실행 (백그라운드)."""
    temp_dir = _doc_updater_sessions.get(req.session_id)
    if not temp_dir:
        return {"error": "세션이 없습니다."}

    api_key = _get_api_key()
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    log_usage(user["id"], "/doc-updater/run", model)

    # Find original file
    original_files = [
        f for f in os.listdir(temp_dir)
        if os.path.isfile(os.path.join(temp_dir, f))
    ]
    if not original_files:
        return {"error": "원본 문서가 없습니다."}
    original_path = os.path.join(temp_dir, original_files[0])

    # Find supplementary files
    sup_dir = os.path.join(temp_dir, "supplementary")
    sup_paths = []
    if os.path.isdir(sup_dir):
        sup_paths = [
            os.path.join(sup_dir, f) for f in os.listdir(sup_dir)
            if os.path.isfile(os.path.join(sup_dir, f))
        ]

    task_id = create_task()

    def _run():
        task = get_task(task_id)
        task["status"] = "generating"
        try:
            import core_doc_updater
            output_path, summary, preview, changes = core_doc_updater.update_document(
                original_path, sup_paths, req.supplementary_text,
                req.instruction, req.mode, api_key, model,
            )
            task["result"] = json.dumps({
                "output_path": output_path,
                "output_filename": os.path.basename(output_path),
                "summary": summary,
                "preview": preview,
                "changes": changes,
            }, ensure_ascii=False)
            task["status"] = "complete"
        except Exception as e:
            task["error"] = str(e)
            task["status"] = "error"

    import threading
    threading.Thread(target=_run, daemon=True).start()
    return {"task_id": task_id}


@router.get("/doc-updater/download")
def doc_updater_download(path: str, user: dict = Depends(get_current_user)):
    """업데이트된 문서 다운로드."""
    from fastapi.responses import FileResponse
    if not os.path.exists(path):
        return {"error": "파일을 찾을 수 없습니다."}
    return FileResponse(
        path,
        filename=os.path.basename(path),
        media_type="application/octet-stream",
    )


class DocUpdaterPromoteRequest(BaseModel):
    session_id: str
    output_path: str


@router.post("/doc-updater/promote-output")
def doc_updater_promote_output(req: DocUpdaterPromoteRequest, user: dict = Depends(get_current_user)):
    """업데이트된 출력 파일을 새 원본으로 승격 (추가 수정용)."""
    import shutil
    temp_dir = _doc_updater_sessions.get(req.session_id)
    if not temp_dir:
        return {"error": "세션이 없습니다."}
    if not req.output_path or not os.path.exists(req.output_path):
        return {"error": "출력 파일을 찾을 수 없습니다."}

    # 기존 원본 파일 삭제 (supplementary 폴더 제외)
    for f in os.listdir(temp_dir):
        fp = os.path.join(temp_dir, f)
        if os.path.isfile(fp):
            os.remove(fp)

    # 업데이트된 파일을 원본 위치로 복사 (_updated 접미사 제거)
    new_name = os.path.basename(req.output_path).replace("_updated", "")
    new_path = os.path.join(temp_dir, new_name)
    shutil.copy2(req.output_path, new_path)

    import core_doc_updater
    doc_type, indexed = core_doc_updater.index_document(new_path)
    preview = core_doc_updater.format_document_map(indexed)

    return {
        "filename": new_name,
        "doc_type": doc_type,
        "paragraph_count": len(indexed),
        "preview": preview[:5000],
    }


# ──────────────────────────────────────────
# NPS (국민연금 가입 사업장 조회)
# ──────────────────────────────────────────

import sys as _sys
# Try local project copy first, then fallback to ~/nps_query
_nps_project_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_nps_home_path = os.path.join(os.path.expanduser("~"), "nps_query")
for _p in [_nps_project_path, _nps_home_path]:
    if _p not in _sys.path:
        _sys.path.insert(0, _p)

try:
    from nps_query import NpsQuery, ENDPOINTS as NPS_ENDPOINTS
    _NPS_KEY = "SEtf7JkRv42rZno75HZqEZwMESrUT5LGXWyqJulYgDT3soMupX4WBEf3FayvqhHJqHKz0KO36R63obsB6xlHqg=="
    _nps = NpsQuery(_NPS_KEY)
    _NPS_AVAILABLE = True
except ImportError:
    _NPS_AVAILABLE = False
    _nps = None

# 캐시
import time as _time
from threading import Lock as _Lock
from concurrent.futures import ThreadPoolExecutor, as_completed

_nps_cache: dict = {}
_nps_cache_lock = _Lock()
_NPS_CACHE_TTL = 300


def _nps_cache_get(key: str):
    with _nps_cache_lock:
        entry = _nps_cache.get(key)
        if entry and _time.time() - entry["t"] < _NPS_CACHE_TTL:
            return entry["v"]
        _nps_cache.pop(key, None)
        return None


def _nps_cache_set(key: str, value):
    with _nps_cache_lock:
        _nps_cache[key] = {"v": value, "t": _time.time()}
        if len(_nps_cache) > 500:
            oldest = min(_nps_cache, key=lambda k: _nps_cache[k]["t"])
            del _nps_cache[oldest]


def _nps_fetch_one(uddi: str, name, page: int, per_page: int):
    cache_key = f"{uddi}|{name}|{page}|{per_page}"
    cached = _nps_cache_get(cache_key)
    if cached is not None:
        return cached
    resp = _nps._call_api(uddi, name=name, page=page, per_page=per_page)
    _nps_cache_set(cache_key, resp)
    return resp


@router.get("/nps/periods")
def nps_periods(user: dict = Depends(get_current_user)):
    if not _NPS_AVAILABLE:
        return {"periods": [], "error": "NPS module not available"}
    return {"periods": NpsQuery.available_periods()}


@router.get("/nps/search")
def nps_search(name: str = "", year: int = None, month: int = None,
               page: int = 1, perPage: int = 50, user: dict = Depends(get_current_user)):
    if not _NPS_AVAILABLE:
        return {"error": "NPS module not available", "data": [], "total": 0}
    if not name and year is None:
        return {"error": "사업장명 또는 연도를 입력하세요", "data": [], "total": 0}

    endpoints = _nps._get_endpoints(year, month)
    if not endpoints:
        return {"data": [], "total": 0, "page": page, "perPage": perPage}

    name_param = name if name else None

    if len(endpoints) == 1:
        ym, uddi = endpoints[0]
        try:
            resp = _nps_fetch_one(uddi, name_param, page, perPage)
            return {
                "data": resp.get("data", []),
                "total": resp.get("matchCount", 0),
                "page": page,
                "perPage": perPage,
            }
        except Exception as e:
            return {"error": str(e), "data": [], "total": 0}
    else:
        all_data = []
        total = 0
        with ThreadPoolExecutor(max_workers=12) as pool:
            futures = {
                pool.submit(_nps_fetch_one, uddi, name_param, 1, perPage): ym
                for ym, uddi in endpoints
            }
            for fut in as_completed(futures):
                try:
                    resp = fut.result()
                    all_data.extend(resp.get("data", []))
                    total += resp.get("matchCount", 0)
                except Exception:
                    logger.debug("NPS parallel fetch failed for one endpoint")
        all_data.sort(key=lambda r: r.get("자료생성년월", ""))
        return {
            "data": all_data,
            "total": total,
            "page": 1,
            "perPage": len(all_data),
        }


# ──────────────────────────────────────────
# QuickMail (AI 이메일 작성)
# ──────────────────────────────────────────

class QuickMailRequest(BaseModel):
    prompt: str
    context: str = ""
    tone: str = "professional"
    language: str = "한국어"

_TONE_MAP = {
    "formal": "격식체 (합쇼체)",
    "casual": "비격식 (해요체)",
    "professional": "비즈니스 (합니다체)",
}

@router.post("/quickmail/generate")
def quickmail_generate(req: QuickMailRequest, user: dict = Depends(get_current_user)):
    """AI 이메일 생성 (task 기반 스트리밍)."""
    api_key = _get_api_key()
    if not api_key:
        return {"error": "API key not configured"}
    model = _load_settings_for_user(user["id"]).get("model_name", "gemini-2.5-flash")
    is_reply = bool(req.context.strip())
    tone_label = _TONE_MAP.get(req.tone, "비즈니스 (합니다체)")

    system_prompt = (
        "당신은 비즈니스 이메일 작성 전문가입니다.\n"
        f"{'주어진 원본 메일에 대한 답장을 작성하세요.' if is_reply else '사용자의 요청에 맞는 이메일 본문을 작성하세요.'}\n\n"
        "[필수 규칙 - 반드시 지킬 것]\n"
        f"- {tone_label} 어조로 {req.language}(으)로 작성\n"
        "- 제목(Subject)은 포함하지 말 것 (별도 요청 시에만 포함)\n"
        "- 절대로 마크다운 서식(**, *, #, - 등)을 사용하지 말 것. 순수 텍스트만 사용\n"
        "- 절대로 [이름], [회의 장소], [Your Name], [안건] 등 대괄호 플레이스홀더를 사용하지 말 것\n"
        "- 구체적 정보가 없으면 일반적이고 자연스러운 표현으로 대체하여 바로 보낼 수 있게 작성\n"
        "- 불릿 포인트 대신 자연스러운 문장으로 서술\n"
        f"{'- 답장 본문만 작성하고 원본 인용은 포함하지 말 것' if is_reply else ''}\n\n"
        "[톤 & 매너 - 공손하고 정중하게]\n"
        "- 바쁜 상대방에 대한 배려를 반드시 표현하세요\n"
        "  예: '바쁘신 와중에 죄송합니다만', '업무에 바쁘신 줄 알지만'\n"
        "- 협조/참석/회신에 대한 감사를 자연스럽게 포함하세요\n"
        "  예: '바쁘신 가운데 참석해 주셔서 감사합니다', '협조해 주셔서 대단히 감사드립니다'\n"
        "- 부탁이나 요청 시 공손한 완곡 표현을 사용하세요\n"
        "  예: '번거로우시겠지만', '혹시 가능하시다면', '양해 부탁드립니다'\n"
        "- 마무리에 안부 인사를 포함하세요\n"
        "  예: '좋은 하루 보내시기 바랍니다', '항상 건강하시길 바랍니다'\n"
        "- 전체적으로 따뜻하고 정중한 느낌을 유지하되, 과도하게 장황하지 않게 작성하세요\n"
    )

    user_message = (
        f"원본 메일:\n---\n{req.context}\n---\n\n요청: {req.prompt}"
        if is_reply else req.prompt
    )

    task_id = create_task()
    task = get_task(task_id)
    task["status"] = "generating"

    def _run():
        try:
            from ai_client import AIClient
            from google.genai import types
            client = AIClient(api_key=api_key)
            config = types.GenerateContentConfig(
                max_output_tokens=2048,
                temperature=0.7,
                system_instruction=system_prompt,
            )
            stream = client.models.generate_content_stream(
                model=model, contents=user_message, config=config,
            )
            full_text = ""
            for chunk in stream:
                text = chunk.text or "" if hasattr(chunk, "text") else ""
                if text:
                    full_text += text
                    task["chunks"].append(text)
            task["result"] = full_text
            task["status"] = "complete"
        except Exception as e:
            task["error"] = str(e)
            task["status"] = "error"

    import threading
    threading.Thread(target=_run, daemon=True).start()
    log_usage(user.get("user_id", 0), "quickmail", model)
    return {"task_id": task_id}


# ──────────────────────────────────────────
# DartWings (DART 전자공시 기업분석)
# ──────────────────────────────────────────

try:
    from dartwings import dart_service as _dw_dart
    from dartwings import stock_service as _dw_stock
    from dartwings import analysis_service as _dw_analysis
    _DW_AVAILABLE = True
    print("[DartWings] Loaded successfully from local package")
except Exception as _e:
    print(f"[DartWings] Failed to load: {_e}")
    _DW_AVAILABLE = False

from fastapi import Query as _Query
from datetime import datetime as _dt, timedelta as _td
from collections import Counter as _Counter


@router.get("/dartwings/search")
def dw_search(q: str = _Query(..., min_length=1), limit: int = _Query(20, le=50),
              user: dict = Depends(get_current_user)):
    if not _DW_AVAILABLE:
        return {"error": "DartWings module not available"}
    return _dw_dart.search_corps(q, limit=limit)


@router.get("/dartwings/popular")
def dw_popular(user: dict = Depends(get_current_user)):
    return [
        {"corpName": "삼성전자", "stockCode": "005930"},
        {"corpName": "SK하이닉스", "stockCode": "000660"},
        {"corpName": "현대차", "stockCode": "005380"},
        {"corpName": "NAVER", "stockCode": "035420"},
        {"corpName": "카카오", "stockCode": "035720"},
    ]


@router.get("/dartwings/analyze")
def dw_analyze(stockCode: str = _Query(...), user: dict = Depends(get_current_user)):
    if not _DW_AVAILABLE:
        return {"error": "DartWings module not available"}
    from fastapi import HTTPException
    corp_code = _dw_dart.stock_code_to_corp_code(stockCode)
    if not corp_code:
        raise HTTPException(404, "종목코드를 찾을 수 없습니다")

    company = _dw_dart.get_company_info(corp_code)
    current_year = _dt.today().year
    month = _dt.today().month
    start_offset = 2 if month <= 3 else 1
    years = [str(current_year - i) for i in range(start_offset, start_offset + 3)]

    financials_by_year = []
    for yr in years:
        fs = _dw_dart.get_financial_statements(corp_code, yr)
        extracted = _dw_analysis.extract_key_financials(fs)
        extracted["year"] = yr
        financials_by_year.append(extracted)
    financials_by_year.reverse()

    financials = {
        "years": [f["year"] for f in financials_by_year],
        "revenue": [f.get("revenue", 0) for f in financials_by_year],
        "operatingProfit": [f.get("operatingProfit", 0) for f in financials_by_year],
        "netIncome": [f.get("netIncome", 0) for f in financials_by_year],
        "totalAssets": [f.get("totalAssets", 0) for f in financials_by_year],
        "totalEquity": [f.get("totalEquity", 0) for f in financials_by_year],
        "totalDebt": [f.get("totalDebt", 0) for f in financials_by_year],
    }

    price = _dw_stock.get_current_price(stockCode)
    price_history = _dw_stock.get_price_history(stockCode)
    market_cap = _dw_stock.get_market_cap(stockCode)

    dividend_list = _dw_dart.get_dividend_info(corp_code, years[0])
    dps = 0
    for item in dividend_list:
        se = item.get("se", "")
        if "주당" in se and "현금" in se:
            dps = _dw_analysis.parse_amount(item.get("thstrm", "0"))
            break
    dividend_yield = round(dps / price["current"] * 100, 2) if price["current"] else 0

    return {
        "company": company,
        "financials": financials,
        "stockPrice": {
            "current": price["current"],
            "change": price["change"],
            "volume": price["volume"],
            "marketCap": market_cap,
        },
        "dividend": {"dps": dps, "dividendYield": dividend_yield},
        "priceHistory": price_history,
    }


@router.get("/dartwings/disclosure-chart")
def dw_disclosure_chart(stockCode: str = _Query(...), user: dict = Depends(get_current_user)):
    if not _DW_AVAILABLE:
        return {"error": "DartWings module not available"}
    from fastapi import HTTPException
    corp_code = _dw_dart.stock_code_to_corp_code(stockCode)
    if not corp_code:
        raise HTTPException(404, "종목코드를 찾을 수 없습니다")

    end = _dt.today().strftime("%Y%m%d")
    bgn = (_dt.today() - _td(days=365 * 3)).strftime("%Y%m%d")
    disclosures = _dw_dart.get_disclosures(corp_code, bgn, end)

    monthly_counter: _Counter = _Counter()
    type_counter: _Counter = _Counter()
    recent = []

    for d in disclosures:
        date_str = d.get("rcept_dt", "")
        report_nm = d.get("report_nm", "")
        if len(date_str) == 8:
            monthly_counter[f"{date_str[:4]}-{date_str[4:6]}"] += 1
        if any(k in report_nm for k in ("사업보고서", "분기보고서", "반기보고서")):
            dtype = "정기보고서"
        elif "주요사항" in report_nm:
            dtype = "주요사항보고"
        elif any(k in report_nm for k in ("임원", "주식등", "지분")):
            dtype = "지분공시"
        else:
            dtype = "기타공시"
        type_counter[dtype] += 1
        if len(recent) < 20:
            recent.append({
                "date": f"{date_str[:4]}-{date_str[4:6]}-{date_str[6:8]}" if len(date_str) == 8 else date_str,
                "title": report_nm, "type": dtype,
                "link": f"https://dart.fss.or.kr/dsaf001/main.do?rcpNo={d.get('rcept_no', '')}",
            })

    sorted_months = sorted(monthly_counter.keys())
    return {
        "monthly": {"labels": sorted_months, "counts": [monthly_counter[m] for m in sorted_months]},
        "byType": {"labels": list(type_counter.keys()), "counts": list(type_counter.values())},
        "recentDisclosures": recent,
    }


@router.get("/dartwings/valuation-models")
def dw_valuation_models(stockCode: str = _Query(...), user: dict = Depends(get_current_user)):
    if not _DW_AVAILABLE:
        return {"error": "DartWings module not available"}
    from fastapi import HTTPException
    corp_code = _dw_dart.stock_code_to_corp_code(stockCode)
    if not corp_code:
        raise HTTPException(404, "종목코드를 찾을 수 없습니다")

    current_year = _dt.today().year
    month = _dt.today().month
    start_offset = 2 if month <= 3 else 1
    years = [str(current_year - i) for i in range(start_offset, start_offset + 3)]

    financials_by_year = []
    for yr in years:
        fs = _dw_dart.get_financial_statements(corp_code, yr)
        extracted = _dw_analysis.extract_key_financials(fs)
        extracted["year"] = yr
        financials_by_year.append(extracted)
    financials_by_year.reverse()

    latest = financials_by_year[-1] if financials_by_year else {}
    fundamentals = _dw_stock.get_fundamentals(stockCode)
    market_cap = _dw_stock.get_market_cap(stockCode)
    price = _dw_stock.get_current_price(stockCode)

    multiples = _dw_analysis.calc_valuation_multiples(latest, market_cap)
    multiples["per"] = fundamentals["per"]
    multiples["pbr"] = fundamentals["pbr"]

    multiples_history = {"years": [f["year"] for f in financials_by_year], "roe": []}
    for f in financials_by_year:
        ni, eq = f.get("netIncome", 0), f.get("totalEquity", 0)
        multiples_history["roe"].append(round(ni / eq * 100, 2) if eq else 0)

    shares = int(market_cap / price["current"]) if price["current"] else 0
    dcf = _dw_analysis.calc_dcf(financials_by_year, market_cap, shares)
    dcf["currentPrice"] = price["current"]
    if dcf["fairValue"] and price["current"]:
        dcf["upside"] = round((dcf["fairValue"] - price["current"]) / price["current"] * 100, 1)

    srim = _dw_analysis.calc_srim(fundamentals["bps"], multiples["roe"])
    srim["currentPrice"] = price["current"]
    if srim["fairValue"] and price["current"]:
        srim["upside"] = round((srim["fairValue"] - price["current"]) / price["current"] * 100, 1)

    avg_upside, count = 0, 0
    if dcf.get("fairValue"):
        avg_upside += dcf.get("upside", 0); count += 1
    if srim.get("fairValue"):
        avg_upside += srim.get("upside", 0); count += 1
    avg_upside = avg_upside / count if count else 0
    grade = "저평가" if avg_upside > 15 else ("고평가" if avg_upside < -15 else "적정")

    return {"multiples": multiples, "multiplesHistory": multiples_history, "dcf": dcf, "srim": srim, "grade": grade}


# ========================================
# Q&A Sessions
# ========================================

@router.get("/qa/sessions")
def list_qa_sessions(project: str, user: dict = Depends(get_current_user)):
    """List Q&A sessions for a project."""
    _verify_project_ownership(project, user["id"])
    from backend.database import get_db
    with get_db() as conn:
        proj = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (project, user["id"])
        ).fetchone()
        if not proj:
            return []
        sessions = conn.execute(
            "SELECT id, title, created_at, updated_at FROM qa_sessions WHERE project_id = ? ORDER BY updated_at DESC",
            (proj["id"],)
        ).fetchall()
    return [dict(s) for s in sessions]


@router.post("/qa/sessions")
def create_qa_session(data: dict, user: dict = Depends(get_current_user)):
    """Create a new Q&A session."""
    project_name = data.get("project")
    _verify_project_ownership(project_name, user["id"])
    from backend.database import get_db
    with get_db() as conn:
        proj = conn.execute(
            "SELECT id FROM projects WHERE name = ? AND owner_id = ?",
            (project_name, user["id"])
        ).fetchone()
        if not proj:
            raise HTTPException(status_code=404, detail="프로젝트를 찾을 수 없습니다.")
        cur = conn.execute(
            "INSERT INTO qa_sessions (project_id, title) VALUES (?, ?)",
            (proj["id"], data.get("title", "새 대화"))
        )
        session_id = cur.lastrowid
    return {"id": session_id}


@router.patch("/qa/sessions/{session_id}")
def update_qa_session(session_id: int, data: dict, user: dict = Depends(get_current_user)):
    """Update session title."""
    from backend.database import get_db
    with get_db() as conn:
        title = data.get("title")
        if title:
            conn.execute(
                "UPDATE qa_sessions SET title = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?",
                (title, session_id)
            )
    return {"ok": True}


@router.delete("/qa/sessions/{session_id}")
def delete_qa_session(session_id: int, user: dict = Depends(get_current_user)):
    """Delete a Q&A session and its messages."""
    from backend.database import get_db
    with get_db() as conn:
        conn.execute("DELETE FROM qa_sessions WHERE id = ?", (session_id,))
    return {"ok": True}


@router.get("/qa/sessions/{session_id}/messages")
def get_session_messages(session_id: int, user: dict = Depends(get_current_user)):
    """Get all messages in a session."""
    from backend.database import get_db
    with get_db() as conn:
        msgs = conn.execute(
            "SELECT id, role, content, created_at FROM qa_messages WHERE session_id = ? ORDER BY created_at",
            (session_id,)
        ).fetchall()
    return [dict(m) for m in msgs]


@router.post("/qa/sessions/{session_id}/messages")
def add_session_message(session_id: int, data: dict, user: dict = Depends(get_current_user)):
    """Add a message to a session."""
    from backend.database import get_db
    with get_db() as conn:
        conn.execute(
            "INSERT INTO qa_messages (session_id, role, content) VALUES (?, ?, ?)",
            (session_id, data["role"], data["content"])
        )
        conn.execute(
            "UPDATE qa_sessions SET updated_at = CURRENT_TIMESTAMP WHERE id = ?",
            (session_id,)
        )
    return {"ok": True}
