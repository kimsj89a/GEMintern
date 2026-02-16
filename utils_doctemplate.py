"""
Document template utilities.
Extracted from ui_doctemplate.py to remove streamlit dependency.
"""
import io
import os
from google import genai

# 지원 파일 형식
SUPPORTED_FILE_TYPES = ['docx', 'pdf', 'pptx', 'xlsx', 'xls', 'txt', 'md']


def extract_text_from_file(file_path: str) -> str:
    """로컬 파일에서 텍스트 추출 (PyQt6용 - 파일 경로 기반)"""
    filename = os.path.basename(file_path).lower()
    content = ""

    try:
        if filename.endswith('.txt') or filename.endswith('.md'):
            raw_bytes = open(file_path, 'rb').read()
            for encoding in ['utf-8-sig', 'utf-8', 'cp949', 'euc-kr']:
                try:
                    content = raw_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            if not content:
                content = raw_bytes.decode('utf-8', errors='replace')

        elif filename.endswith('.docx'):
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = '\n\n'.join(paragraphs)

        elif filename.endswith('.pdf'):
            try:
                import fitz
                pdf_doc = fitz.open(file_path)
                text_parts = []
                for page in pdf_doc:
                    text_parts.append(page.get_text())
                content = '\n'.join(text_parts)
                pdf_doc.close()
            except ImportError:
                content = "[PDF 읽기 실패: PyMuPDF 필요]"

        elif filename.endswith('.pptx'):
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                content = '\n\n'.join(text_parts)
            except ImportError:
                content = "[PPT 읽기 실패: python-pptx 필요]"

        elif filename.endswith('.xlsx') or filename.endswith('.xls'):
            try:
                import pandas as pd
                df = pd.read_excel(file_path, sheet_name=None)
                text_parts = []
                for sheet_name, sheet_df in df.items():
                    text_parts.append(f"[{sheet_name}]\n{sheet_df.to_string()}")
                content = '\n\n'.join(text_parts)
            except ImportError:
                content = "[Excel 읽기 실패: pandas 필요]"

    except Exception as e:
        content = f"[파일 읽기 오류: {str(e)}]"

    return content


def analyze_document_format(raw_text: str, api_key: str, model: str) -> str:
    """AI로 문서 형식을 마크다운으로 분석"""
    client = genai.Client(api_key=api_key)

    prompt = """다음 문서의 **형식과 구조**를 분석해서 마크다운 템플릿으로 변환해주세요.

[분석 원칙]
1. 문서의 전체적인 구조(섹션, 제목, 소제목 등)를 파악
2. 각 섹션에 어떤 종류의 내용이 들어가는지 파악 (예: 회사 개요, 재무 현황, 결론 등)
3. 반복되는 패턴이나 표 형식이 있다면 구조 파악
4. 실제 내용은 {{섹션명}} 형태의 플레이스홀더로 대체

[출력 형식]
- 마크다운 형식으로 출력
- 각 섹션의 제목은 그대로 유지
- 내용이 들어갈 부분은 {{해당_섹션_설명}} 형태로 표시
- 문서의 톤앤매너, 형식적 특징도 주석으로 메모

[문서 내용]
""" + raw_text[:15000]

    try:
        response = client.models.generate_content(model=model, contents=prompt)
        return response.text.strip() if response.text else ""
    except Exception as e:
        return f"[형식 분석 오류: {str(e)}]"


def extract_content_as_markdown(raw_text: str, api_key: str, model: str) -> str:
    """AI로 콘텐츠를 마크다운으로 정리"""
    client = genai.Client(api_key=api_key)

    prompt = """다음 문서의 **내용**을 마크다운 형식으로 깔끔하게 정리해주세요.

[정리 원칙]
1. 핵심 정보와 데이터를 빠짐없이 추출
2. 논리적 구조로 재정리 (제목, 소제목, 불릿 포인트 활용)
3. 숫자, 고유명사, 날짜 등은 정확히 유지
4. 불필요한 반복이나 형식적 문구는 제거
5. 표 형식 데이터는 마크다운 테이블로 변환

[문서 내용]
""" + raw_text[:15000]

    try:
        response = client.models.generate_content(model=model, contents=prompt)
        return response.text.strip() if response.text else ""
    except Exception as e:
        return f"[콘텐츠 추출 오류: {str(e)}]"


def generate_final_document(format_md: str, content_md: str, api_key: str, model: str) -> str:
    """형식 템플릿 + 콘텐츠 = 최종 문서 생성"""
    client = genai.Client(api_key=api_key)

    prompt = f"""다음 두 가지를 결합하여 완성된 문서를 마크다운으로 작성해주세요.

[문서 형식 템플릿]
{format_md}

---

[채워 넣을 콘텐츠]
{content_md}

---

[작성 원칙]
1. 형식 템플릿의 구조와 스타일을 따름
2. 콘텐츠의 정보를 적절한 섹션에 배치
3. 플레이스홀더({{{{...}}}})를 실제 내용으로 채움
4. 콘텐츠에 없는 정보는 추가하지 않음
5. 자연스럽고 전문적인 문체 유지
6. 최종 결과는 마크다운 형식으로 출력
"""

    try:
        response = client.models.generate_content(model=model, contents=prompt)
        return response.text.strip() if response.text else ""
    except Exception as e:
        return f"[문서 생성 오류: {str(e)}]"
