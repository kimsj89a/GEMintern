"""
기존 문서 업데이트 모듈
원본 파일을 직접 수정하여 서식을 보존한 채 내용을 업데이트한다.
"""
import os
import io
import json
import copy
import shutil
from google.genai import types
from ai_client import AIClient
from prompts import DOC_UPDATER_PROMPTS


# ── 1. 문서 인덱싱: paragraph index → text 매핑 ──

def index_docx(file_path: str) -> list[dict]:
    """Word 문서의 paragraph를 인덱싱하여 [{index, text, style}] 반환."""
    from docx import Document
    doc = Document(file_path)
    result = []
    for i, para in enumerate(doc.paragraphs):
        result.append({
            "index": i,
            "text": para.text,
            "style": para.style.name if para.style else "",
        })
    return result


def index_pptx(file_path: str) -> list[dict]:
    """PPT 문서의 slide/shape/paragraph를 인덱싱."""
    from pptx import Presentation
    prs = Presentation(file_path)
    result = []
    idx = 0
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                result.append({
                    "index": idx,
                    "text": para.text,
                    "slide": si,
                    "shape_name": shape.name,
                    "para_in_shape": pi,
                })
                idx += 1
    return result


def index_text_file(file_path: str) -> list[dict]:
    """TXT/MD 파일을 줄 단위로 인덱싱."""
    raw = open(file_path, 'rb').read()
    text = ""
    for enc in ['utf-8-sig', 'utf-8', 'cp949', 'euc-kr']:
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if not text:
        text = raw.decode('utf-8', errors='replace')

    lines = text.split('\n')
    return [{"index": i, "text": line} for i, line in enumerate(lines)]


def index_document(file_path: str) -> tuple[str, list[dict]]:
    """파일 확장자에 따라 적절한 인덱서를 선택. (doc_type, indexed_paragraphs) 반환."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.docx':
        return 'docx', index_docx(file_path)
    elif ext == '.pptx':
        return 'pptx', index_pptx(file_path)
    elif ext in ('.txt', '.md'):
        return 'text', index_text_file(file_path)
    elif ext == '.pdf':
        return 'pdf', _index_pdf(file_path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")


def _index_pdf(file_path: str) -> list[dict]:
    """PDF에서 텍스트 추출하여 인덱싱."""
    import fitz
    doc = fitz.open(file_path)
    result = []
    idx = 0
    for page_num, page in enumerate(doc):
        blocks = page.get_text("blocks")
        for block in blocks:
            text = block[4].strip()
            if text:
                result.append({"index": idx, "text": text, "page": page_num})
                idx += 1
    doc.close()
    return result


def format_document_map(indexed: list[dict]) -> str:
    """인덱싱된 문서를 AI에게 전달할 문자열로 포맷."""
    lines = []
    for item in indexed:
        text_preview = item["text"][:200]
        if len(item["text"]) > 200:
            text_preview += "..."
        lines.append(f"[{item['index']}] {text_preview}")
    return "\n".join(lines)


# ── 2. 추가 자료 파싱 ──

def parse_supplementary_files(file_paths: list[str]) -> str:
    """추가 자료 파일들을 파싱하여 하나의 텍스트로 합침."""
    from utils_doctemplate import extract_text_from_file
    parts = []
    for path in file_paths:
        name = os.path.basename(path)
        text = extract_text_from_file(path)
        if text:
            parts.append(f"--- [{name}] ---\n{text}")
    return "\n\n".join(parts)


# ── 3. AI 호출 ──

def _repair_truncated_json(raw: str) -> dict:
    """잘린 JSON 응답을 복구 시도. 실패 시 None 반환."""
    # 열린 문자열 닫기: 마지막 열린 " 찾아서 닫기
    text = raw.rstrip()
    # 마지막 완전한 항목까지 잘라내기
    # updated_paragraphs 또는 new_paragraphs 배열 중간에서 잘렸을 가능성
    for attempt in range(5):
        # 각 시도마다 마지막 불완전 요소를 점점 더 제거
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass

        # 전략 1: 열린 문자열 닫고 구조 닫기
        closers = ""
        in_string = False
        brace_depth = 0
        bracket_depth = 0
        for ch in text:
            if ch == '"' and (not text or text[text.index(ch) - 1:text.index(ch)] != '\\'):
                in_string = not in_string
            elif not in_string:
                if ch == '{':
                    brace_depth += 1
                elif ch == '}':
                    brace_depth -= 1
                elif ch == '[':
                    bracket_depth += 1
                elif ch == ']':
                    bracket_depth -= 1

        if in_string:
            closers += '"'
        closers += ']' * bracket_depth
        closers += '}' * brace_depth

        try:
            return json.loads(text + closers)
        except json.JSONDecodeError:
            pass

        # 전략 2: 마지막 불완전 JSON 객체/항목 제거 후 재시도
        # 마지막 '{' 부터 잘라내기
        last_brace = text.rfind('{')
        last_comma = text.rfind(',', 0, last_brace)
        if last_comma > 0:
            text = text[:last_comma]
        else:
            break

    return None


def call_update_ai(document_map: str, supplementary: str, instruction: str,
                   mode: str, api_key: str, model: str) -> dict:
    """Gemini API를 호출하여 업데이트 지시 JSON을 받는다."""
    client = AIClient(api_key=api_key)

    if mode == "full":
        user_prompt = DOC_UPDATER_PROMPTS['full_update'].format(
            document_map=document_map,
            supplementary=supplementary,
            instruction=instruction,
        )
    else:
        user_prompt = DOC_UPDATER_PROMPTS['partial_update'].format(
            document_map=document_map,
            supplementary=supplementary,
            instruction=instruction,
        )

    config = types.GenerateContentConfig(
        system_instruction=DOC_UPDATER_PROMPTS['system'],
        temperature=0.3,
        response_mime_type="application/json",
        max_output_tokens=65536,
    )

    response = client.models.generate_content(
        model=model, contents=user_prompt, config=config
    )

    raw = response.text.strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        # 잘린 JSON 복구 시도
        repaired = _repair_truncated_json(raw)
        if repaired is not None:
            return repaired
        raise ValueError(
            f"AI 응답이 불완전한 JSON입니다 (출력이 잘렸을 수 있습니다). "
            f"문서 크기를 줄이거나 부분 수정 모드를 사용해 보세요. 원본 오류: {e}"
        )


# ── 4. 원본 파일 직접 수정 ──

def apply_updates_docx(file_path: str, updates: dict, output_path: str):
    """Word 문서의 paragraph를 직접 수정하여 output_path에 저장."""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    doc = Document(file_path)

    for item in updates.get("updated_paragraphs", []):
        idx = item["index"]
        new_text = item["new_text"]
        if 0 <= idx < len(doc.paragraphs):
            para = doc.paragraphs[idx]
            if para.runs:
                para.runs[0].text = new_text
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.text = new_text

    new_paras = sorted(updates.get("new_paragraphs", []),
                       key=lambda x: x["after_index"], reverse=True)
    for item in new_paras:
        after_idx = item["after_index"]
        text = item["text"]
        if 0 <= after_idx < len(doc.paragraphs):
            ref_para = doc.paragraphs[after_idx]
            new_para = copy.deepcopy(ref_para._element)
            for r in new_para.findall(qn('w:r')):
                new_para.remove(r)
            r_elem = OxmlElement('w:r')
            t_elem = OxmlElement('w:t')
            t_elem.text = text
            r_elem.append(t_elem)
            new_para.append(r_elem)
            ref_para._element.addnext(new_para)

    doc.save(output_path)


def apply_updates_pptx(file_path: str, updates: dict, indexed: list[dict], output_path: str):
    """PPT 문서의 shape 텍스트를 직접 수정하여 output_path에 저장."""
    from pptx import Presentation
    prs = Presentation(file_path)

    index_map = {}
    for item in indexed:
        index_map[item["index"]] = item

    for item in updates.get("updated_paragraphs", []):
        idx = item["index"]
        new_text = item["new_text"]
        if idx not in index_map:
            continue
        meta = index_map[idx]
        slide = prs.slides[meta["slide"]]
        for shape in slide.shapes:
            if shape.name == meta["shape_name"] and shape.has_text_frame:
                pi = meta["para_in_shape"]
                if pi < len(shape.text_frame.paragraphs):
                    para = shape.text_frame.paragraphs[pi]
                    if para.runs:
                        para.runs[0].text = new_text
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = new_text
                break

    prs.save(output_path)


def apply_updates_text(file_path: str, updates: dict, output_path: str):
    """TXT/MD 파일의 줄을 수정하여 output_path에 저장."""
    raw = open(file_path, 'rb').read()
    text = ""
    for enc in ['utf-8-sig', 'utf-8', 'cp949', 'euc-kr']:
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if not text:
        text = raw.decode('utf-8', errors='replace')

    lines = text.split('\n')

    for item in updates.get("updated_paragraphs", []):
        idx = item["index"]
        if 0 <= idx < len(lines):
            lines[idx] = item["new_text"]

    new_items = sorted(updates.get("new_paragraphs", []),
                       key=lambda x: x["after_index"], reverse=True)
    for item in new_items:
        after_idx = item["after_index"]
        if 0 <= after_idx < len(lines):
            lines.insert(after_idx + 1, item["text"])

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))


def apply_updates_pdf_as_docx(file_path: str, updates: dict, indexed: list[dict], output_path: str):
    """PDF는 직접 수정 불가 → 텍스트를 Word로 변환하여 저장."""
    from docx import Document
    doc = Document()

    update_map = {item["index"]: item["new_text"] for item in updates.get("updated_paragraphs", [])}

    for item in indexed:
        idx = item["index"]
        text = update_map.get(idx, item["text"])
        doc.add_paragraph(text)

    for item in updates.get("new_paragraphs", []):
        doc.add_paragraph(item["text"])

    doc.save(output_path)


# ── 5. 통합 실행 함수 ──

def update_document(original_path: str, supplementary_paths: list[str],
                    supplementary_text: str, instruction: str,
                    mode: str, api_key: str, model: str) -> tuple[str, str, str]:
    """
    문서 업데이트 전체 파이프라인.
    Returns: (output_path, summary, preview_text)
    """
    doc_type, indexed = index_document(original_path)
    document_map = format_document_map(indexed)

    supplementary = ""
    if supplementary_paths:
        supplementary = parse_supplementary_files(supplementary_paths)
    if supplementary_text:
        supplementary += f"\n\n--- [직접 입력 텍스트] ---\n{supplementary_text}"

    if not supplementary.strip():
        supplementary = "(추가 자료 없음 — 지시사항에 따라 기존 문서만 수정)"

    updates = call_update_ai(document_map, supplementary, instruction, mode, api_key, model)

    base, ext = os.path.splitext(original_path)
    if doc_type == 'pdf':
        output_path = base + "_updated.docx"
    else:
        output_path = base + "_updated" + ext

    if doc_type == 'docx':
        apply_updates_docx(original_path, updates, output_path)
    elif doc_type == 'pptx':
        apply_updates_pptx(original_path, updates, indexed, output_path)
    elif doc_type == 'text':
        apply_updates_text(original_path, updates, output_path)
    elif doc_type == 'pdf':
        apply_updates_pdf_as_docx(original_path, updates, indexed, output_path)

    summary = updates.get("summary", "업데이트 완료")
    preview_lines = [f"## 변경 사항 요약\n{summary}\n"]
    for item in updates.get("updated_paragraphs", []):
        reason = item.get("reason", "")
        preview_lines.append(f"**[{item['index']}번 문단 수정]** {reason}")
        preview_lines.append(f"> {item['new_text'][:150]}...\n")
    for item in updates.get("new_paragraphs", []):
        reason = item.get("reason", "")
        preview_lines.append(f"**[{item['after_index']}번 뒤 신규 삽입]** {reason}")
        preview_lines.append(f"> {item['text'][:150]}...\n")
    preview = "\n".join(preview_lines)

    return output_path, summary, preview
