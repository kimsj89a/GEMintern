"""
IM PPT 변환 엔진 — Noh & Partners IM Style
마크다운 형식의 IM 보고서를 NP 브랜드 PPT로 변환

Design Reference: noh-partners-im-style SKILL (Redvelvet IM v10 실측 기반)
Slide Dimensions: 11.93" × 8.50" (커스텀 비율 ≈1.40:1)
Layout: 5-Zone (A: Title, B: Section Tab, C: Governing, D: Content, E: Footer)
"""

import io
import re
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ══════════════════════════════════════════════════
# §1. 디자인 상수 — NP IM Style
# ══════════════════════════════════════════════════

# --- 브랜드 컬러 (6색 팔레트) ---
COLOR_NAVY = RGBColor(0x0C, 0x30, 0x64)       # Main: 테이블 헤더, 서브섹션 바, 파트 구분 배경
COLOR_BLUE = RGBColor(0x00, 0x5D, 0xA2)       # 보조1: 차트 2nd, 보조 강조
COLOR_SKY = RGBColor(0x00, 0xA2, 0xE8)        # 보조2: 차트 3rd, 아이콘
COLOR_GOLD = RGBColor(0xCC, 0xA0, 0x00)       # 액센트: 서브섹션 바 좌측 스트라이프, 하이라이트
COLOR_RED = RGBColor(0xC0, 0x00, 0x00)        # 강조: 핵심 수치, 경고
COLOR_DARK_GREY = RGBColor(0x40, 0x40, 0x40)  # 본문 텍스트 (#000000은 미사용)

# --- 보조 컬러 ---
COLOR_DEEP_NAVY = RGBColor(0x00, 0x1D, 0x3D)  # 슬라이드 제목 텍스트
COLOR_LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2) # 테이블 교차행, 박스 배경
COLOR_MUTED_SLATE = RGBColor(0x64, 0x74, 0x8B) # 각주, 소스, 캡션
COLOR_BORDER = RGBColor(0xBF, 0xBF, 0xBF)     # 테이블 테두리, 구분선
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)

# 하위 호환용 별칭
COLOR_PRIMARY = COLOR_NAVY
COLOR_ACCENT = COLOR_BLUE
COLOR_MID_GREY = COLOR_MUTED_SLATE
COLOR_TABLE_HEADER = COLOR_NAVY
COLOR_TABLE_ALT = COLOR_LIGHT_GREY
COLOR_CHART_BG = RGBColor(0xE6, 0xE6, 0xE6)

# --- 폰트 ---
DEFAULT_FONT = "Arial"  # pptxgenjs 호환; 한글은 OS가 맑은 고딕으로 fallback

# --- 슬라이드 크기: 11.93" × 8.50" (NP IM 커스텀) ---
SLIDE_WIDTH = Inches(11.93)
SLIDE_HEIGHT = Inches(8.50)

# --- Zone 좌표 (본문 슬라이드) ---
ZONE_B_X, ZONE_B_Y, ZONE_B_W, ZONE_B_H = 0.40, 0.24, 9.39, 0.30   # 섹션 탭
ZONE_A_X, ZONE_A_Y, ZONE_A_W, ZONE_A_H = 0.42, 0.54, 9.37, 0.53   # 타이틀
ZONE_C_X, ZONE_C_Y, ZONE_C_W, ZONE_C_H = 0.41, 1.33, 11.09, 0.53  # 거버닝 문구
ZONE_D_X, ZONE_D_Y, ZONE_D_W = 0.41, 2.05, 11.09                    # 본문 영역 시작
ZONE_D_END_Y = 7.55                                                    # 본문 영역 끝

# 서브섹션 바
SUBSEC_BAR_H = 0.32
SUBSEC_BAR_W = 5.35
SUBSEC_GOLD_W = 0.08

# 2열 레이아웃
COL_LEFT_X = 0.41
COL_RIGHT_X = 6.16
COL_W = 5.35

# 푸터
FOOTER_SOURCE_Y = 7.88
FOOTER_PAGE_Y = 8.15
FOOTER_CONF_Y = 8.23

# --- 차트 시리즈 컬러 ---
_CHART_COLORS = [
    COLOR_NAVY, COLOR_BLUE, COLOR_SKY,
    COLOR_GOLD, COLOR_RED, COLOR_DARK_GREY,
]


# ══════════════════════════════════════════════════
# §2. 헬퍼 함수
# ══════════════════════════════════════════════════

def clean_text(text):
    """Markdown 문법 제거."""
    text = text.replace('**', '').replace('*', '')
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    return text.strip()


def _set_font(run, size, bold=False, color=None, font_name=DEFAULT_FONT):
    """폰트 설정 헬퍼."""
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _set_paragraph_font(paragraph, size, bold=False, color=None, font_name=DEFAULT_FONT):
    """Paragraph 레벨 폰트 설정."""
    for run in paragraph.runs:
        _set_font(run, size, bold, color, font_name)
    paragraph.font.name = font_name
    paragraph.font.size = size
    paragraph.font.bold = bold
    if color:
        paragraph.font.color.rgb = color


# ══════════════════════════════════════════════════
# §3. Zone 헬퍼 — 모든 본문 슬라이드 공통
# ══════════════════════════════════════════════════

def _add_footer(slide, source="", page_num=0):
    """ZONE E: 푸터 — Source / 페이지번호 / 기밀표시."""
    # Source / Note
    if source:
        tb = slide.shapes.add_textbox(
            Inches(0.42), Inches(FOOTER_SOURCE_Y), Inches(9.68), Inches(0.31))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = source
        _set_font(run, Pt(9), color=COLOR_MUTED_SLATE)

    # 페이지 번호 (center)
    if page_num > 0:
        tb2 = slide.shapes.add_textbox(
            Inches(5.75), Inches(FOOTER_PAGE_Y), Inches(0.43), Inches(0.32))
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = str(page_num)
        _set_font(run2, Pt(9), color=COLOR_DARK_GREY)

    # 기밀표시 (right)
    tb3 = slide.shapes.add_textbox(
        Inches(9.32), Inches(FOOTER_CONF_Y), Inches(2.19), Inches(0.15))
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    run3 = p3.add_run()
    run3.text = "Strictly Private and Confidential"
    _set_font(run3, Pt(8), color=COLOR_DARK_GREY)


def _add_section_tab(slide, section_label):
    """ZONE B: 섹션 탭 — 슬라이드 최상단 파트 표시."""
    if not section_label:
        return
    tb = slide.shapes.add_textbox(
        Inches(ZONE_B_X), Inches(ZONE_B_Y), Inches(ZONE_B_W), Inches(ZONE_B_H))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = section_label
    _set_font(run, Pt(9), color=COLOR_NAVY)


def _add_title(slide, title):
    """ZONE A: 슬라이드 타이틀 (24pt Bold Deep Navy)."""
    tb = slide.shapes.add_textbox(
        Inches(ZONE_A_X), Inches(ZONE_A_Y), Inches(ZONE_A_W), Inches(ZONE_A_H))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean_text(title)
    _set_font(run, Pt(24), bold=True, color=COLOR_DEEP_NAVY)


def _add_governing(slide, text):
    """ZONE C: 거버닝 문구 — 핵심 메시지 1~2줄 (11pt Bold Black)."""
    if not text:
        return
    tb = slide.shapes.add_textbox(
        Inches(ZONE_C_X), Inches(ZONE_C_Y), Inches(ZONE_C_W), Inches(ZONE_C_H))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean_text(text)
    _set_font(run, Pt(11), bold=True, color=COLOR_BLACK)


def _add_subsection_bar(slide, title, x=None, y=None):
    """서브섹션 바 — Gold 스트라이프 + Navy 본체 + White 텍스트."""
    bx = x if x is not None else COL_LEFT_X
    by = y if y is not None else ZONE_D_Y

    # Gold 스트라이프 (좌측 0.08")
    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(bx), Inches(by), Inches(SUBSEC_GOLD_W), Inches(SUBSEC_BAR_H))
    gold.fill.solid()
    gold.fill.fore_color.rgb = COLOR_GOLD
    gold.line.fill.background()

    # Navy 본체
    navy = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(bx + SUBSEC_GOLD_W), Inches(by),
        Inches(SUBSEC_BAR_W - SUBSEC_GOLD_W), Inches(SUBSEC_BAR_H))
    navy.fill.solid()
    navy.fill.fore_color.rgb = COLOR_NAVY
    navy.line.fill.background()

    # 텍스트 (11pt Bold White)
    tb = slide.shapes.add_textbox(
        Inches(bx + 0.15), Inches(by), Inches(SUBSEC_BAR_W - 0.15), Inches(SUBSEC_BAR_H))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    tb.text_frame.paragraphs[0].space_before = Pt(0)
    run = p.add_run()
    run.text = clean_text(title)
    _set_font(run, Pt(11), bold=True, color=COLOR_WHITE)
    # Vertical center
    try:
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        from pptx.oxml.ns import qn
        tb.text_frame._txBody.attrib[qn('anchor')] = 'ctr'
    except Exception:
        pass


# ══════════════════════════════════════════════════
# §4. 슬라이드 생성 함수
# ══════════════════════════════════════════════════

def _create_cover_slide(prs, project_name, gp_name, date_str):
    """표지 슬라이드 — Navy 전체 배경 + 프로젝트명."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full Navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY
    bg.line.fill.background()

    # Gold accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(3.2), Inches(1.5), Emu(38100))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_GOLD
    bar.line.fill.background()

    # Project name (36pt Bold White)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = project_name if project_name else "Information Memorandum"
    _set_font(run, Pt(36), bold=True, color=COLOR_WHITE)

    # Subtitle: "Information Memorandum" (20pt light)
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(10.5), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Information Memorandum"
    _set_font(run2, Pt(20), color=RGBColor(0xB4, 0xC6, 0xE7))

    # GP name (14pt White)
    if gp_name:
        tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(10.5), Inches(0.5))
        p3 = tb3.text_frame.paragraphs[0]
        run3 = p3.add_run()
        run3.text = gp_name
        _set_font(run3, Pt(14), color=COLOR_WHITE)

    # Date (14pt light)
    if not date_str:
        date_str = datetime.date.today().strftime("%Y년 %m월")
    tb4 = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(10.5), Inches(0.5))
    p4 = tb4.text_frame.paragraphs[0]
    run4 = p4.add_run()
    run4.text = date_str
    _set_font(run4, Pt(14), color=RGBColor(0xB4, 0xC6, 0xE7))

    # Confidential (10pt muted)
    tb5 = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(10.5), Inches(0.4))
    p5 = tb5.text_frame.paragraphs[0]
    run5 = p5.add_run()
    run5.text = "STRICTLY PRIVATE AND CONFIDENTIAL"
    _set_font(run5, Pt(10), color=RGBColor(0x96, 0xAA, 0xD2))

    return slide


def _create_disclaimer_slide(prs):
    """Disclaimer 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title (24pt Bold Deep Navy)
    tb = slide.shapes.add_textbox(
        Inches(ZONE_A_X), Inches(ZONE_A_Y), Inches(ZONE_A_W), Inches(ZONE_A_H))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Disclaimer"
    _set_font(run, Pt(24), bold=True, color=COLOR_DEEP_NAVY)

    # Navy line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.42), Inches(1.20), Inches(11.09), Emu(19050))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_NAVY
    line.line.fill.background()

    # Disclaimer text (11pt Dark Grey)
    disclaimer_text = (
        "본 자료는 정보 제공 목적으로만 작성되었으며, 투자 권유를 구성하지 않습니다.\n\n"
        "본 자료에 포함된 정보는 신뢰할 수 있다고 판단되는 자료에 기초하였으나, "
        "그 정확성이나 완전성에 대해 보증하지 않습니다.\n\n"
        "본 자료는 기밀 정보를 포함하고 있으며, 수신자의 사전 서면 동의 없이 "
        "제3자에게 공개, 복사, 배포할 수 없습니다.\n\n"
        "투자에는 원금 손실의 위험이 있으며, 과거 실적이 미래 수익을 보장하지 않습니다.\n\n"
        "본 자료에 포함된 미래 전망에 관한 진술은 현재 시점의 추정치이며, "
        "실제 결과는 다양한 요인에 의해 달라질 수 있습니다."
    )
    tb2 = slide.shapes.add_textbox(
        Inches(0.42), Inches(1.50), Inches(11.09), Inches(5.5))
    tf = tb2.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    run2 = p2.add_run()
    run2.text = disclaimer_text
    _set_font(run2, Pt(11), color=COLOR_DARK_GREY)
    p2.line_spacing = Pt(18)

    _add_footer(slide, "", 1)
    return slide


def _create_section_divider(prs, section_num, section_title, page_num=0):
    """파트 구분 슬라이드 — Navy 전체 배경 + 화이트 텍스트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Navy 전체 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY
    bg.line.fill.background()

    # 장식 라인 (short navy + long grey)
    short_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.17), Inches(1.86), Inches(0.79), Emu(25400))
    short_line.fill.solid()
    short_line.fill.fore_color.rgb = COLOR_GOLD
    short_line.line.fill.background()

    long_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.21), Inches(2.10), Inches(5.72), Emu(12700))
    long_line.fill.solid()
    long_line.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
    long_line.line.fill.background()

    # 섹션 제목 (텍스트)
    title_text = clean_text(section_title)
    if section_num:
        title_text = f"{section_num}. {title_text}"

    tb = slide.shapes.add_textbox(
        Inches(6.17), Inches(2.30), Inches(5.66), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    _set_font(run, Pt(36), bold=True, color=COLOR_WHITE)

    return slide


def _create_content_slide(prs, title, items, section_label="", page_num=0,
                          governing="", subsections=None, source=""):
    """본문 슬라이드 — Zone A~E 구조.

    Args:
        title: ZONE A 슬라이드 제목
        items: ZONE D 본문 항목 (리스트)
        section_label: ZONE B 섹션 탭 텍스트
        page_num: 페이지 번호
        governing: ZONE C 거버닝 문구 (선택)
        subsections: [{title, x}] 서브섹션 바 목록 (선택)
        source: ZONE E 소스 텍스트 (선택)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ZONE B — 섹션 탭
    _add_section_tab(slide, section_label)

    # ZONE A — 타이틀
    _add_title(slide, title)

    # ZONE C — 거버닝 문구 (첫 아이템이 Bold로 시작하면 자동 추출)
    gov_text = governing
    content_items = list(items) if items else []
    if not gov_text and content_items:
        # 첫 줄이 ### 이 아니고 - 도 아닌 일반 텍스트면 거버닝으로 사용
        first = content_items[0].strip()
        if (not first.startswith('### ') and not first.startswith('- ')
                and not first.startswith('* ') and not first.startswith('|')
                and len(first) > 20):
            gov_text = first
            content_items = content_items[1:]
    _add_governing(slide, gov_text)

    # ZONE D — 서브섹션 바
    if subsections:
        for sub in subsections:
            _add_subsection_bar(slide, sub.get('title', ''),
                                sub.get('x'), sub.get('y'))

    # 서브헤딩 감지 → 자동 서브섹션 바 생성
    auto_subsec_y = ZONE_D_Y
    has_subsec = False
    for item in content_items:
        stripped = item.strip()
        if stripped.startswith('### '):
            sub_title = clean_text(stripped.replace('### ', ''))
            _add_subsection_bar(slide, sub_title, COL_LEFT_X, auto_subsec_y)
            auto_subsec_y += SUBSEC_BAR_H + 0.05
            has_subsec = True

    # ZONE D — 본문 텍스트
    content_y = ZONE_D_Y + (SUBSEC_BAR_H + 0.10 if has_subsec or subsections else 0)
    content_h = ZONE_D_END_Y - content_y

    body_items = [item for item in content_items if not item.strip().startswith('### ')]

    if body_items:
        tb = slide.shapes.add_textbox(
            Inches(ZONE_D_X), Inches(content_y),
            Inches(ZONE_D_W), Inches(content_h))
        tf = tb.text_frame
        tf.word_wrap = True

        for i, item in enumerate(body_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            cleaned = clean_text(item)

            if item.strip().startswith('- ') or item.strip().startswith('* '):
                text = cleaned.lstrip('- ').lstrip('* ')
                # key: value 패턴
                bold_match = re.match(r'^([^:]+):\s*(.+)$', text)
                if bold_match:
                    run_bold = p.add_run()
                    run_bold.text = bold_match.group(1) + ": "
                    _set_font(run_bold, Pt(11), bold=True, color=COLOR_DARK_GREY)
                    run_normal = p.add_run()
                    run_normal.text = bold_match.group(2)
                    _set_font(run_normal, Pt(11), color=COLOR_DARK_GREY)
                else:
                    run = p.add_run()
                    run.text = "  •  " + text
                    _set_font(run, Pt(11), color=COLOR_DARK_GREY)
                p.space_before = Pt(3)
                p.space_after = Pt(3)
            else:
                run = p.add_run()
                run.text = cleaned
                _set_font(run, Pt(11), color=COLOR_DARK_GREY)
                p.space_before = Pt(2)
                p.space_after = Pt(2)

    # ZONE E — 푸터
    _add_footer(slide, source, page_num)
    return slide


def _create_table_slide(prs, title, table_data, section_label="", page_num=0, source=""):
    """테이블 슬라이드 — NP 테이블 스타일."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ZONE B
    _add_section_tab(slide, section_label)
    # ZONE A
    _add_title(slide, title)

    if not table_data or len(table_data) < 2:
        _add_footer(slide, source, page_num)
        return slide

    rows = len(table_data)
    cols = len(table_data[0])
    cols = min(cols, 10)
    rows = min(rows, 18)

    # 테이블 위치: Zone D 시작점
    table_x = Inches(ZONE_D_X)
    table_y = Inches(ZONE_D_Y)
    table_w = Inches(ZONE_D_W)
    row_h = 0.32
    table_h = Inches(min(rows * row_h, ZONE_D_END_Y - ZONE_D_Y))

    table_shape = slide.shapes.add_table(rows, cols, table_x, table_y, table_w, table_h)
    table = table_shape.table

    # 컬럼 너비 균등 배분
    col_width = int(table_w / cols)
    for c in range(cols):
        table.columns[c].width = col_width

    for r in range(min(rows, len(table_data))):
        for c in range(min(cols, len(table_data[r]))):
            cell = table.cell(r, c)
            cell.text = clean_text(table_data[r][c])

            # 셀 패딩
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = DEFAULT_FONT
                paragraph.font.size = Pt(11)

                if r == 0:
                    # 헤더: Navy 배경, White Bold, center
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLOR_WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = COLOR_DARK_GREY
                    # 첫 열: left, 나머지: center (숫자면 right)
                    if c == 0:
                        paragraph.alignment = PP_ALIGN.CENTER
                    else:
                        paragraph.alignment = PP_ALIGN.CENTER

            # 셀 배경
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_NAVY
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GREY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

    _add_footer(slide, source, page_num)
    return slide


# ══════════════════════════════════════════════════
# §5. 차트 관련
# ══════════════════════════════════════════════════

def _detect_chart_type(chart_desc):
    """차트 설명에서 차트 타입 추론."""
    desc_lower = chart_desc.lower()
    if any(k in desc_lower for k in ['pie', '파이', '비중', '구성비', '비율']):
        return 'pie'
    if any(k in desc_lower for k in ['line', '추이', '추세', '변화', '성장', 'trend', 'growth']):
        return 'line'
    return 'bar'


_CHART_TYPE_MAP = {
    'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'line': XL_CHART_TYPE.LINE_MARKERS,
    'pie': XL_CHART_TYPE.PIE,
}


def _create_chart_slide(prs, title, chart_desc, table_data=None,
                        section_label="", page_num=0):
    """차트 슬라이드 — 선행 테이블 데이터 있으면 실제 차트, 없으면 플레이스홀더."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_tab(slide, section_label)
    _add_title(slide, title)

    chart_type = _detect_chart_type(chart_desc)

    # 실제 차트 생성 시도
    if table_data and len(table_data) >= 2:
        categories, series_list = _table_to_chart_data(table_data, chart_type)
        if categories and series_list:
            xl_type = _CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for s in series_list:
                chart_data.add_series(s['name'], s['values'])

            chart_frame = slide.shapes.add_chart(
                xl_type,
                Inches(ZONE_D_X), Inches(ZONE_D_Y),
                Inches(ZONE_D_W), Inches(4.8),
                chart_data)
            chart = chart_frame.chart
            chart.has_legend = len(series_list) > 1
            if chart.has_legend:
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(9)
                chart.legend.font.name = DEFAULT_FONT
            chart.font.name = DEFAULT_FONT
            chart.font.size = Pt(9)

            # NP 차트 컬러 적용
            try:
                plot = chart.plots[0]
                for idx, series in enumerate(plot.series):
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = _CHART_COLORS[idx % len(_CHART_COLORS)]
            except Exception:
                pass

            # 차트 설명
            if chart_desc:
                tb2 = slide.shapes.add_textbox(
                    Inches(ZONE_D_X), Inches(ZONE_D_Y + 5.0),
                    Inches(ZONE_D_W), Inches(0.4))
                p2 = tb2.text_frame.paragraphs[0]
                run2 = p2.add_run()
                run2.text = clean_text(chart_desc)
                _set_font(run2, Pt(9), color=COLOR_MUTED_SLATE)

            _add_footer(slide, "", page_num)
            return slide

    # 플레이스홀더
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(ZONE_D_Y), Inches(8.93), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CHART_BG
    box.line.color.rgb = COLOR_BORDER
    box.line.width = Pt(1)

    tb2 = slide.shapes.add_textbox(
        Inches(2.5), Inches(3.5), Inches(6.93), Inches(1.5))
    tf = tb2.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"[Chart: {chart_type.upper()}]\n{clean_text(chart_desc)}"
    _set_font(run2, Pt(14), color=COLOR_MUTED_SLATE)

    _add_footer(slide, "", page_num)
    return slide


# ══════════════════════════════════════════════════
# §6. 파싱 유틸리티
# ══════════════════════════════════════════════════

def _table_to_chart_data(table_data, chart_type='bar'):
    """테이블 데이터(2D 배열)를 차트 데이터로 변환."""
    if not table_data or len(table_data) < 2:
        return None, None

    headers = table_data[0]
    rows = table_data[1:]

    numeric_cols = []
    for ci in range(1, len(headers)):
        numeric_count = 0
        for row in rows:
            if ci < len(row):
                val = re.sub(r'[,%원$₩억만]', '', str(row[ci]).strip())
                try:
                    float(val)
                    numeric_count += 1
                except ValueError:
                    pass
        if numeric_count >= len(rows) * 0.5:
            numeric_cols.append(ci)

    if not numeric_cols:
        return None, None

    categories = [str(row[0]).strip() for row in rows if row]
    series_list = []
    for ci in numeric_cols:
        values = []
        for row in rows:
            if ci < len(row):
                val = re.sub(r'[,%원$₩억만]', '', str(row[ci]).strip())
                try:
                    values.append(float(val))
                except ValueError:
                    values.append(0)
            else:
                values.append(0)
        series_list.append({'name': str(headers[ci]).strip(), 'values': values})

    return categories, series_list


def _parse_table(lines):
    """마크다운 테이블 라인들을 2D 배열로 파싱."""
    table_data = []
    for line in lines:
        line = line.strip()
        if not line.startswith('|'):
            continue
        if re.match(r'^\|[\s\-:]+\|', line):
            continue
        cells = [c.strip() for c in line.split('|')[1:-1]]
        if cells:
            table_data.append(cells)
    return table_data


def _extract_section_num(title):
    """제목에서 섹션 번호 추출."""
    match = re.match(r'^([IVX]+\.?|[0-9]+\.?)\s', title.strip())
    if match:
        return match.group(1).rstrip('.')
    return ""


# ══════════════════════════════════════════════════
# §7. 메인 진입점
# ══════════════════════════════════════════════════

def create_im_ppt(markdown_text, project_name="", gp_name="", date_str="",
                  template_path=None):
    """
    IM 마크다운을 NP 스타일 PPT로 변환하는 메인 함수.

    Args:
        markdown_text: 마크다운 형식의 IM 텍스트
        project_name: 프로젝트명 (표지용)
        gp_name: GP사명 (표지용)
        date_str: 날짜 (표지용)
        template_path: .pptx 템플릿 파일 경로 (선택)

    Returns:
        bytes: PPT 파일 바이트
    """
    import os
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    # Cover slide
    _create_cover_slide(prs, project_name, gp_name, date_str)

    # Disclaimer slide
    _create_disclaimer_slide(prs)

    page_num = 2
    current_section_label = ""
    current_slide_title = ""
    current_items = []
    is_first_h1 = True
    last_table_data = None

    lines = markdown_text.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped or stripped == '---':
            i += 1
            continue

        # H1: Section divider
        if stripped.startswith('# ') and not stripped.startswith('## '):
            if current_slide_title and current_items:
                page_num += 1
                _create_content_slide(prs, current_slide_title, current_items,
                                      current_section_label, page_num)
                current_items = []
                current_slide_title = ""

            h1_text = stripped[2:].strip()

            if is_first_h1:
                is_first_h1 = False
            else:
                section_num = _extract_section_num(h1_text)
                section_title = re.sub(r'^[IVX]+\.?\s*|^[0-9]+\.?\s*', '', h1_text).strip()
                page_num += 1
                _create_section_divider(prs, section_num, section_title or h1_text, page_num)

            current_section_label = h1_text
            i += 1
            continue

        # H2: New content slide
        if stripped.startswith('## '):
            if current_slide_title and current_items:
                page_num += 1
                _create_content_slide(prs, current_slide_title, current_items,
                                      current_section_label, page_num)
                current_items = []

            current_slide_title = stripped[3:].strip()
            i += 1
            continue

        # Table detection
        if stripped.startswith('|'):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                table_lines.append(lines[i])
                i += 1

            table_data = _parse_table(table_lines)
            if table_data and len(table_data) >= 2:
                if current_items:
                    page_num += 1
                    _create_content_slide(prs, current_slide_title, current_items,
                                          current_section_label, page_num)
                    current_items = []

                page_num += 1
                _create_table_slide(prs, current_slide_title, table_data,
                                    current_section_label, page_num)
                last_table_data = table_data
                current_slide_title = ""
            continue

        # Chart placeholder
        chart_match = re.match(r'\[차트:\s*(.+?)\]', stripped)
        if chart_match:
            if current_items:
                page_num += 1
                _create_content_slide(prs, current_slide_title, current_items,
                                      current_section_label, page_num)
                current_items = []

            page_num += 1
            _create_chart_slide(prs, current_slide_title, chart_match.group(1),
                                table_data=last_table_data,
                                section_label=current_section_label, page_num=page_num)
            last_table_data = None
            i += 1
            continue

        # Regular content
        current_items.append(stripped)
        i += 1

    # Flush last slide
    if current_slide_title and current_items:
        page_num += 1
        _create_content_slide(prs, current_slide_title, current_items,
                              current_section_label, page_num)

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def im_markdown_to_slide_json(markdown_text, project_name="", gp_name="", date_str=""):
    """IM 마크다운을 NP 렌더러 호환 JSON 슬라이드 배열로 변환.

    Returns:
        dict: {"slides": [...]} NP 렌더러 호환 JSON
    """
    slides = []

    slides.append({
        "slide_type": "title",
        "title": project_name or "Information Memorandum",
        "subtitle": gp_name or "Information Memorandum",
    })

    current_section_label = ""
    current_slide_title = ""
    current_items = []
    is_first_h1 = True
    section_num_counter = 0
    last_table_data = None

    def _flush_content():
        nonlocal current_items, current_slide_title
        if current_slide_title and current_items:
            slide = {
                "slide_type": "two_column",
                "title": clean_text(current_slide_title),
                "subtitle": current_section_label,
                "left": {"items": [clean_text(item) for item in current_items]},
            }
            slides.append(slide)
            current_items = []
            current_slide_title = ""

    lines = markdown_text.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped or stripped == '---':
            i += 1
            continue

        if stripped.startswith('# ') and not stripped.startswith('## '):
            _flush_content()
            h1_text = stripped[2:].strip()

            if is_first_h1:
                is_first_h1 = False
            else:
                section_num_counter += 1
                section_num = _extract_section_num(h1_text)
                section_title = re.sub(r'^[IVX]+\.?\s*|^[0-9]+\.?\s*', '', h1_text).strip()
                slides.append({
                    "slide_type": "divider",
                    "title": section_title or h1_text,
                    "section_number": section_num or str(section_num_counter),
                })

            current_section_label = h1_text
            i += 1
            continue

        if stripped.startswith('## '):
            _flush_content()
            current_slide_title = stripped[3:].strip()
            i += 1
            continue

        if stripped.startswith('|'):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                table_lines.append(lines[i])
                i += 1

            table_data = _parse_table(table_lines)
            if table_data and len(table_data) >= 2:
                _flush_content()
                headers = table_data[0]
                rows = table_data[1:]
                slides.append({
                    "slide_type": "data_table",
                    "title": clean_text(current_slide_title) if current_slide_title else "",
                    "subtitle": current_section_label,
                    "table": {"headers": headers, "rows": rows},
                })
                last_table_data = table_data
                current_slide_title = ""
            continue

        chart_match = re.match(r'\[차트:\s*(.+?)\]', stripped)
        if chart_match:
            _flush_content()
            chart_desc = chart_match.group(1)
            chart_type = _detect_chart_type(chart_desc)

            chart_slide = {
                "slide_type": "chart_table",
                "title": clean_text(current_slide_title) if current_slide_title else chart_desc,
                "subtitle": current_section_label,
            }

            if last_table_data:
                categories, series_list = _table_to_chart_data(last_table_data, chart_type)
                if categories and series_list:
                    chart_slide["chart"] = {
                        "chart_type": chart_type,
                        "categories": categories,
                        "series": series_list,
                    }
                    chart_slide["table"] = {
                        "headers": last_table_data[0],
                        "rows": last_table_data[1:],
                    }
                last_table_data = None

            slides.append(chart_slide)
            i += 1
            continue

        current_items.append(stripped)
        i += 1

    _flush_content()

    return {"slides": slides}
