import io
import json
import re
import math
import copy
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ============================================================
# Design Constants (NP Theme — 16:9)
# ============================================================
SLIDE_WIDTH = Inches(10.0)
SLIDE_HEIGHT = Inches(5.625)

# NP Color Palette
COLOR_NAVY = RGBColor(0x1A, 0x27, 0x44)
COLOR_DARK_NAVY = RGBColor(0x11, 0x1C, 0x33)
COLOR_GOLD = RGBColor(0xC5, 0x97, 0x3B)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_OFF_WHITE = RGBColor(0xF5, 0xF6, 0xF8)
COLOR_LIGHT_GRAY = RGBColor(0xE8, 0xEA, 0xF0)
COLOR_MID_GRAY = RGBColor(0x88, 0x92, 0xA5)
COLOR_DARK_GRAY = RGBColor(0x4A, 0x55, 0x68)
COLOR_TEXT = RGBColor(0x2D, 0x37, 0x48)
COLOR_BLUE = RGBColor(0x25, 0x63, 0xEB)
COLOR_GREEN = RGBColor(0x05, 0x96, 0x69)
COLOR_RED = RGBColor(0xDC, 0x26, 0x26)
COLOR_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
COLOR_ORANGE = RGBColor(0xD9, 0x77, 0x06)
COLOR_ROW_EVEN = RGBColor(0xF8, 0xF9, 0xFB)
COLOR_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)

# Legacy aliases
COLOR_PRIMARY = COLOR_BLUE
COLOR_HEADER_BG = COLOR_NAVY
COLOR_GREY = COLOR_MID_GRAY
COLOR_LIGHT_BLUE = COLOR_LIGHT_GRAY
COLOR_LIGHT_GREY = COLOR_OFF_WHITE
COLOR_DARK_TEXT = COLOR_TEXT

# Fonts
FONT_HEADING = "Georgia"
FONT_BODY = "Calibri"
DEFAULT_FONT = FONT_BODY

# Layout grid (inches)
MARGIN_LEFT = 0.5
MARGIN_RIGHT = 0.5
HEADER_BAR_Y = 0.35
TITLE_Y = 0.55
SUBTITLE_Y = 1.05
CONTENT_START_Y = 1.5
CONTENT_END_Y = 5.0
FOOTER_Y = 5.25
LEFT_X = 0.5
LEFT_W = 4.3
RIGHT_X = 5.2
RIGHT_W = 4.3
FULL_W = 9.0

# Legacy content area constants (used by compute_layout)
CONTENT_X = MARGIN_LEFT
CONTENT_Y = CONTENT_START_Y
CONTENT_W = FULL_W
CONTENT_H = CONTENT_END_Y - CONTENT_START_Y
CONTENT_PAD = 0.15

# ============================================================
# Presets & Mappings
# ============================================================
ROLE_PRESETS = {
    "title":      {"font_size": 22, "bold": True,  "color": "#1A2744"},
    "subtitle":   {"font_size": 11, "bold": False, "color": "#8892A5"},
    "body":       {"font_size": 11, "bold": False, "color": "#2D3748"},
    "label":      {"font_size": 13, "bold": True,  "color": "#1A2744"},
    "kpi_number": {"font_size": 36, "bold": True,  "color": "#C5973B"},
}

SHAPE_TYPE_MAP = {
    "rectangle":     MSO_SHAPE.RECTANGLE,
    "rounded_rect":  MSO_SHAPE.ROUNDED_RECTANGLE,
    "circle":        MSO_SHAPE.OVAL,
    "oval":          MSO_SHAPE.OVAL,
    "diamond":       MSO_SHAPE.DIAMOND,
    "hexagon":       MSO_SHAPE.HEXAGON,
    "triangle":      MSO_SHAPE.ISOSCELES_TRIANGLE,
    "arrow_right":   MSO_SHAPE.RIGHT_ARROW,
    "chevron_right": MSO_SHAPE.CHEVRON,
    "line":          MSO_SHAPE.RECTANGLE,
}

CHART_TYPE_MAP = {
    "bar":   XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":  XL_CHART_TYPE.LINE_MARKERS,
    "pie":   XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
}

ICON_MAP = {
    "check": "\u2713", "warning": "\u26A0", "star": "\u2605",
    "chart": "\U0001F4CA", "people": "\U0001F465", "money": "\U0001F4B0",
    "trending_up": "\u2197", "building": "\U0001F3E2", "globe": "\U0001F310",
    "lightbulb": "\U0001F4A1", "arrow": "\u2192", "target": "\U0001F3AF",
}

# Map old layout_hint values to new NP template types
_LEGACY_LAYOUT_MAP = {
    "single_column":  "two_column",
    "two_column":     "two_column",
    "three_column":   "two_column",
    "kpi_row":        "kpi_dashboard",
    "chart_with_text": "chart_table",
    "text_with_chart": "chart_table",
    "process_flow":   "timeline_flow",
    "timeline":       "timeline_flow",
    "quote":          "two_column",
    "grid":           "kpi_dashboard",
    "full_image":     "two_column",
    "auto":           "two_column",
}

# Map for legacy slide_type + layout_hint -> NP template
_NP_TEMPLATE_TYPES = {
    "title", "divider", "data_table", "chart_table",
    "two_column", "kpi_dashboard", "risk_matrix",
    "timeline_flow", "comparison",
}

# NP accent color cycle for cards/badges
_ACCENT_COLORS = [COLOR_BLUE, COLOR_GREEN, COLOR_PURPLE, COLOR_ORANGE, COLOR_RED, COLOR_GOLD]


# ============================================================
# Helper Functions
# ============================================================
def clean_text(text):
    """Markdown 문법 제거"""
    if not text:
        return ""
    return text.replace('**', '').strip()


def parse_hex_color(hex_str):
    """'#1E3A8A' -> RGBColor"""
    if not hex_str or not isinstance(hex_str, str):
        return None
    hex_str = hex_str.lstrip('#')
    if len(hex_str) != 6:
        return None
    try:
        r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return None


def set_font(paragraph, size, bold=False, color=None, font_name=DEFAULT_FONT):
    """폰트 설정 헬퍼"""
    paragraph.font.name = font_name
    paragraph.font.size = size
    paragraph.font.bold = bold
    if color:
        paragraph.font.color.rgb = color


def _set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _color_for_name(name):
    """Resolve a color name string to RGBColor."""
    name_map = {
        "blue": COLOR_BLUE, "green": COLOR_GREEN, "red": COLOR_RED,
        "purple": COLOR_PURPLE, "orange": COLOR_ORANGE, "gold": COLOR_GOLD,
        "navy": COLOR_NAVY, "gray": COLOR_MID_GRAY, "white": COLOR_WHITE,
    }
    if isinstance(name, str) and name.startswith("#"):
        return parse_hex_color(name) or COLOR_BLUE
    return name_map.get(str(name).lower(), COLOR_BLUE)


# ============================================================
# NP Component Functions
# ============================================================
def add_np_header(slide, title, subtitle=""):
    """Navy accent bar at top + title + optional subtitle."""
    # Navy bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(HEADER_BAR_Y),
        SLIDE_WIDTH, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_NAVY
    bar.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(
        Inches(MARGIN_LEFT), Inches(TITLE_Y),
        Inches(FULL_W), Inches(0.45)
    )
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = clean_text(title)
    set_font(p, Pt(22), bold=True, color=COLOR_NAVY, font_name=FONT_HEADING)

    # Subtitle
    if subtitle:
        sb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(SUBTITLE_Y),
            Inches(FULL_W), Inches(0.35)
        )
        sb.text_frame.word_wrap = True
        ps = sb.text_frame.paragraphs[0]
        ps.text = clean_text(subtitle)
        set_font(ps, Pt(11), color=COLOR_MID_GRAY, font_name=FONT_BODY)


def add_np_summary(slide, text, y=None):
    """Render 1-2 line insight summary between header and content.
    Returns the Y position where content should start after the summary.
    """
    if not text:
        return CONTENT_START_Y
    summary_y = y or (SUBTITLE_Y + 0.35)
    tb = slide.shapes.add_textbox(
        Inches(MARGIN_LEFT), Inches(summary_y),
        Inches(FULL_W), Inches(0.55)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(str(text).split('\n')[:2]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = clean_text(line)
        set_font(p, Pt(9.5), bold=True, color=COLOR_DARK_GRAY, font_name=FONT_BODY)
        p.space_after = Pt(2)
    return summary_y + 0.55


def add_np_footer(slide, page_num, draft=True):
    """Dark navy footer bar with CONFIDENTIAL + page number + optional DRAFT badge."""
    # Footer bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(FOOTER_Y),
        SLIDE_WIDTH, Inches(0.375)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_DARK_NAVY
    bar.line.fill.background()

    # CONFIDENTIAL label
    conf = slide.shapes.add_textbox(
        Inches(MARGIN_LEFT), Inches(FOOTER_Y + 0.06),
        Inches(3.0), Inches(0.28)
    )
    pc = conf.text_frame.paragraphs[0]
    pc.text = "CONFIDENTIAL"
    set_font(pc, Pt(7), bold=True, color=COLOR_MID_GRAY, font_name=FONT_BODY)

    # Page number
    pn = slide.shapes.add_textbox(
        Inches(8.0), Inches(FOOTER_Y + 0.06),
        Inches(1.5), Inches(0.28)
    )
    pp = pn.text_frame.paragraphs[0]
    pp.text = str(page_num)
    pp.alignment = PP_ALIGN.RIGHT
    set_font(pp, Pt(7), color=COLOR_MID_GRAY, font_name=FONT_BODY)

    # DRAFT badge
    if draft:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.8), Inches(0.15),
            Inches(0.8), Inches(0.25)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_RED
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = "DRAFT"
        p.alignment = PP_ALIGN.CENTER
        set_font(p, Pt(7), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)


def add_np_table(slide, headers, rows, x, y, w, col_widths=None,
                 has_total_row=False, row_height=0.24):
    """NP-styled table with navy header, alternating rows, optional total row."""
    if not headers:
        return
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    table_h = Inches(row_height * n_rows)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), table_h
    )
    table = tbl_shape.table

    # Column widths
    if col_widths and len(col_widths) == n_cols:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = Inches(cw)
    else:
        equal_w = int(Inches(w) / n_cols)
        for ci in range(n_cols):
            table.columns[ci].width = equal_w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = clean_text(str(hdr))
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_font(p, Pt(8), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

    # Data rows
    for ri, row_data in enumerate(rows):
        is_total = has_total_row and ri == len(rows) - 1
        for ci in range(n_cols):
            cell = table.cell(ri + 1, ci)
            cell_text = str(row_data[ci]) if ci < len(row_data) else ""
            cell.text = clean_text(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            if is_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_NAVY
                set_font(p, Pt(7.5), bold=True, color=COLOR_GOLD, font_name=FONT_BODY)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ROW_EVEN if ri % 2 == 0 else COLOR_ROW_ODD
                set_font(p, Pt(7.5), color=COLOR_TEXT, font_name=FONT_BODY)


def add_np_kpi_cards(slide, metrics, y=4.45):
    """Render a row of KPI cards. metrics: list of {label, value, sub, color}."""
    if not metrics:
        return
    n = min(len(metrics), 6)
    total_w = FULL_W
    gap = 0.15
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 0.5

    for i, m in enumerate(metrics[:n]):
        cx = MARGIN_LEFT + i * (card_w + gap)
        accent = _color_for_name(m.get("color", "blue"))

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx), Inches(y), Inches(card_w), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_LIGHT_GRAY
        card.line.width = Pt(0.5)

        # Left color accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx), Inches(y), Inches(0.04), Inches(card_h)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        # Value
        vb = slide.shapes.add_textbox(
            Inches(cx + 0.12), Inches(y + 0.02),
            Inches(card_w - 0.2), Inches(0.26)
        )
        pv = vb.text_frame.paragraphs[0]
        pv.text = str(m.get("value", ""))
        set_font(pv, Pt(14), bold=True, color=accent, font_name=FONT_BODY)

        # Label
        lb = slide.shapes.add_textbox(
            Inches(cx + 0.12), Inches(y + 0.26),
            Inches(card_w - 0.2), Inches(0.18)
        )
        pl = lb.text_frame.paragraphs[0]
        pl.text = str(m.get("label", ""))
        set_font(pl, Pt(8), color=COLOR_DARK_GRAY, font_name=FONT_BODY)

        # Sub text (optional)
        sub = m.get("sub")
        if sub:
            sb = slide.shapes.add_textbox(
                Inches(cx + card_w - 1.0), Inches(y + 0.04),
                Inches(0.9), Inches(0.18)
            )
            ps = sb.text_frame.paragraphs[0]
            ps.text = str(sub)
            ps.alignment = PP_ALIGN.RIGHT
            is_positive = str(sub).startswith("+")
            sub_color = COLOR_GREEN if is_positive else COLOR_RED
            set_font(ps, Pt(7), bold=True, color=sub_color, font_name=FONT_BODY)


def add_np_banner(slide, y, label, text, full_width=True):
    """Navy background banner with gold label + white text."""
    bw = FULL_W if full_width else LEFT_W
    bx = MARGIN_LEFT

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(bx), Inches(y), Inches(bw), Inches(0.32)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(bx + 0.15), Inches(y + 0.04),
        Inches(bw - 0.3), Inches(0.24)
    )
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    # Gold label + white text in same paragraph via runs
    run_label = p.add_run()
    run_label.text = str(label) + " "
    run_label.font.name = FONT_BODY
    run_label.font.size = Pt(8)
    run_label.font.bold = True
    run_label.font.color.rgb = COLOR_GOLD

    run_text = p.add_run()
    run_text.text = str(text)
    run_text.font.name = FONT_BODY
    run_text.font.size = Pt(8)
    run_text.font.color.rgb = COLOR_WHITE


def add_np_source_line(slide, text, y=5.0):
    """Small gray italic source attribution line."""
    tb = slide.shapes.add_textbox(
        Inches(MARGIN_LEFT), Inches(y),
        Inches(FULL_W), Inches(0.2)
    )
    p = tb.text_frame.paragraphs[0]
    p.text = str(text)
    set_font(p, Pt(6.5), color=COLOR_MID_GRAY, font_name=FONT_BODY)
    p.font.italic = True


# ============================================================
# Element Renderers (kept for legacy/atomic element support)
# ============================================================
def render_text_box(slide, el):
    """text_box element -> pptx textbox"""
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    text = clean_text(el.get("text", ""))
    role = el.get("role", "body")
    preset = ROLE_PRESETS.get(role, ROLE_PRESETS["body"])

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True

    lines = text.split("\\n") if "\\n" in text else text.split("\n")
    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_bullet = line_text.lstrip().startswith("- ")
        if is_bullet:
            p.text = "\u2022 " + line_text.lstrip()[2:]
        else:
            p.text = line_text
        fs = Pt(el.get("font_size", preset["font_size"]))
        bold = el.get("bold", preset["bold"])
        color = parse_hex_color(el.get("color", preset["color"])) or COLOR_TEXT
        fn = el.get("font_name", FONT_BODY)
        set_font(p, fs, bold=bold, color=color, font_name=fn)
        align = el.get("alignment", "left")
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                        "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        if is_bullet:
            p.space_after = Pt(2)


def render_shape(slide, el):
    """shape element -> pptx autoshape"""
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    shape_type_str = el.get("shape_type", "rounded_rect")
    mso = SHAPE_TYPE_MAP.get(shape_type_str, MSO_SHAPE.ROUNDED_RECTANGLE)

    if shape_type_str == "line":
        h = 0.03

    shape = slide.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))

    fill_hex = el.get("fill", "#F5F6F8")
    fill_color = parse_hex_color(fill_hex)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color

    border_hex = el.get("border_color")
    if border_hex:
        border_color = parse_hex_color(border_hex)
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    text = el.get("text")
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        lines = text.split("\\n") if "\\n" in text else text.split("\n")
        for i, line_text in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line_text.strip()
            p.alignment = PP_ALIGN.CENTER
            set_font(p, Pt(11), bold=False, color=COLOR_TEXT)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].space_before = Pt(0)


def render_chart(slide, el):
    """chart element -> pptx chart"""
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    chart_type_str = el.get("chart_type", "bar")
    data = el.get("data", {})
    categories = data.get("categories", [])
    series_list = data.get("series", [])

    if not categories or not series_list:
        el["kind"] = "text_box"
        el["text"] = f"[Chart: {chart_type_str} - no data]"
        el["role"] = "body"
        render_text_box(slide, el)
        return

    xl_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for s in series_list:
        name = s.get("name", "Series")
        values = s.get("values", [])
        safe_values = []
        for v in values:
            try:
                safe_values.append(float(v))
            except (ValueError, TypeError):
                safe_values.append(0)
        chart_data.add_series(name, safe_values)

    chart_frame = slide.shapes.add_chart(
        xl_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data
    )

    chart = chart_frame.chart
    chart.has_legend = el.get("show_legend", True)
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT_BODY

    chart.font.name = FONT_BODY
    chart.font.size = Pt(9)

    # Apply NP colors to chart series
    np_colors = [COLOR_NAVY, COLOR_GOLD, COLOR_BLUE, COLOR_GREEN, COLOR_PURPLE, COLOR_ORANGE]
    try:
        plot = chart.plots[0]
        for idx, series in enumerate(plot.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = np_colors[idx % len(np_colors)]
    except Exception:
        pass


def render_callout(slide, el):
    """callout (KPI card) -> rounded rect + value + label + delta"""
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    value = el.get("value", "")
    label = el.get("label", "")
    delta = el.get("delta", "")
    accent = parse_hex_color(el.get("accent_color", "#C5973B")) or COLOR_GOLD

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_LIGHT_GRAY
    card.line.width = Pt(0.5)

    # Left accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(0.04), Inches(h)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    # Icon
    icon_name = el.get("icon")
    icon_char = ICON_MAP.get(icon_name, "")
    icon_offset = 0.0
    if icon_char:
        ib = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.1), Inches(0.4), Inches(0.4)
        )
        p = ib.text_frame.paragraphs[0]
        p.text = icon_char
        set_font(p, Pt(16), color=accent)
        icon_offset = 0.1

    # Value
    val_y = y + 0.08 + icon_offset
    vb = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(val_y), Inches(w - 0.3), Inches(0.5)
    )
    vb.text_frame.word_wrap = True
    p_val = vb.text_frame.paragraphs[0]
    p_val.text = str(value)
    p_val.alignment = PP_ALIGN.LEFT
    set_font(p_val, Pt(min(28, max(18, int(36 - len(str(value)))))),
             bold=True, color=accent)

    # Label
    lb = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(val_y + 0.45), Inches(w - 0.3), Inches(0.3)
    )
    p_lbl = lb.text_frame.paragraphs[0]
    p_lbl.text = label
    set_font(p_lbl, Pt(9), color=COLOR_DARK_GRAY)

    # Delta
    if delta:
        db = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(val_y + 0.7), Inches(w - 0.3), Inches(0.2)
        )
        p_d = db.text_frame.paragraphs[0]
        p_d.text = delta
        is_positive = delta.startswith("+")
        delta_color = COLOR_GREEN if is_positive else COLOR_RED
        set_font(p_d, Pt(9), bold=True, color=delta_color)


def render_icon(slide, el):
    """icon element -> textbox with unicode character"""
    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    name = el.get("name", "star")
    char = ICON_MAP.get(name, "\u2605")
    color = parse_hex_color(el.get("color", "#C5973B")) or COLOR_GOLD

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = tb.text_frame.paragraphs[0]
    p.text = char
    p.alignment = PP_ALIGN.CENTER
    set_font(p, Pt(el.get("size", 24)), color=color)


# Dispatcher
RENDER_MAP = {
    "text_box": render_text_box,
    "shape":    render_shape,
    "chart":    render_chart,
    "callout":  render_callout,
    "icon":     render_icon,
}


def render_element(slide, el):
    """Dispatch element rendering by kind."""
    kind = el.get("kind", "text_box")
    fn = RENDER_MAP.get(kind)
    if fn:
        try:
            fn(slide, el)
        except Exception as e:
            fb = slide.shapes.add_textbox(
                Inches(el.get("x", 0.3)), Inches(el.get("y", 1.2)),
                Inches(el.get("w", 4)), Inches(el.get("h", 0.5))
            )
            fb.text_frame.paragraphs[0].text = f"[{kind} render error: {e}]"


# ============================================================
# Layout Engine (adapted for NP grid)
# ============================================================
def compute_layout(layout_hint, elements):
    """Assign x, y, w, h to elements that lack coordinates."""
    if not elements:
        return elements

    needs_layout = [e for e in elements if not all(k in e for k in ("x", "y", "w", "h"))]
    if not needs_layout:
        return elements

    cx, cy, cw, ch = CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H
    pad = CONTENT_PAD

    if layout_hint in ("auto", None, ""):
        layout_hint = _auto_detect_layout(elements)

    if layout_hint == "single_column":
        _layout_single_column(needs_layout, cx, cy, cw, ch, pad)
    elif layout_hint == "two_column":
        _layout_columns(needs_layout, cx, cy, cw, ch, pad, n_cols=2)
    elif layout_hint == "three_column":
        _layout_columns(needs_layout, cx, cy, cw, ch, pad, n_cols=3)
    elif layout_hint == "kpi_row":
        _layout_kpi_row(needs_layout, cx, cy, cw, ch, pad)
    elif layout_hint == "chart_with_text":
        _layout_chart_text(needs_layout, cx, cy, cw, ch, pad, chart_left=True)
    elif layout_hint == "text_with_chart":
        _layout_chart_text(needs_layout, cx, cy, cw, ch, pad, chart_left=False)
    elif layout_hint == "process_flow":
        _layout_process_flow(needs_layout, cx, cy, cw, ch, pad)
    elif layout_hint == "timeline":
        _layout_timeline(needs_layout, cx, cy, cw, ch, pad)
    elif layout_hint == "quote":
        _layout_quote(needs_layout, cx, cy, cw, ch)
    elif layout_hint == "grid":
        _layout_grid(needs_layout, cx, cy, cw, ch, pad)
    elif layout_hint == "full_image":
        _layout_full_image(needs_layout, cx, cy, cw, ch, pad)
    else:
        _layout_single_column(needs_layout, cx, cy, cw, ch, pad)

    return elements


def _auto_detect_layout(elements):
    kinds = [e.get("kind") for e in elements]
    callout_count = kinds.count("callout")
    chart_count = kinds.count("chart")
    shape_count = kinds.count("shape")
    text_count = kinds.count("text_box")

    if callout_count >= 3:
        return "kpi_row"
    if chart_count >= 1 and text_count >= 1:
        return "chart_with_text"
    if shape_count >= 3 and chart_count == 0:
        return "process_flow"
    if text_count >= 4:
        return "two_column"
    return "single_column"


def _layout_single_column(elems, cx, cy, cw, ch, pad):
    n = len(elems)
    item_h = min(1.0, (ch - pad * (n - 1)) / max(n, 1))
    for i, el in enumerate(elems):
        el.setdefault("x", cx)
        el.setdefault("y", cy + i * (item_h + pad))
        el.setdefault("w", cw)
        el.setdefault("h", item_h)


def _layout_columns(elems, cx, cy, cw, ch, pad, n_cols=2):
    col_w = (cw - pad * (n_cols - 1)) / n_cols
    per_col = math.ceil(len(elems) / n_cols)
    for i, el in enumerate(elems):
        col = min(i // per_col, n_cols - 1)
        row = i % per_col
        item_h = min(1.0, (ch - pad * (per_col - 1)) / max(per_col, 1))
        el.setdefault("x", cx + col * (col_w + pad))
        el.setdefault("y", cy + row * (item_h + pad))
        el.setdefault("w", col_w)
        el.setdefault("h", item_h)


def _layout_kpi_row(elems, cx, cy, cw, ch, pad):
    callouts = [e for e in elems if e.get("kind") == "callout"]
    others = [e for e in elems if e.get("kind") != "callout"]

    if callouts:
        card_w = (cw - pad * (len(callouts) - 1)) / len(callouts)
        card_h = min(1.4, ch * 0.35)
        for i, el in enumerate(callouts):
            el.setdefault("x", cx + i * (card_w + pad))
            el.setdefault("y", cy)
            el.setdefault("w", card_w)
            el.setdefault("h", card_h)
        below_y = cy + card_h + pad
    else:
        below_y = cy

    remaining_h = cy + ch - below_y
    if others:
        item_h = min(0.8, (remaining_h - pad * (len(others) - 1)) / max(len(others), 1))
        for i, el in enumerate(others):
            el.setdefault("x", cx)
            el.setdefault("y", below_y + i * (item_h + pad))
            el.setdefault("w", cw)
            el.setdefault("h", item_h)


def _layout_chart_text(elems, cx, cy, cw, ch, pad, chart_left=True):
    charts = [e for e in elems if e.get("kind") == "chart"]
    others = [e for e in elems if e.get("kind") != "chart"]

    chart_w = cw * 0.58
    text_w = cw - chart_w - pad

    if chart_left:
        chart_x, text_x = cx, cx + chart_w + pad
    else:
        text_x, chart_x = cx, cx + text_w + pad

    if charts:
        ch_h = (ch - pad * (len(charts) - 1)) / max(len(charts), 1)
        for i, el in enumerate(charts):
            el.setdefault("x", chart_x)
            el.setdefault("y", cy + i * (ch_h + pad))
            el.setdefault("w", chart_w)
            el.setdefault("h", ch_h)

    if others:
        item_h = min(1.0, (ch - pad * (len(others) - 1)) / max(len(others), 1))
        for i, el in enumerate(others):
            el.setdefault("x", text_x)
            el.setdefault("y", cy + i * (item_h + pad))
            el.setdefault("w", text_w)
            el.setdefault("h", item_h)


def _layout_process_flow(elems, cx, cy, cw, ch, pad):
    n = len(elems)
    if n == 0:
        return
    arrow_w = 0.3
    total_arrows = max(n - 1, 0)
    shape_w = (cw - total_arrows * arrow_w - pad * total_arrows) / max(n, 1)
    shape_h = min(1.5, ch * 0.5)
    mid_y = cy + (ch - shape_h) / 2

    curr_x = cx
    for i, el in enumerate(elems):
        el.setdefault("x", curr_x)
        el.setdefault("y", mid_y)
        el.setdefault("w", shape_w)
        el.setdefault("h", shape_h)
        if el.get("kind") != "shape":
            el["kind"] = "shape"
            el.setdefault("shape_type", "rounded_rect")
        curr_x += shape_w + pad
        if i < n - 1:
            curr_x += arrow_w


def _layout_timeline(elems, cx, cy, cw, ch, pad):
    n = len(elems)
    if n == 0:
        return
    node_w = min(1.2, (cw - pad * (n - 1)) / max(n, 1))
    node_h = node_w
    mid_y = cy + ch * 0.4
    spacing = (cw - node_w) / max(n - 1, 1) if n > 1 else 0

    for i, el in enumerate(elems):
        el.setdefault("x", cx + i * spacing)
        el.setdefault("y", mid_y)
        el.setdefault("w", node_w)
        el.setdefault("h", node_h)
        if el.get("kind") != "shape":
            el["kind"] = "shape"
        el.setdefault("shape_type", "circle")


def _layout_quote(elems, cx, cy, cw, ch):
    if elems:
        el = elems[0]
        el.setdefault("x", cx + 0.5)
        el.setdefault("y", cy + ch * 0.2)
        el.setdefault("w", cw - 1.0)
        el.setdefault("h", ch * 0.4)
        el.setdefault("role", "title")
        el.setdefault("alignment", "center")
        el.setdefault("font_size", 28)
    if len(elems) > 1:
        el2 = elems[1]
        el2.setdefault("x", cx + 1.0)
        el2.setdefault("y", cy + ch * 0.65)
        el2.setdefault("w", cw - 2.0)
        el2.setdefault("h", 0.5)
        el2.setdefault("role", "subtitle")
        el2.setdefault("alignment", "center")


def _layout_grid(elems, cx, cy, cw, ch, pad):
    n = len(elems)
    if n == 0:
        return
    cols = math.ceil(math.sqrt(n))
    rows = math.ceil(n / cols)
    cell_w = (cw - pad * (cols - 1)) / cols
    cell_h = (ch - pad * (rows - 1)) / rows

    for i, el in enumerate(elems):
        r, c = divmod(i, cols)
        el.setdefault("x", cx + c * (cell_w + pad))
        el.setdefault("y", cy + r * (cell_h + pad))
        el.setdefault("w", cell_w)
        el.setdefault("h", cell_h)


def _layout_full_image(elems, cx, cy, cw, ch, pad):
    if elems:
        elems[0].setdefault("x", cx)
        elems[0].setdefault("y", cy)
        elems[0].setdefault("w", cw)
        elems[0].setdefault("h", ch)
    for el in elems[1:]:
        el.setdefault("x", cx + 0.5)
        el.setdefault("y", cy + ch - 1.2)
        el.setdefault("w", cw - 1.0)
        el.setdefault("h", 0.8)


# ============================================================
# NP Template Renderers
# ============================================================
def render_title(prs, data, page_num):
    """Cover slide — full navy background with title/subtitle/date."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY
    bg.line.fill.background()

    # Gold accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(2.3), Inches(1.5), Inches(0.04)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_GOLD
    accent.line.fill.background()

    # Title
    title_text = data.get("title", "")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9.0), Inches(1.0))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = clean_text(title_text)
    set_font(p, Pt(36), bold=True, color=COLOR_WHITE, font_name=FONT_HEADING)

    # Subtitle
    subtitle_text = data.get("subtitle", "")
    if subtitle_text:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9.0), Inches(0.5))
        sb.text_frame.word_wrap = True
        ps = sb.text_frame.paragraphs[0]
        ps.text = clean_text(subtitle_text)
        set_font(ps, Pt(16), color=COLOR_LIGHT_GRAY, font_name=FONT_BODY)

    # Date
    date_str = datetime.date.today().strftime("%Y.%m.%d")
    db = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9.0), Inches(0.3))
    pd = db.text_frame.paragraphs[0]
    pd.text = date_str
    set_font(pd, Pt(12), color=COLOR_MID_GRAY, font_name=FONT_BODY)

    return slide


def render_divider(prs, data, page_num):
    """Section divider — navy background, gold circle with number, titles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY
    bg.line.fill.background()

    # Section number circle
    section_num = data.get("section_number", str(page_num))
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.25), Inches(1.5), Inches(1.5), Inches(1.5)
    )
    circle.fill.background()
    circle.line.color.rgb = COLOR_GOLD
    circle.line.width = Pt(2)
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.alignment = PP_ALIGN.CENTER
    set_font(p, Pt(32), bold=True, color=COLOR_GOLD, font_name=FONT_HEADING)
    tf.paragraphs[0].space_before = Pt(12)

    # Korean title
    title = data.get("title", "")
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9.0), Inches(0.7))
    tb.text_frame.word_wrap = True
    pt = tb.text_frame.paragraphs[0]
    pt.text = clean_text(title)
    pt.alignment = PP_ALIGN.CENTER
    set_font(pt, Pt(28), bold=True, color=COLOR_WHITE, font_name=FONT_HEADING)

    # English title / subtitle
    en_title = data.get("subtitle", data.get("en_title", ""))
    if en_title:
        eb = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9.0), Inches(0.4))
        pe = eb.text_frame.paragraphs[0]
        pe.text = clean_text(en_title)
        pe.alignment = PP_ALIGN.CENTER
        set_font(pe, Pt(14), color=COLOR_MID_GRAY, font_name=FONT_BODY)

    # Optional items list
    items = data.get("items", [])
    if items:
        ib = slide.shapes.add_textbox(Inches(2.0), Inches(4.6), Inches(6.0), Inches(0.8))
        ib.text_frame.word_wrap = True
        for idx, item in enumerate(items):
            p = ib.text_frame.paragraphs[0] if idx == 0 else ib.text_frame.add_paragraph()
            p.text = f"\u2022 {clean_text(item)}"
            p.alignment = PP_ALIGN.CENTER
            set_font(p, Pt(11), color=COLOR_LIGHT_GRAY, font_name=FONT_BODY)
            p.space_after = Pt(2)

    return slide


def render_data_table(prs, data, page_num):
    """Full-width table + optional KPI cards + source."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    content_y = add_np_summary(slide, data.get("summary", ""))

    # Table
    table_data = data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    has_total = table_data.get("has_total_row", False)
    col_widths = table_data.get("col_widths")

    if headers and rows:
        add_np_table(slide, headers, rows,
                     x=MARGIN_LEFT, y=content_y, w=FULL_W,
                     col_widths=col_widths, has_total_row=has_total)

    # KPI cards
    metrics = data.get("metrics", data.get("kpi_cards", []))
    if metrics:
        add_np_kpi_cards(slide, metrics)

    # Source
    source = data.get("source", "")
    if source:
        add_np_source_line(slide, source)

    add_np_footer(slide, page_num)
    return slide


def render_chart_table(prs, data, page_num):
    """Left chart (60%) + right table (40%) + optional banners."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    content_y = add_np_summary(slide, data.get("summary", ""))

    # Left: chart
    chart_data = data.get("chart", {})
    chart_type_str = chart_data.get("chart_type", "bar")
    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])

    chart_x = LEFT_X
    chart_y = content_y
    chart_w = LEFT_W
    chart_h = min(2.8, CONTENT_END_Y - content_y - 0.6)

    if categories and series_list:
        el = {
            "x": chart_x, "y": chart_y, "w": chart_w, "h": chart_h,
            "kind": "chart", "chart_type": chart_type_str,
            "data": {"categories": categories, "series": series_list},
            "show_legend": chart_data.get("show_legend", True),
        }
        render_chart(slide, el)
    else:
        # Placeholder
        tb = slide.shapes.add_textbox(
            Inches(chart_x), Inches(chart_y), Inches(chart_w), Inches(chart_h)
        )
        tb.text_frame.paragraphs[0].text = "[Chart: no data]"
        set_font(tb.text_frame.paragraphs[0], Pt(11), color=COLOR_MID_GRAY)

    # Right: table
    table_data = data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    has_total = table_data.get("has_total_row", False)

    if headers and rows:
        add_np_table(slide, headers, rows,
                     x=RIGHT_X, y=content_y, w=RIGHT_W,
                     has_total_row=has_total)

    # Optional banner
    banner = data.get("banner", {})
    if banner:
        banner_y = chart_y + chart_h + 0.15
        add_np_banner(slide, banner_y,
                      banner.get("label", ""), banner.get("text", ""))

    # Source
    source = data.get("source", "")
    if source:
        add_np_source_line(slide, source)

    add_np_footer(slide, page_num)
    return slide


def render_two_column(prs, data, page_num):
    """Left/right independent content columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    content_y = add_np_summary(slide, data.get("summary", ""))

    def _render_column_content(col_data, col_x, col_w, start_y):
        """Render a column's content (text items, elements, or table)."""
        if not col_data:
            return

        # If col_data is a list of strings, render as text
        if isinstance(col_data, list):
            tb = slide.shapes.add_textbox(
                Inches(col_x), Inches(start_y),
                Inches(col_w), Inches(CONTENT_END_Y - start_y - 0.3)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(col_data):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                text = clean_text(str(item))
                if text.startswith("- "):
                    p.text = "\u2022 " + text[2:]
                else:
                    p.text = text
                set_font(p, Pt(10), color=COLOR_TEXT, font_name=FONT_BODY)
                p.space_after = Pt(3)
            return

        # If col_data is a dict with title + items
        if isinstance(col_data, dict):
            col_title = col_data.get("title", "")
            if col_title:
                header_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(col_x), Inches(start_y),
                    Inches(col_w), Inches(0.3)
                )
                header_bar.fill.solid()
                header_bar.fill.fore_color.rgb = COLOR_NAVY
                header_bar.line.fill.background()
                ht = slide.shapes.add_textbox(
                    Inches(col_x + 0.1), Inches(start_y + 0.03),
                    Inches(col_w - 0.2), Inches(0.24)
                )
                ph = ht.text_frame.paragraphs[0]
                ph.text = clean_text(col_title)
                set_font(ph, Pt(10), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)
                start_y += 0.35

            # Table inside column
            tbl = col_data.get("table", {})
            if tbl and tbl.get("headers"):
                add_np_table(slide, tbl["headers"], tbl.get("rows", []),
                             x=col_x, y=start_y, w=col_w,
                             has_total_row=tbl.get("has_total_row", False))
                return

            # Items
            items = col_data.get("items", [])
            if items:
                tb = slide.shapes.add_textbox(
                    Inches(col_x), Inches(start_y),
                    Inches(col_w), Inches(CONTENT_END_Y - start_y - 0.3)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                for idx, item in enumerate(items):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    text = clean_text(str(item))
                    if text.startswith("- "):
                        p.text = "\u2022 " + text[2:]
                    else:
                        p.text = text
                    set_font(p, Pt(10), color=COLOR_TEXT, font_name=FONT_BODY)
                    p.space_after = Pt(3)

            # Elements (atomic)
            elements = col_data.get("elements", [])
            if elements:
                elems = copy.deepcopy(elements)
                _layout_single_column(elems, col_x, start_y, col_w,
                                      CONTENT_END_Y - start_y - 0.3, CONTENT_PAD)
                for el in elems:
                    render_element(slide, el)

    left_data = data.get("left", data.get("left_column", {}))
    right_data = data.get("right", data.get("right_column", {}))

    _render_column_content(left_data, LEFT_X, LEFT_W, content_y)
    _render_column_content(right_data, RIGHT_X, RIGHT_W, content_y)

    add_np_footer(slide, page_num)
    return slide


def render_kpi_dashboard(prs, data, page_num):
    """Header + main content area + bottom KPI card row."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    content_y = add_np_summary(slide, data.get("summary", ""))

    # Main content: table or chart
    table_data = data.get("table", {})
    chart_data = data.get("chart", {})

    content_h = min(2.5, CONTENT_END_Y - content_y - 1.0)
    if table_data and table_data.get("headers"):
        add_np_table(slide, table_data["headers"], table_data.get("rows", []),
                     x=MARGIN_LEFT, y=content_y, w=FULL_W,
                     has_total_row=table_data.get("has_total_row", False))
    elif chart_data and chart_data.get("categories"):
        el = {
            "x": MARGIN_LEFT, "y": content_y, "w": FULL_W, "h": content_h,
            "kind": "chart",
            "chart_type": chart_data.get("chart_type", "bar"),
            "data": {
                "categories": chart_data.get("categories", []),
                "series": chart_data.get("series", []),
            },
            "show_legend": chart_data.get("show_legend", True),
        }
        render_chart(slide, el)

    # Text content
    content_text = data.get("content", "")
    if content_text:
        tb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT), Inches(CONTENT_START_Y),
            Inches(FULL_W), Inches(content_h)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        lines = content_text.split("\n")
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = clean_text(line)
            set_font(p, Pt(10), color=COLOR_TEXT, font_name=FONT_BODY)

    # KPI cards at bottom
    metrics = data.get("metrics", data.get("kpi_cards", []))
    if metrics:
        add_np_kpi_cards(slide, metrics)

    # Banner
    banner = data.get("banner", {})
    if banner:
        add_np_banner(slide, 4.1,
                      banner.get("label", ""), banner.get("text", ""))

    source = data.get("source", "")
    if source:
        add_np_source_line(slide, source)

    add_np_footer(slide, page_num)
    return slide


def render_risk_matrix(prs, data, page_num):
    """Categorized risk cards with colored badges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    add_np_summary(slide, data.get("summary", ""))

    risks = data.get("risks", data.get("items", []))
    if not risks:
        add_np_footer(slide, page_num)
        return slide

    # Layout: up to 6 risk cards in 2x3 grid
    n = min(len(risks), 6)
    cols = min(n, 3)
    rows_count = math.ceil(n / cols)
    card_w = (FULL_W - 0.15 * (cols - 1)) / cols
    card_h = (CONTENT_END_Y - CONTENT_START_Y - 0.15 * (rows_count - 1)) / rows_count

    severity_colors = {
        "high": COLOR_RED, "critical": COLOR_RED,
        "medium": COLOR_ORANGE, "moderate": COLOR_ORANGE,
        "low": COLOR_GREEN, "info": COLOR_BLUE,
    }

    for idx, risk in enumerate(risks[:n]):
        r, c = divmod(idx, cols)
        cx = MARGIN_LEFT + c * (card_w + 0.15)
        cy = CONTENT_START_Y + r * (card_h + 0.15)

        severity = str(risk.get("severity", risk.get("level", "medium"))).lower()
        sev_color = severity_colors.get(severity, COLOR_ORANGE)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx), Inches(cy), Inches(card_w), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_LIGHT_GRAY
        card.line.width = Pt(0.5)

        # Top color bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx), Inches(cy), Inches(card_w), Inches(0.04)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = sev_color
        top_bar.line.fill.background()

        # Severity badge
        badge_text = severity.upper()
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx + card_w - 0.8), Inches(cy + 0.08),
            Inches(0.7), Inches(0.2)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = sev_color
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = badge_text
        bp.alignment = PP_ALIGN.CENTER
        set_font(bp, Pt(6), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

        # Category / title
        category = risk.get("category", risk.get("title", ""))
        if category:
            ct = slide.shapes.add_textbox(
                Inches(cx + 0.1), Inches(cy + 0.08),
                Inches(card_w - 1.0), Inches(0.2)
            )
            pc = ct.text_frame.paragraphs[0]
            pc.text = clean_text(category)
            set_font(pc, Pt(8), bold=True, color=COLOR_NAVY, font_name=FONT_BODY)

        # Description
        desc = risk.get("description", risk.get("text", ""))
        if desc:
            dt = slide.shapes.add_textbox(
                Inches(cx + 0.1), Inches(cy + 0.32),
                Inches(card_w - 0.2), Inches(card_h - 0.42)
            )
            dt.text_frame.word_wrap = True
            pd = dt.text_frame.paragraphs[0]
            pd.text = clean_text(desc)
            set_font(pd, Pt(7.5), color=COLOR_TEXT, font_name=FONT_BODY)

    add_np_footer(slide, page_num)
    return slide


def render_timeline_flow(prs, data, page_num):
    """Horizontal timeline/process with connected nodes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    add_np_summary(slide, data.get("summary", ""))

    nodes = data.get("nodes", data.get("items", data.get("steps", [])))
    if not nodes:
        add_np_footer(slide, page_num)
        return slide

    n = min(len(nodes), 8)
    area_w = FULL_W
    area_y = CONTENT_START_Y + 0.5
    node_w = min(1.1, (area_w - 0.3 * (n - 1)) / n)
    node_h = 0.6
    spacing = (area_w - node_w * n) / max(n - 1, 1) if n > 1 else 0

    # Connector line
    if n > 1:
        line_y = area_y + node_h / 2
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(MARGIN_LEFT + node_w / 2), Inches(line_y - 0.015),
            Inches(area_w - node_w), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        line.line.fill.background()

    for i, node in enumerate(nodes[:n]):
        nx = MARGIN_LEFT + i * (node_w + spacing)
        accent = _ACCENT_COLORS[i % len(_ACCENT_COLORS)]

        # Circle node
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(nx), Inches(area_y), Inches(node_w), Inches(node_h)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()

        # Node number/label
        cf = circle.text_frame
        cf.word_wrap = False
        cp = cf.paragraphs[0]
        node_label = node.get("label", str(i + 1)) if isinstance(node, dict) else str(i + 1)
        cp.text = str(node_label)
        cp.alignment = PP_ALIGN.CENTER
        set_font(cp, Pt(11), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

        # Title below
        node_title = node.get("title", node) if isinstance(node, dict) else str(node)
        tt = slide.shapes.add_textbox(
            Inches(nx - 0.15), Inches(area_y + node_h + 0.1),
            Inches(node_w + 0.3), Inches(0.35)
        )
        tt.text_frame.word_wrap = True
        tp = tt.text_frame.paragraphs[0]
        tp.text = clean_text(str(node_title))
        tp.alignment = PP_ALIGN.CENTER
        set_font(tp, Pt(8), bold=True, color=COLOR_NAVY, font_name=FONT_BODY)

        # Description below title
        node_desc = node.get("description", "") if isinstance(node, dict) else ""
        if node_desc:
            db = slide.shapes.add_textbox(
                Inches(nx - 0.15), Inches(area_y + node_h + 0.45),
                Inches(node_w + 0.3), Inches(0.6)
            )
            db.text_frame.word_wrap = True
            dp = db.text_frame.paragraphs[0]
            dp.text = clean_text(str(node_desc))
            dp.alignment = PP_ALIGN.CENTER
            set_font(dp, Pt(7), color=COLOR_DARK_GRAY, font_name=FONT_BODY)

    add_np_footer(slide, page_num)
    return slide


def render_comparison(prs, data, page_num):
    """VS comparison table with two sides."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    add_np_summary(slide, data.get("summary", ""))

    left = data.get("left", {})
    right = data.get("right", {})

    # VS badge in center
    vs = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.5), Inches(CONTENT_START_Y - 0.05),
        Inches(1.0), Inches(0.35)
    )
    vs.fill.solid()
    vs.fill.fore_color.rgb = COLOR_GOLD
    vs.line.fill.background()
    vp = vs.text_frame.paragraphs[0]
    vp.text = "VS"
    vp.alignment = PP_ALIGN.CENTER
    set_font(vp, Pt(11), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

    def _render_side(side_data, sx, sw, accent):
        side_title = side_data.get("title", side_data.get("name", ""))
        # Title bar
        hbar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(sx), Inches(CONTENT_START_Y + 0.4),
            Inches(sw), Inches(0.3)
        )
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = accent
        hbar.line.fill.background()
        ht = slide.shapes.add_textbox(
            Inches(sx + 0.1), Inches(CONTENT_START_Y + 0.43),
            Inches(sw - 0.2), Inches(0.24)
        )
        hp = ht.text_frame.paragraphs[0]
        hp.text = clean_text(side_title)
        set_font(hp, Pt(10), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

        # Items
        items = side_data.get("items", [])
        if items:
            tb = slide.shapes.add_textbox(
                Inches(sx + 0.1), Inches(CONTENT_START_Y + 0.8),
                Inches(sw - 0.2), Inches(CONTENT_END_Y - CONTENT_START_Y - 1.3)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(items):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                text = clean_text(str(item))
                p.text = f"\u2022 {text}"
                set_font(p, Pt(9), color=COLOR_TEXT, font_name=FONT_BODY)
                p.space_after = Pt(3)

        # Table
        tbl = side_data.get("table", {})
        if tbl and tbl.get("headers"):
            add_np_table(slide, tbl["headers"], tbl.get("rows", []),
                         x=sx, y=CONTENT_START_Y + 0.8, w=sw,
                         has_total_row=tbl.get("has_total_row", False))

    _render_side(left, LEFT_X, LEFT_W, COLOR_BLUE)
    _render_side(right, RIGHT_X, RIGHT_W, COLOR_NAVY)

    # Verdict banner
    verdict = data.get("verdict", data.get("conclusion", ""))
    if verdict:
        add_np_banner(slide, CONTENT_END_Y - 0.35, "결론:", verdict)

    add_np_footer(slide, page_num)
    return slide


def render_numbered_blocks(prs, data, page_num):
    """Numbered insight blocks (01, 02, 03...) in a 2-column grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    content_y = add_np_summary(slide, data.get("summary", ""))

    blocks = data.get("blocks", data.get("items", []))
    if not blocks:
        add_np_footer(slide, page_num)
        return slide

    n = len(blocks)
    cols = 2 if n > 1 else 1
    rows_count = math.ceil(n / cols)

    col_w = (FULL_W - 0.3) / cols
    avail_h = CONTENT_END_Y - content_y - 0.15
    block_h = min(avail_h / rows_count - 0.1, 1.4)

    for idx, block in enumerate(blocks):
        col = idx % cols
        row = idx // cols
        bx = LEFT_X + col * (col_w + 0.3)
        by = content_y + row * (block_h + 0.1)

        # Number badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bx), Inches(by),
            Inches(0.45), Inches(0.35)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_NAVY
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = str(block.get("number", f"{idx + 1:02d}"))
        bp.alignment = PP_ALIGN.CENTER
        set_font(bp, Pt(12), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

        # Title
        title_tb = slide.shapes.add_textbox(
            Inches(bx + 0.55), Inches(by),
            Inches(col_w - 0.55), Inches(0.3)
        )
        tp = title_tb.text_frame.paragraphs[0]
        tp.text = clean_text(block.get("title", ""))
        set_font(tp, Pt(11), bold=True, color=COLOR_NAVY, font_name=FONT_BODY)

        # Description
        desc = block.get("description", "")
        if desc:
            desc_tb = slide.shapes.add_textbox(
                Inches(bx + 0.55), Inches(by + 0.32),
                Inches(col_w - 0.55), Inches(block_h - 0.35)
            )
            tf = desc_tb.text_frame
            tf.word_wrap = True
            for li, line in enumerate(str(desc).split('\n')):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = clean_text(line)
                set_font(p, Pt(8.5), color=COLOR_DARK_GRAY, font_name=FONT_BODY)
                p.space_after = Pt(2)

        # Subtle divider line under each block
        if row < rows_count - 1 or col == 0:
            slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(bx), Inches(by + block_h),
                Inches(col_w), Inches(0.015)
            ).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = COLOR_LIGHT_GRAY
            slide.shapes[-1].line.fill.background()

    source = data.get("source", "")
    if source:
        add_np_source_line(slide, source)
    add_np_footer(slide, page_num)
    return slide


def render_grid_cards(prs, data, page_num):
    """Card grid layout (2-4 cards in a row) for product lineups, use cases, etc."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, data.get("title", ""), data.get("subtitle", ""))
    content_y = add_np_summary(slide, data.get("summary", ""))

    cards = data.get("cards", [])
    if not cards:
        add_np_footer(slide, page_num)
        return slide

    n = len(cards)
    cols = min(n, 4)
    rows_count = math.ceil(n / cols)
    gap = 0.2
    card_w = (FULL_W - gap * (cols - 1)) / cols
    avail_h = CONTENT_END_Y - content_y - 0.15
    card_h = min(avail_h / rows_count - 0.1, 2.8)

    for idx, card in enumerate(cards):
        col = idx % cols
        row = idx // cols
        cx = LEFT_X + col * (card_w + gap)
        cy = content_y + row * (card_h + 0.1)

        # Card background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx), Inches(cy),
            Inches(card_w), Inches(card_h)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_WHITE
        bg.line.color.rgb = COLOR_LIGHT_GRAY
        bg.line.width = Pt(0.75)

        # Card header bar
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cx), Inches(cy),
            Inches(card_w), Inches(0.35)
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = COLOR_NAVY
        hdr.line.fill.background()

        # Card title
        ttb = slide.shapes.add_textbox(
            Inches(cx + 0.1), Inches(cy + 0.03),
            Inches(card_w - 0.2), Inches(0.3)
        )
        tp = ttb.text_frame.paragraphs[0]
        tp.text = clean_text(card.get("title", ""))
        set_font(tp, Pt(10), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

        # Card subtitle
        sub = card.get("subtitle", "")
        text_y = cy + 0.4
        if sub:
            stb = slide.shapes.add_textbox(
                Inches(cx + 0.1), Inches(text_y),
                Inches(card_w - 0.2), Inches(0.2)
            )
            sp = stb.text_frame.paragraphs[0]
            sp.text = clean_text(sub)
            set_font(sp, Pt(8), color=COLOR_MID_GRAY, font_name=FONT_BODY)
            text_y += 0.22

        # Card items
        items = card.get("items", [])
        if items:
            itb = slide.shapes.add_textbox(
                Inches(cx + 0.1), Inches(text_y),
                Inches(card_w - 0.2), Inches(card_h - (text_y - cy) - 0.1)
            )
            tf = itb.text_frame
            tf.word_wrap = True
            for li, item in enumerate(items):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                text = clean_text(str(item))
                p.text = f"\u2022 {text}" if not text.startswith("\u2022") else text
                set_font(p, Pt(8), color=COLOR_TEXT, font_name=FONT_BODY)
                p.space_after = Pt(2)

        # Card table (alternative to items)
        tbl = card.get("table", {})
        if tbl and tbl.get("headers") and not items:
            tbl_y = text_y + 0.05
            add_np_table(slide, tbl["headers"], tbl.get("rows", []),
                         x=cx + 0.05, y=tbl_y, w=card_w - 0.1,
                         has_total_row=tbl.get("has_total_row", False))

    source = data.get("source", "")
    if source:
        add_np_source_line(slide, source)
    add_np_footer(slide, page_num)
    return slide


# NP template dispatch
_NP_RENDERERS = {
    "title":            render_title,
    "divider":          render_divider,
    "data_table":       render_data_table,
    "chart_table":      render_chart_table,
    "two_column":       render_two_column,
    "kpi_dashboard":    render_kpi_dashboard,
    "risk_matrix":      render_risk_matrix,
    "timeline_flow":    render_timeline_flow,
    "comparison":       render_comparison,
    "numbered_blocks":  render_numbered_blocks,
    "grid_cards":       render_grid_cards,
}


# ============================================================
# Legacy Slide Renderers (backward compatibility)
# ============================================================
def add_master_design(slide, title_text=""):
    """Legacy header bar — now delegates to NP header style."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.35)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_NAVY
    bar.line.fill.background()

    if title_text:
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(7), Inches(0.25))
        p = tb.text_frame.paragraphs[0]
        p.text = title_text
        set_font(p, Pt(9), bold=True, color=COLOR_WHITE, font_name=FONT_BODY)

    date_str = datetime.date.today().strftime("%Y-%m-%d")
    date_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.05), Inches(2.2), Inches(0.25))
    p_date = date_box.text_frame.paragraphs[0]
    p_date.text = date_str
    set_font(p_date, Pt(8), color=COLOR_WHITE, font_name=FONT_BODY)
    p_date.alignment = PP_ALIGN.RIGHT


def create_title_slide(prs, title_text, subtitle_text="Generated by GEM Intern AI"):
    """Legacy title slide — now uses NP style."""
    data = {"title": title_text, "subtitle": subtitle_text}
    return render_title(prs, data, 0)


def create_section_slide(prs, text):
    """Legacy section slide — alias for render_divider."""
    data = {"title": text, "section_number": ""}
    return render_divider(prs, data, 0)


def create_table_slide(prs, title_text, table_data):
    """Legacy table slide — updated to NP style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, title_text)

    if not table_data or len(table_data) < 2:
        add_np_footer(slide, 0)
        return slide

    headers = table_data[0]
    rows = table_data[1:]
    add_np_table(slide, headers, rows,
                 x=MARGIN_LEFT, y=CONTENT_START_Y, w=FULL_W)

    add_np_footer(slide, 0)
    return slide


# ============================================================
# Content Slide Renderer (legacy elements mode)
# ============================================================
def render_content_slide(prs, slide_dict):
    """Render a content slide using layout_hint + elements (legacy)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, slide_dict.get("title", ""))

    elements = slide_dict.get("elements", [])
    if not elements:
        return slide

    elems = copy.deepcopy(elements)
    layout_hint = slide_dict.get("layout_hint", "auto")
    compute_layout(layout_hint, elems)

    # Process flow arrows
    if layout_hint == "process_flow":
        shapes = [e for e in elems if e.get("kind") == "shape"]
        for i in range(len(shapes) - 1):
            a, b = shapes[i], shapes[i + 1]
            ax = a["x"] + a["w"]
            bx = b["x"]
            mid_y = a["y"] + a["h"] / 2
            arrow_w = bx - ax
            if arrow_w > 0.05:
                arr = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(ax + 0.02), Inches(mid_y - 0.1),
                    Inches(arrow_w - 0.04), Inches(0.2)
                )
                arr.fill.solid()
                arr.fill.fore_color.rgb = COLOR_NAVY
                arr.line.fill.background()

    # Timeline connector
    if layout_hint == "timeline" and len(elems) >= 2:
        first, last = elems[0], elems[-1]
        line_y = first.get("y", CONTENT_Y) + first.get("h", 1.0) / 2
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(first.get("x", CONTENT_X)),
            Inches(line_y - 0.015),
            Inches(last.get("x", 5) + last.get("w", 1) - first.get("x", CONTENT_X)),
            Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        line.line.fill.background()

    for el in elems:
        render_element(slide, el)

    return slide


# ============================================================
# Legacy Two-Column Renderer
# ============================================================
def _render_legacy_two_column(prs, slide_title, summary_text,
                              left_title, left_items, right_title, right_items):
    """Legacy 2-column rendering for backward compatibility."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_OFF_WHITE)
    add_np_header(slide, slide_title)

    if summary_text:
        sbg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(MARGIN_LEFT), Inches(CONTENT_START_Y - 0.3),
            Inches(FULL_W), Inches(0.5)
        )
        sbg.fill.solid()
        sbg.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        sbg.line.fill.background()
        sb = slide.shapes.add_textbox(
            Inches(MARGIN_LEFT + 0.1), Inches(CONTENT_START_Y - 0.25),
            Inches(FULL_W - 0.2), Inches(0.4)
        )
        sb.text_frame.word_wrap = True
        p_sum = sb.text_frame.paragraphs[0]
        p_sum.text = summary_text
        set_font(p_sum, Pt(10), color=COLOR_TEXT, font_name=FONT_BODY)

    for col_x, col_w, col_title, col_items in [
        (LEFT_X, LEFT_W, left_title, left_items),
        (RIGHT_X, RIGHT_W, right_title, right_items),
    ]:
        col_y = CONTENT_START_Y + 0.25
        # Column header bar
        tbg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(col_x), Inches(col_y),
            Inches(col_w), Inches(0.3)
        )
        tbg.fill.solid()
        tbg.fill.fore_color.rgb = COLOR_NAVY
        tbg.line.fill.background()
        tt = slide.shapes.add_textbox(
            Inches(col_x + 0.1), Inches(col_y + 0.03),
            Inches(col_w - 0.2), Inches(0.24)
        )
        tt.text_frame.paragraphs[0].text = col_title or ""
        set_font(tt.text_frame.paragraphs[0], Pt(10), bold=True,
                 color=COLOR_WHITE, font_name=FONT_BODY)

        cb = slide.shapes.add_textbox(
            Inches(col_x + 0.1), Inches(col_y + 0.4),
            Inches(col_w - 0.1), Inches(CONTENT_END_Y - col_y - 0.6)
        )
        cb.text_frame.word_wrap = True
        for idx, item in enumerate(col_items or []):
            p = cb.text_frame.paragraphs[0] if idx == 0 else cb.text_frame.add_paragraph()
            if item.get('type') == 'bullet':
                p.text = "\u2022 " + item['text']
                set_font(p, Pt(10), font_name=FONT_BODY)
                p.space_after = Pt(3)
            else:
                p.text = item.get('text', '')
                set_font(p, Pt(10), font_name=FONT_BODY)

    return slide


# ============================================================
# Main Entry Point
# ============================================================
def create_deck_from_json(json_data):
    """JSON -> PPTX bytes.

    Supports:
      - New NP schema (slide_type in NP template types)
      - Legacy schema (slide_type: content + layout_hint + elements)
      - Legacy schema (slide_type: content + left/right columns)
    """
    if isinstance(json_data, str):
        try:
            cleaned = json_data.strip()
            if cleaned.startswith("```"):
                cleaned = re.sub(r'^```\w*\n?', '', cleaned)
                cleaned = re.sub(r'\n?```$', '', cleaned)
            data = json.loads(cleaned)
        except json.JSONDecodeError:
            print("Invalid JSON for PPT")
            return None
    else:
        data = json_data

    slides = data.get('slides', [])
    if not slides:
        return None

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    page_num = 0

    for sd in slides:
        s_type = sd.get('slide_type', sd.get('type', 'content'))
        page_num += 1

        # Check if it's a new NP template type
        if s_type in _NP_RENDERERS:
            renderer = _NP_RENDERERS[s_type]
            renderer(prs, sd, page_num)
            continue

        # Legacy: slide_type == "content" with layout_hint -> map to NP template
        if s_type == 'content':
            layout_hint = sd.get('layout_hint', '')

            # If it has elements, use legacy element renderer
            if 'elements' in sd:
                render_content_slide(prs, sd)

            # If it has left/right (old legacy format)
            elif 'left' in sd or 'right' in sd:
                # Try to detect if it's structured for NP two_column
                left = sd.get('left', {})
                right = sd.get('right', {})

                if isinstance(left, dict) and 'items' not in left and 'title' not in left:
                    # Old-style left with sub-keys
                    summary = sd.get('summary', '')
                    l_items = [{'type': 'bullet', 'text': item}
                               for item in left.get('items', [])]
                    r_items = [{'type': 'bullet', 'text': item}
                               for item in right.get('items', [])]
                    _render_legacy_two_column(
                        prs, sd.get('title', ''), summary,
                        left.get('title', ''), l_items,
                        right.get('title', ''), r_items
                    )
                else:
                    # Use NP two_column renderer
                    render_two_column(prs, sd, page_num)

            else:
                # Bare content slide — render as two_column with content
                render_two_column(prs, sd, page_num)

        elif s_type == 'title':
            render_title(prs, sd, page_num)

        elif s_type == 'section':
            render_divider(prs, sd, page_num)

        else:
            # Unknown type: try NP mapping from layout_hint
            mapped = _LEGACY_LAYOUT_MAP.get(s_type, "two_column")
            if mapped in _NP_RENDERERS:
                _NP_RENDERERS[mapped](prs, sd, page_num)
            else:
                render_content_slide(prs, sd)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()
